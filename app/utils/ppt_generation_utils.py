import pandas as pd
import os
import datetime
from app.utils.common_utils import protect_pptx
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Pt
import warnings
from app.connections.pylogger import log_message
from app.utils.constants import LevelType
from app.utils.constants import LARGE_KEYS

DEFAULT_PT = Pt(10)
LARGE_PT = Pt(24)


warnings.filterwarnings('ignore')

# --------- Replace Text Placeholders ---------
def process_text_shape(shape, mapping):
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original = run.text
            for key, val in mapping.items():
                if key in original:
                    # replace and set size
                    new_text = original.replace(key, str(val))
                    run.text = new_text
                    run.font.size = LARGE_PT if key in LARGE_KEYS else DEFAULT_PT
                    # Bold only for <placeholders>
                    if key.startswith("<"):
                        run.font.bold = True
                    # update search text for multiple tags in same run
                    original = new_text

# --------- Currency Formatter ----------
def format_currency(num):
    try:
        num = float(num)
    except Exception:
        return str(num)
    
    abs_num = abs(num)
    if abs_num >= 1_000_000_000:
        return f"${num / 1_000_000_000:.1f}B"
    elif abs_num >= 1_000_000:
        return f"${num / 1_000_000:.1f}M"
    elif abs_num >= 1_000:
        return f"${num / 1_000:.1f}K"
    else:
        return f"${num:.2f}" 
        
# --------- Dollar Spend Drilldown (4a/4b/4c) ----------
def trim_drill_down(analysis_sku):
    columns_to_select = [
        ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 4_level_1', 'Current Instance', 'Unnamed: 4_level_3'),
        ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 8_level_1', 'Current Annual Cost', 'Unnamed: 8_level_3')
    ]
    analysis_sku_item4 = analysis_sku[columns_to_select].copy()
    analysis_sku_item4.columns = ['_'.join([str(i) for i in col if i]).strip() for col in analysis_sku_item4.columns]
    analysis_sku_item4 = analysis_sku_item4.rename(columns={
        'EPYC Cloud Cost Advisory Recommendations_Unnamed: 4_level_1_Current Instance_Unnamed: 4_level_3': 'Current Instance',
        'EPYC Cloud Cost Advisory Recommendations_Unnamed: 8_level_1_Current Annual Cost_Unnamed: 8_level_3': 'Current Annual Cost'
    })
    grouped_sum = analysis_sku_item4.groupby('Current Instance')['Current Annual Cost'].sum().reset_index()
    grouped_sum = grouped_sum.rename(columns={'Current Annual Cost': 'Total Annual Cost'})
    sorted_grouped = grouped_sum.sort_values(by='Total Annual Cost', ascending=False)
    return sorted_grouped

def dollar_spend_eval(sorted_grouped, total_cost):
    first_cost = sorted_grouped.iloc[0]['Total Annual Cost']
    max_spend = sorted_grouped.iloc[0]['Current Instance']
    item_4a = round((first_cost/total_cost)*100, 1) if total_cost != 0 else 0
    next_10_sum = sorted_grouped.iloc[1:11]['Total Annual Cost'].sum()
    item_4b = round((next_10_sum/total_cost)*100, 1) if total_cost != 0 else 0
    remaining_sum = sorted_grouped.iloc[11:]['Total Annual Cost'].sum()
    item_4c = round((remaining_sum/total_cost)*100, 1) if total_cost != 0 else 0
    return max_spend, item_4a, item_4b, item_4c

def generate_ppt(excel_file, email, pptx_template, results_path, password=None):
    try:
        if '@' in email and '.' in email.split('@')[-1]:
            customer = email.split('@')[1].split('.')[0].upper()
        else:
            customer = email.upper()

        # Get File Creation/Modification Date
        creation_time = os.path.getctime(excel_file)
        dt = datetime.datetime.fromtimestamp(creation_time)
        formatted = dt.strftime("%B %d, %Y")
        customer_short = os.path.splitext(os.path.basename(excel_file))[0]


        # --------- Read Data ---------
        # xls = pd.ExcelFile(excel_file)
        # Read recommended instance sheet with 4-level header
        data = pd.read_excel(excel_file, sheet_name='Recommended-Instance', header=[0, 1, 2, 3])

        # CSP
        item_2a = data.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 1_level_1', 'Cloud', 'Unnamed: 1_level_3')]

        # Grand Total
        analysis_total = data[data[('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 0_level_1', 'UUID/Instance Name', 'Unnamed: 0_level_3')] == "Grand Total"].reset_index(drop=True)

        # Main recommendation rows: CSP matches, no Remarks
        analysis_sku = data[
            (data[('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 27_level_1', 'Remark', 'Unnamed: 27_level_3')].isna()) &
            (data[('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 1_level_1', 'Cloud', 'Unnamed: 1_level_3')] == item_2a)
        ].reset_index(drop=True)

        # --------- "3a", "3b", "3c" -----------
        item_3a = int(
            data.loc[
                data[('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 27_level_1', 'Remark', 'Unnamed: 27_level_3')].isna() &
                (data[('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 0_level_1', 'UUID/Instance Name', 'Unnamed: 0_level_3')] != "Grand Total"),
                ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 5_level_1', 'Quantity', 'Unnamed: 5_level_3')
            ].sum()
        )

        item_3b = analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 8_level_1', 'Current Annual Cost', 'Unnamed: 8_level_3')]
        item_3bF = format_currency(item_3b)
        item_3c = len(analysis_sku[('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 3_level_1', 'Region', 'Unnamed: 3_level_3')].unique())

        sorted_group = trim_drill_down(analysis_sku)
        if not sorted_group.empty:
            max_spend, item_4a, item_4b, item_4c = dollar_spend_eval(sorted_group, item_3b)
        else:
            max_spend, item_4a, item_4b, item_4c = "N/A", 0.0, 0.0, 0.0

        # --------- Recommendations ---------
        item_5a = format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 12_level_1', 'Hourly Cost Optimization', 'Annual Cost ($)')])
        item_5b = format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 18_level_1', 'Modernize', 'Annual Cost ($)')])
        item_5c = format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 24_level_1', 'Modernize & Downsize', 'Annual Cost ($)')])

        item_6a = format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 13_level_1', 'Hourly Cost Optimization', 'Annual Savings ($)')])
        item_6b = f"{round(float(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 14_level_1', 'Hourly Cost Optimization', 'Performance Improvement')]), 2)}x"

        item_7a = format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 19_level_1', 'Modernize', 'Annual Savings ($)')])
        item_7b = f"{round(float(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 20_level_1', 'Modernize', 'Performance Improvement')]), 2)}x"

        item_8a = format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 25_level_1', 'Modernize & Downsize', 'Annual Savings ($)')])
        item_8b = f"{round(float(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 26_level_1', 'Modernize & Downsize', 'Performance Improvement')]), 2)}x"

        current_instance_col = ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 4_level_1', 'Current Instance', 'Unnamed: 4_level_3')
        if not analysis_sku.empty and current_instance_col in analysis_sku.columns:
            ph1 = str(analysis_sku.iloc[0][current_instance_col])
        else:
            ph1 = ""
            
        recommend_1_instance_col = ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 9_level_1', 'Hourly Cost Optimization', 'Instance')
        if not analysis_sku.empty and recommend_1_instance_col in analysis_sku.columns:
            ph4 = str(analysis_sku.iloc[0][recommend_1_instance_col])
        else:
            ph4 = ""

        recommend_2_instance_col = ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 15_level_1', 'Modernize', 'Instance')
        if not analysis_sku.empty and recommend_2_instance_col in analysis_sku.columns:
            ph3 = str(analysis_sku.iloc[0][recommend_2_instance_col])
        else:
            ph3 = ""

        recommend_3_instance_col = ('EPYC Cloud Cost Advisory Recommendations', 'Unnamed: 21_level_1', 'Modernize & Downsize', 'Instance')
        if not analysis_sku.empty and recommend_3_instance_col in analysis_sku.columns:
            ph2 = str(analysis_sku.iloc[0][recommend_3_instance_col])
        else:
            ph2 = ""

        # --------- Load PPT and Replace Chart ---------
        prs = Presentation(pptx_template)
        slide = prs.slides[0]

        # Find existing chart shape
        chart_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "has_chart") and shape.has_chart:  # Robustness
                chart_shape = shape
                break

        if chart_shape:
            left, top, width, height = chart_shape.left, chart_shape.top, chart_shape.width, chart_shape.height
            sp = chart_shape._element
            sp.getparent().remove(sp)

            chart_data = CategoryChartData()
            chart_data.categories = [max_spend, 'Next 10', 'Rest']
            chart_data.add_series('Dollar Spend', (float(item_4a),float(item_4b),float(item_4c)))
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                left, top, width, height,
                chart_data
            ).chart
            chart.has_title = False
            chart.value_axis.has_major_gridlines = False
            chart.category_axis.tick_labels.font.size = Pt(7)
            chart.value_axis.tick_labels.font.size = Pt(7)
            for series in chart.series:
                series.has_data_labels = True
                data_labels = series.data_labels
                data_labels.number_format = '0.0"%"'
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.font.size = Pt(12)
            chart.plots[0].has_data_labels = True
            
            # Ensure plot-level labels (already set above)
            plot = chart.plots[0]
            plot.has_data_labels = True

            # Ensure series-level labels
            s = chart.series[0]
            s.has_data_labels = True

            # Force point-level labels for each column
            for p in s.points:
                dl = p.data_label
                dl.show_value = True            # critical to make labels appear
                dl.number_format = '0"%"'
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
                dl.font.size = Pt(12)

        # Mapping for all replacement values
        tag_map = {
            "customer_name": customer,
            "cloud_provider": item_2a,
            "Date_Format": formatted,
            "<3a>": item_3a,
            "<3b>": item_3bF,
            "<3bF>": item_3bF,
            "<3c>": item_3c,
            "<4a>": item_4a,
            "<4b>": item_4b,
            "<4c>": item_4c,
            "<5a>": item_5a,
            "<5b>": item_5b,
            "<5c>": item_5c,
            "<6a>": item_6a,
            "<6b>": item_6b,
            "<7a>": item_7a,
            "<7b>": item_7b,
            "<8a>": item_8a,
            "<8b>": item_8b,
            "ph1": ph1,
            "ph2": ph2,
            "ph3": ph3,
            "ph4": ph4,
        }

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 6:  # group shape
                    for sub_shape in shape.shapes:
                        process_text_shape(sub_shape, tag_map)
                else:
                    process_text_shape(shape, tag_map)

        # --------- Save PPTX ----------
        output_name = f"{results_path}/{customer_short}.pptx"
        if os.path.exists(output_name):
            os.remove(output_name)
        prs.save(output_name)
         # --------- Apply password protection if password provided ----------
        if password:
            protect_pptx(output_name, password)
            log_message(LevelType.INFO, f"ppt protected for {output_name}", ErrorCode=1)
        return customer_short + '.pptx'
    except Exception as err:
        log_message(LevelType.ERROR, f"Error: {str(err)}", ErrorCode=-1)
