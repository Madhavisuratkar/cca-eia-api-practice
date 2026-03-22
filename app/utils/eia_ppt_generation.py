"""
EIA PowerPoint Generation Module

This module generates PowerPoint presentations from Excel data using a template.
It populates various sections including Cloud Footprint, Business Value, Path to Savings,
Dollar Spend Distribution, and Power/Carbon charts.
"""

import re
from openpyxl import load_workbook
from datetime import datetime
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from app.connections.pylogger import log_message
from app.utils.constants import LevelType
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from app.utils.common_utils import protect_pptx



def safe_float(value, default=0.0):
    """Safely convert a value to float, returning default if conversion fails."""
    try:
        if value is None:
            return default
        return float(value)
    except (ValueError, TypeError):
        return default


def format_value(value):
    """
    Format numeric values with K (thousands) or M (millions) suffix.
    Examples:
        1234 -> 1.23K
        1234567 -> 1.23M
        123 -> 123
    """
    try:
        num = float(value)
        if abs(num) >= 1_000_000_000:
            return f"{num / 1_000_000_000:.2f}B"
        elif abs(num) >= 1_000_000:
            return f"{num / 1_000_000:.2f}M"
        elif abs(num) >= 1_000:
            return f"{num / 1_000:.2f}k"
        else:
            return f"{num:.2f}"
    except (ValueError, TypeError):
        return str(value)



def extract_excel_data(excel_path):
    """
    Extract all required data from the Excel file.
    
    Args:
        excel_path: Path to the Excel file
        
    Returns:
        dict: Dictionary containing all extracted data
    """
    try:
        log_message(LevelType.INFO, f"Reading Excel file: {excel_path}")
        
        wb = load_workbook(excel_path)
        ws = wb['Recommended-Instance']
        
        data = {
            'instances': [],
            'grand_total': {},
            'unique_regions': set()
        }
        
        # Find data rows (start from row 5, skip headers in rows 1-4)
        data_row_start = 5
        
        for row_idx in range(data_row_start, ws.max_row + 1):
            # Check if this is the Grand Total row
            if ws.cell(row_idx, 1).value == "Grand Total":
                data['grand_total'] = {
                    'current_cost': safe_float(ws.cell(row_idx, 3).value),
                    'current_power': safe_float(ws.cell(row_idx, 4).value),
                    'current_carbon': safe_float(ws.cell(row_idx, 5).value),
                    'optimal_cost': safe_float(ws.cell(row_idx, 13).value),
                    'optimal_power': safe_float(ws.cell(row_idx, 14).value),
                    'optimal_carbon': safe_float(ws.cell(row_idx, 15).value),
                    'optimal_savings': safe_float(ws.cell(row_idx, 16).value),
                    'optimal_perf': ws.cell(row_idx, 17).value or 0,
                    'good_cost': safe_float(ws.cell(row_idx, 21).value),
                    'good_power': safe_float(ws.cell(row_idx, 22).value),
                    'good_carbon': safe_float(ws.cell(row_idx, 23).value),
                    'good_savings': safe_float(ws.cell(row_idx, 24).value),
                    'good_perf': ws.cell(row_idx, 25).value or 0,
                }
                break
            
            # Skip empty rows or note rows
            region = ws.cell(row_idx, 1).value
            if not region or region.startswith("Note"):
                continue
                
            instance_data = {
                'region': region,
                'current_instance': ws.cell(row_idx, 2).value,
                'current_cost': safe_float(ws.cell(row_idx, 3).value),
                'current_power': safe_float(ws.cell(row_idx, 4).value),
                'current_carbon': safe_float(ws.cell(row_idx, 5).value),
                'uuid': ws.cell(row_idx, 6).value,
                'cloud': ws.cell(row_idx, 7).value,
                'pricing_model': ws.cell(row_idx, 8).value,
                'vcpu': safe_float(ws.cell(row_idx, 9).value),
                'remark': ws.cell(row_idx, 10).value,
                'optimal_instance': ws.cell(row_idx, 11).value,
                'optimal_vcpu': safe_float(ws.cell(row_idx, 12).value),
                'optimal_cost': safe_float(ws.cell(row_idx, 13).value),
                'optimal_power': safe_float(ws.cell(row_idx, 14).value),
                'optimal_carbon': safe_float(ws.cell(row_idx, 15).value),
                'optimal_savings': safe_float(ws.cell(row_idx, 16).value),
                'optimal_perf': ws.cell(row_idx, 17).value or 0,
                'good_instance': ws.cell(row_idx, 19).value,
                'good_vcpu': safe_float(ws.cell(row_idx, 20).value),
                'good_cost': safe_float(ws.cell(row_idx, 21).value),
                'good_power': safe_float(ws.cell(row_idx, 22).value),
                'good_carbon': safe_float(ws.cell(row_idx, 23).value),
                'good_savings': safe_float(ws.cell(row_idx, 24).value),
                'good_perf': ws.cell(row_idx, 25).value or 0,
            }
            
            data['instances'].append(instance_data)
            data['unique_regions'].add(region)
        
        log_message(LevelType.INFO, f"Extracted {len(data['instances'])} instances from Excel")
        return data
        
    except Exception as e:
        log_message(LevelType.ERROR, f"Error extracting Excel data: {e}", ErrorCode=-1)
        raise


def replace_text_preserving_formatting(shape, new_text):
    """
    Replace text in a shape while preserving the formatting of the first run.
    
    Args:
        shape: The PowerPoint shape object
        new_text: The new text to insert
    """
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        return

    # Target the first paragraph
    p = text_frame.paragraphs[0]
    
    if not p.runs:
        p.add_run().text = new_text
        return

    # Set the text of the first run to the new text
    p.runs[0].text = new_text

    # Clear text from all subsequent runs in this paragraph
    for i in range(1, len(p.runs)):
        p.runs[i].text = ""
        
    # Clear all subsequent paragraphs completely
    # We loop through remaining paragraphs and clear their text
    # Note: We cannot easily 'remove' paragraphs in python-pptx without low-level XML manipulation
    # so we just empty their text.
    for i in range(1, len(text_frame.paragraphs)):
        for run in text_frame.paragraphs[i].runs:
            run.text = ""


def update_cloud_footprint(slide, data):
    """
    Update Cloud Footprint section with instance count, monthly spend, and regions.
    
    Args:
        slide: PowerPoint slide object
        data: Extracted Excel data
    """
    try:
        # Count instances with recommendations (remark is empty/None, similar to CCA logic)
        instances_analyzed = sum(1 for inst in data['instances'] if not inst.get('remark'))
        
        # Monthly spend (directly from current cost)
        monthly_spend = data['grand_total'].get('current_cost', 0)
        # Annualize if needed (User requested Annual Spend in Cloud Footprint)
        annual_spend = monthly_spend * 12
        
        # Unique regions count
        region_count = len(data['unique_regions'])
        
        # Format monthly spend using format_value function
        # Using Annual Spend as per request
        annual_spend_str = f"${format_value(annual_spend)}"
            
        # Update text shapes by targeting specific placeholders
        for shape in slide.shapes:
            if not hasattr(shape, 'text'):
                continue
                
            text = shape.text.strip()
            
            # Target "02" -> Instances Analyzed count
            if text == "<3a>":
                replace_text_preserving_formatting(shape, str(instances_analyzed).zfill(2))
                
            # Target "$2.5 K" or "$ 2.5 K" -> Annual Spend
            # The placeholder <3bF> is now for Annual Spend
            elif text == "<3bF>":
                replace_text_preserving_formatting(shape, annual_spend_str)
                
            # Target Summary text with region count
            elif "Summary : Current infrastructure across" in text:
                # Use regex to replace [X] with actual count
                new_text = re.sub(r'\[\d+\]', f'[{region_count}]', text)
                replace_text_preserving_formatting(shape, new_text)
        
        log_message(LevelType.INFO, f"Updated Cloud Footprint: {instances_analyzed} instances, {annual_spend_str} annual spend, {region_count} regions")
        
    except Exception as e:
        log_message(LevelType.ERROR, f"Error updating Cloud Footprint: {e}", ErrorCode=-1)
        raise




def style_chart_for_dark_theme(chart, bar_color_rgb=(0, 176, 240), is_percentage=False, font_size=6):
    """
    Apply dark theme styling to chart: white fonts, white lines, no gridlines, custom bar colors, no legend, data labels.
    
    Args:
        chart: Chart object to style
        bar_color_rgb: RGB tuple for bar color (default: cyan/turquoise)
        is_percentage: If True, format Y-axis as percentage (0-100), else use K format
        font_size: Font size for labels (default: 6)
    """
    try:
        # Remove chart title if it exists
        chart.has_title = False
        
        # Remove legend
        chart.has_legend = False
        
        # Style value axis (Y-axis)
        if hasattr(chart, 'value_axis'):
            value_axis = chart.value_axis
            value_axis.tick_labels.font.color.rgb = RGBColor(255, 255, 255)
            value_axis.tick_labels.font.size = Pt(font_size)
            # Remove gridlines
            value_axis.has_major_gridlines = False
            value_axis.has_minor_gridlines = False
            value_axis.format.line.color.rgb = RGBColor(255, 255, 255)
            # Format numbers based on chart type
            if is_percentage:
                # For percentage charts (0-100), use simple number format
                value_axis.tick_labels.number_format = '0'
            else:
                # Conditional format: M for millions, K for thousands
                # [>=1000000]#,##0.00,,"M";[>=1000]#,##0.00,"K";0.00
                value_axis.tick_labels.number_format = '[>=1000000]#,##0.00,,"M";[>=1000]#,##0.00,"K";0.00'
            
            # Ensure labels are at the bottom (Low)
            value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
            
        # Style category axis (X-axis)
        if hasattr(chart, 'category_axis'):
            category_axis = chart.category_axis
            category_axis.tick_labels.font.color.rgb = RGBColor(255, 255, 255)
            category_axis.tick_labels.font.size = Pt(font_size)
            category_axis.format.line.color.rgb = RGBColor(255, 255, 255)
        
        # Set bar colors
        for series in chart.series:
            # Set bar colors
            for point in series.points:
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*bar_color_rgb)
            
            # Add data labels manually per point for B/M/K support
            try:
                # Activate at Plot level to ensure visibility
                plot = chart.plots[0]
                plot.has_data_labels = True
                
                series.has_data_labels = True
                for i, point in enumerate(series.points):
                    data_label = point.data_label
                    data_label.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    
                    # Get the actual value
                    val = series.values[i]
                    
                    # Manual formatting
                    if is_percentage:
                        label_text = f"{val:.2f}%"
                    else:
                        if val >= 1_000_000_000:
                            label_text = f"{val/1_000_000_000:.2f}B"
                        elif val >= 1_000_000:
                            label_text = f"{val/1_000_000:.2f}M"
                        elif val >= 1_000:
                            label_text = f"{val/1_000:.2f}K"
                        else:
                            label_text = f"{val:.2f}"
                    
                    # Set the text directly
                    data_label.text_frame.text = label_text
                    
                    # Force run-level styling to avoid reset to black/default
                    for paragraph in data_label.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(7)
                            run.font.color.rgb = RGBColor(255, 255, 255)
            except Exception as e:
                log_message(LevelType.WARNING, f"Could not manually set data labels: {e}")
            
    except Exception as e:
        log_message(LevelType.WARNING, f"Could not fully style chart: {e}")


def update_business_value(slide, data):
    """
    Update Business Value section with dynamic chart showing Current/Optimal/Good costs.
    Removes old static shapes and inserts a new chart.
    
    Args:
        slide: PowerPoint slide object
        data: Extracted Excel data
    """
    try:
        # Annualize costs for the chart
        current_cost = data['grand_total'].get('current_cost', 0) * 12
        optimal_cost = data['grand_total'].get('optimal_cost', 0) * 12
        good_cost = data['grand_total'].get('good_cost', 0) * 12
        
        # Remove old shapes in the Business Value chart area - EXTREMELY AGGRESSIVE
        # Need to remove Y-axis text far on the left side
        # Based on inspection: title at left=3205163, top=609600
        # Expand FAR left to catch all Y-axis labels: left 2800000-6000000, top 600000-2600000
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                # Check if shape is in the Business Value chart area (including far left Y-axis)
                if (2800000 <= shape.left <= 6000000 and 
                    900000 <= shape.top <= 2900000):
                    # PRESERVE TITLES AND BACKGROUND BOXES
                    if hasattr(shape, 'text'):
                        # Keep the title
                        if 'Business Value' in shape.text:
                            continue  
                    
                    # Keep background boxes (AutoShapes without text or empty text)
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                         if not shape.has_text_frame or not shape.text.strip():
                            # Check dimensions to distinguish background from chart elements
                            # RELAXED FILTER: Keep if reasonably large (Width > 1,000,000)
                            if shape.width > 1000000:
                                continue

                    # Remove everything else (Chart content, labels, small shapes)
                    shapes_to_remove.append(shape)
        
        # Remove shapes (must be done separately to avoid iterator issues)
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
        
        # Define chart data
        chart_data = CategoryChartData()
        chart_data.categories = ['Current', 'Optimal', 'Good']
        chart_data.add_series('Cost ($)', (current_cost, optimal_cost, good_cost))
        
        # Insert new chart - moved down more
        # Insert new chart - moved down MORE
        # Y position increased to 1350000
        x, y, cx, cy = 3250000, 1350000, 2700000, 1350000
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart
        
        # Apply dark theme styling with cyan color and no legend
        style_chart_for_dark_theme(chart, bar_color_rgb=(0, 176, 240))
        
        log_message(LevelType.INFO, f"Inserted Business Value chart: Current=${current_cost}, Optimal=${optimal_cost}, Good=${good_cost}")
        
    except Exception as e:
        log_message(LevelType.ERROR, f"Error updating Business Value: {e}", ErrorCode=-1)
        raise


def update_path_to_savings(slide, data):
    """
    Update Path to Savings section with instance types and savings from first row + Grand Total.
    
    Args:
        slide: PowerPoint slide object
        data: Extracted Excel data
    """
    try:
        # Get first row data for instance types
        first_instance = data['instances'][0] if data['instances'] else {}
        
        current_instance = first_instance.get('current_instance', '')
        optimal_instance = first_instance.get('optimal_instance', '')
        good_instance = first_instance.get('good_instance', '')
        
        # Get savings and perf from Grand Total (Monthly values)
        optimal_savings_monthly = data['grand_total'].get('optimal_savings', 0)
        optimal_perf = data['grand_total'].get('optimal_perf', 0)
        good_savings_monthly = data['grand_total'].get('good_savings', 0)
        good_perf = data['grand_total'].get('good_perf', 0)
        
        # Annualize Savings
        optimal_savings_annual = optimal_savings_monthly * 12
        good_savings_annual = good_savings_monthly * 12
        
        log_message(LevelType.INFO, f"Path to Savings - Optimal Perf: {optimal_perf}, Good Perf: {good_perf}")
        
        # Format values with K/M suffixes
        optimal_savings_str = f"$ {format_value(optimal_savings_annual)}"
        good_savings_str = f"$ {format_value(good_savings_annual)}"
        
        optimal_perf_str = f"{format_value(optimal_perf)} X"
        good_perf_str = f"{format_value(good_perf)} X"
        
        # Update text shapes by targeting specific placeholders
        for shape in slide.shapes:
            if not hasattr(shape, 'text'):
                continue
                
            text = shape.text.strip()
            
            # OPTIMAL SECTION PLACEHOLDERS
            if text == "ph1":  # Current Instance placeholder (Optimal section)
                # Note: This placeholder appears twice (Optimal and Good sections).
                # We need to distinguish based on position (Top vs Bottom)
                # if shape.top < 5000000:  # Top half (Optimal)
                #     replace_text_preserving_formatting(shape, current_instance)
                # else:
                replace_text_preserving_formatting(shape, current_instance)
                
            elif text == "ph2":  # Optimal Instance placeholder
                replace_text_preserving_formatting(shape, optimal_instance)
                
            elif text == "<3os>":  # Optimal Savings placeholder
                replace_text_preserving_formatting(shape, optimal_savings_str)
                
            elif text == "<3op>":  # Optimal Perf placeholder
                replace_text_preserving_formatting(shape, optimal_perf_str)
                # Reduce font size to fit
                if shape.has_text_frame and shape.text_frame.paragraphs:
                    shape.text_frame.paragraphs[0].runs[0].font.size = Pt(5.3)
                
            # GOOD SECTION PLACEHOLDERS
            # elif text == "ph1":  # Current Instance placeholder (Good section)
            #     replace_text_preserving_formatting(shape, current_instance)
                
            elif text == "ph3":  # Good Instance placeholder
                replace_text_preserving_formatting(shape, good_instance)
                
            elif text == "<3gs>":  # Good Savings placeholder
                replace_text_preserving_formatting(shape, good_savings_str)
                
            elif text == "<3gp>":  # Good Perf placeholder
                replace_text_preserving_formatting(shape, good_perf_str)
                # Reduce font size to fit
                if shape.has_text_frame and shape.text_frame.paragraphs:
                    shape.text_frame.paragraphs[0].runs[0].font.size = Pt(5.3)
        
        log_message(LevelType.INFO, "Updated Path to Savings section")
        
    except Exception as e:
        log_message(LevelType.ERROR, f"Error updating Path to Savings: {e}", ErrorCode=-1)
        raise


def update_dollar_spend_distribution(slide, data):
    """
    Update Dollar Spend Distribution section with: 
    1st highest cost instance, Next 10 (sum), Rest (sum).
    Always shows 3 bars even if values are 0.
    Removes old static shapes and inserts a new chart.
    
    Args:
        slide: PowerPoint slide object
        data: Extracted Excel data
    """
    try:
        # Group instances by type and sum their costs
        instance_costs = {}
        for inst in data['instances']:
            name = inst.get('current_instance', 'Unknown')
            cost = inst.get('current_cost', 0)
            instance_costs[name] = instance_costs.get(name, 0) + cost
            
        # Create list of dicts for sorting
        grouped_instances = [{'current_instance': k, 'current_cost': v} for k, v in instance_costs.items()]
        
        # Sort instances by current cost descending
        sorted_instances = sorted(grouped_instances, key=lambda x: x['current_cost'], reverse=True)
        
        total_cost = data['grand_total'].get('current_cost', 0)
        
        # ALWAYS create 3 categories: Top 1, Next 10, Rest
        # First bar: Top 1 (largest instance group)
        if sorted_instances:
            top_instance = sorted_instances[0]
            top_percentage = round((top_instance['current_cost'] / total_cost * 100), 2) if total_cost > 0 else 0
            top_name = top_instance['current_instance']
        else:
            top_percentage = 0
            top_name = "N/A"
        
        # Second bar: Next 10 (groups 2-11)
        if len(sorted_instances) > 1:
            next_10_instances = sorted_instances[1:11]
            next_10_cost = sum(inst['current_cost'] for inst in next_10_instances)
            next_10_percentage = round((next_10_cost / total_cost * 100), 2) if total_cost > 0 else 0
        else:
            next_10_percentage = 0
        
        # Third bar: Rest (groups 12+)
        if len(sorted_instances) > 11:
            rest_instances = sorted_instances[11:]
            rest_cost = sum(inst['current_cost'] for inst in rest_instances)
            rest_percentage = round((rest_cost / total_cost * 100), 2) if total_cost > 0 else 0
        else:
            rest_percentage = 0
        
        # Remove old shapes in the Dollar Spend chart area - EXTREMELY AGGRESSIVE
        # Expand area to catch all Y-axis labels on the FAR left (100, 75, 50, 25, 0)
        # Based on inspection: title at left=185738, top=2805113
        # Chart area: left -100000 to 3200000, top 2800000-4900000
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                # Catch shapes from far left (including Y-axis labels) to right
                if (shape.left <= 3200000 and 
                    3100000 <= shape.top <= 4900000):
                    # PRESERVE TITLES AND BACKGROUND BOXES
                    if hasattr(shape, 'text'):
                        # Keep the title
                        if 'Dollar Spend' in shape.text:
                            continue  
                    
                    # Keep background boxes (AutoShapes without text or empty text)
                    # FILTER BY SIZE: Keep only large backgrounds
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                         if not shape.has_text_frame or not shape.text.strip():
                            if shape.width > 1000000 and shape.height > 1000000:
                                continue
                            
                            # Preserve header background
                            if shape.top < 2900000: # Section starts at 2800000
                                continue

                    # Remove everything else
                    shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
        
        # Define chart data - ALWAYS 3 categories
        chart_data = CategoryChartData()
        chart_data.categories = [top_name, 'Next 10', 'Rest']
        chart_data.add_series('Percentage', (top_percentage, next_10_percentage, rest_percentage))
        
        # Insert new chart - MOVED DOWN MORE within the section
        # Position adjusted: y increased to move down, staying within bounds
        x, y, cx, cy = 250000, 3500000, 2700000, 1350000
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart
        
        # Apply dark theme styling with cyan color - PERCENTAGE FORMAT for Y-axis
        style_chart_for_dark_theme(chart, bar_color_rgb=(0, 176, 240), is_percentage=True)
        
        # Set Y-axis to percentage format (0-100)
        if hasattr(chart, 'value_axis'):
            chart.value_axis.maximum_scale = 100
            chart.value_axis.minimum_scale = 0
        
        log_message(LevelType.INFO, "Inserted Dollar Spend Distribution chart with 3 bars")
        
    except Exception as e:
        log_message(LevelType.ERROR, f"Error updating Dollar Spend Distribution: {e}", ErrorCode=-1)
        raise


def update_power_carbon(slide, data):
    """
    Update Power/Carbon section with dynamic charts for Current/Optimal/Good.
    Also updates the detailed text blocks with calculated savings.
    
    Args:
        slide: PowerPoint slide object
        data: Extracted Excel data
    """
    try:
        current_power = data['grand_total'].get('current_power', 0)
        optimal_power = data['grand_total'].get('optimal_power', 0)
        good_power = data['grand_total'].get('good_power', 0)
        
        current_carbon = data['grand_total'].get('current_carbon', 0)
        optimal_carbon = data['grand_total'].get('optimal_carbon', 0)
        good_carbon = data['grand_total'].get('good_carbon', 0)
        
        current_cost = data['grand_total'].get('current_cost', 0)
        optimal_cost = data['grand_total'].get('optimal_cost', 0)
        good_cost = data['grand_total'].get('good_cost', 0)
        
        # Calculate savings percentages (clamp negative values to 0)
        opt_cost_save_pct = round(((current_cost - optimal_cost) / current_cost * 100), 2) if current_cost > 0 else 0
        opt_power_save_pct = round(((current_power - optimal_power) / current_power * 100), 2) if current_power > 0 else 0
        opt_carbon_save_pct = round(((current_carbon - optimal_carbon) / current_carbon * 100), 2) if current_carbon > 0 else 0
        
        good_cost_save_pct = round(((current_cost - good_cost) / current_cost * 100), 2) if current_cost > 0 else 0
        good_power_save_pct = round(((current_power - good_power) / current_power * 100), 2) if current_power > 0 else 0
        good_carbon_save_pct = round(((current_carbon - good_carbon) / current_carbon * 100), 2) if current_carbon > 0 else 0
        
        # Remove old shapes in the Power/Carbon chart area
        # Based on inspection: title at left=3205163, top=2819400
        # Chart area extended to bottom to catch all extra text
        # Chart area: left 3000000-9000000 (Expanded left and right), top 3000000-5100000
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                if (3000000 <= shape.left <= 9000000 and 
                    3100000 <= shape.top <= 5100000):
                    # PRESERVE TITLES AND BACKGROUND BOXES
                    if hasattr(shape, 'text'):
                        # Keep the section title - check robustly
                        if 'Power' in shape.text and 'Carbon' in shape.text:
                            continue
                        # Keep the detail text boxes (contain Cost: and lower cost)
                        if 'Cost:' in shape.text and 'lower cost' in shape.text:
                            continue
                        # Keep Optimal and Good boxes (headers) - ONLY in the detail area (Right side)
                        if ('Optimal' in shape.text or 'Good' in shape.text) and shape.left > 7000000:
                            continue
                    
                    # Keep background boxes (AutoShapes without text or empty text)
                    # FILTER BY SIZE: Keep only large backgrounds
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                         if not shape.has_text_frame or not shape.text.strip():
                            if shape.width > 1000000 and shape.height > 1000000:
                                continue
                            
                            # Preserve header background
                            if shape.top < 3100000: # Section starts at 3000000
                                continue
                            
                            # Preserve Background Boxes for Optimal/Good Details (Right side)
                            # These might be smaller than the 1M x 1M filter above
                            # Detail area is roughly Left > 7.0M
                            if shape.left > 7000000:
                                continue

                    # Remove everything else (including old chart elements)
                    shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
        
        # --- CHART 1: POWER (Left) ---
        power_chart_data = CategoryChartData()
        # For Bar Chart (Horizontal), categories are bottom-to-top.
        # Template visual: Top=Current, Middle=Optimal, Bottom=Good.
        power_chart_data.categories = ['Current', 'Optimal', 'Good']
        power_chart_data.add_series('Power (kW)', (current_power, optimal_power, good_power))
        
        # Position: Left side of the container
        # Increased height slightly for better visibility
        x1, y1, cx1, cy1 = 3200000, 3300000, 2000000, 1150000
        graphic_frame1 = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x1, y1, cx1, cy1, power_chart_data
        )
        chart1 = graphic_frame1.chart
        style_chart_for_dark_theme(chart1, bar_color_rgb=(255, 153, 51), font_size=5) # Orange

        # Add Manual Title for Power
        # Moved up again to avoid overlap with boxes
        tx_box_1 = slide.shapes.add_textbox(3200000, 4320000, 2200000, 250000)
        tf_1 = tx_box_1.text_frame
        tf_1.text = "Power (kW)"
        tf_1.paragraphs[0].font.size = Pt(8)
        tf_1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf_1.paragraphs[0].alignment = 2 # Centered

        # Sustainability Score Box for Power
        # Move up to Y = 4.52M
        try:
            score_box_y = 4520000
            score_box_cx = 1950000
            score_box_cy = 180000
            shape1 = slide.shapes.add_shape(
                1, 3250000, score_box_y, score_box_cx, score_box_cy
            )
            shape1.fill.solid()
            shape1.fill.fore_color.rgb = RGBColor(25, 25, 25) # Slightly darker for "transparent dark" look
            shape1.line.color.rgb = RGBColor(89, 89, 89) # Subtler border
            shape1.line.width = Pt(0.5)
            
            tf_score1 = shape1.text_frame
            tf_score1.word_wrap = False
            tf_score1.margin_top = 0
            tf_score1.margin_bottom = 0
            tf_score1.margin_left = 0
            tf_score1.margin_right = 0
            p1 = tf_score1.paragraphs[0]
            p1.alignment = 1 # Center
            
            # Add info icon (ⓘ)
            run_i = p1.add_run()
            run_i.text = " ⓘ "
            run_i.font.size = Pt(6)
            run_i.font.color.rgb = RGBColor(180, 180, 180)
            
            # Add score text (Multi-run for conditional coloring)
            run_lbl = p1.add_run()
            # 2 Decimal Places Max as requested
            run_lbl.text = f"Sustainability Score : Optimal {opt_power_save_pct:.2f}% "
            run_lbl.font.size = Pt(5)
            run_lbl.font.bold = True
            run_lbl.font.color.rgb = RGBColor(255, 255, 255)

            # Optimal Arrow
            run_a1 = p1.add_run()
            run_a1.text = "↑" if opt_power_save_pct >= 0 else "↓"
            run_a1.font.size = Pt(5)
            run_a1.font.bold = True
            run_a1.font.color.rgb = RGBColor(0, 176, 80) if opt_power_save_pct >= 0 else RGBColor(255, 0, 0)

            run_lbl2 = p1.add_run()
            run_lbl2.text = f" , Good {good_power_save_pct:.2f}% "
            run_lbl2.font.size = Pt(5)
            run_lbl2.font.bold = True
            run_lbl2.font.color.rgb = RGBColor(255, 255, 255)

            # Good Arrow
            run_a2 = p1.add_run()
            run_a2.text = "↑" if good_power_save_pct >= 0 else "↓"
            run_a2.font.size = Pt(5)
            run_a2.font.bold = True
            run_a2.font.color.rgb = RGBColor(0, 176, 80) if good_power_save_pct >= 0 else RGBColor(255, 0, 0)
        except Exception as e:
            log_message(LevelType.WARNING, f"Could not add Power sustainability score box: {e}")

        
        # Manually color points for Power Chart
        try:
            series = chart1.series[0]
            points = series.points
            points[0].format.fill.solid()
            points[0].format.fill.fore_color.rgb = RGBColor(192, 0, 0)
            points[1].format.fill.solid()
            points[1].format.fill.fore_color.rgb = RGBColor(255, 192, 0)
            points[2].format.fill.solid()
            points[2].format.fill.fore_color.rgb = RGBColor(255, 192, 0)
            
            if hasattr(chart1, 'category_axis'):
                chart1.category_axis.reverse_order = True
            if hasattr(chart1, 'value_axis'):
                chart1.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        except Exception as e:
            log_message(LevelType.WARNING, f"Could not manually color Power chart points: {e}")

        
        # --- CHART 2: CARBON (Right) ---
        carbon_chart_data = CategoryChartData()
        carbon_chart_data.categories = ['Current', 'Optimal', 'Good']
        carbon_chart_data.add_series('Carbon (kgCO2eq)', (current_carbon, optimal_carbon, good_carbon))
        
        # --- VERTICAL DIVIDER LINE ---
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, 5250000, 3150000, 5250000, 4750000
        )
        connector.line.color.rgb = RGBColor(128, 128, 128) # Grey
        connector.line.width = Pt(1)

        #Position: Right side of the container
        # Increased height slightly
        x2, y2, cx2, cy2 = 5300000, 3300000, 2000000, 1150000
        graphic_frame2 = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x2, y2, cx2, cy2, carbon_chart_data
        )
        chart2 = graphic_frame2.chart
        style_chart_for_dark_theme(chart2, bar_color_rgb=(146, 208, 80), font_size=5) # Green

        # Add Manual Title for Carbon
        # Moved up again
        tx_box_2 = slide.shapes.add_textbox(5300000, 4320000, 2200000, 250000)
        tf_2 = tx_box_2.text_frame
        tf_2.text = "Carbon (kgCO2eq)"
        tf_2.paragraphs[0].font.size = Pt(8)
        tf_2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf_2.paragraphs[0].alignment = 2 # Centered

        # Sustainability Score Box for Carbon
        try:
            shape2 = slide.shapes.add_shape(
                1, 5325000, score_box_y, score_box_cx, score_box_cy
            )
            shape2.fill.solid()
            shape2.fill.fore_color.rgb = RGBColor(25, 25, 25)
            shape2.line.color.rgb = RGBColor(89, 89, 89)
            shape2.line.width = Pt(0.5)
            
            tf_score2 = shape2.text_frame
            tf_score2.word_wrap = False
            tf_score2.margin_top = 0
            tf_score2.margin_bottom = 0
            tf_score2.margin_left = 0
            tf_score2.margin_right = 0
            p2 = tf_score2.paragraphs[0]
            p2.alignment = 1
            
            # Add info icon (ⓘ)
            run_i2 = p2.add_run()
            run_i2.text = " ⓘ "
            run_i2.font.size = Pt(6)
            run_i2.font.color.rgb = RGBColor(180, 180, 180)
            
            # Add score text (Multi-run for conditional coloring)
            run_lbl_c = p2.add_run()
            run_lbl_c.text = f"Sustainability Score : Optimal {opt_carbon_save_pct:.2f}% "
            run_lbl_c.font.size = Pt(5)
            run_lbl_c.font.bold = True
            run_lbl_c.font.color.rgb = RGBColor(255, 255, 255)

            # Optimal Arrow
            run_ca1 = p2.add_run()
            run_ca1.text = "↑" if opt_carbon_save_pct >= 0 else "↓"
            run_ca1.font.size = Pt(5)
            run_ca1.font.bold = True
            run_ca1.font.color.rgb = RGBColor(0, 176, 80) if opt_carbon_save_pct >= 0 else RGBColor(255, 0, 0)

            run_lbl_c2 = p2.add_run()
            run_lbl_c2.text = f" , Good {good_carbon_save_pct:.2f}% "
            run_lbl_c2.font.size = Pt(5)
            run_lbl_c2.font.bold = True
            run_lbl_c2.font.color.rgb = RGBColor(255, 255, 255)

            # Good Arrow
            run_ca2 = p2.add_run()
            run_ca2.text = "↑" if good_carbon_save_pct >= 0 else "↓"
            run_ca2.font.size = Pt(5)
            run_ca2.font.bold = True
            run_ca2.font.color.rgb = RGBColor(0, 176, 80) if good_carbon_save_pct >= 0 else RGBColor(255, 0, 0)
        except Exception as e:
            log_message(LevelType.WARNING, f"Could not add Carbon sustainability score box: {e}")

        # --- Calculation Formula at bottom ---
        # Moved up to Y = 4.75M
        try:
            formula_box = slide.shapes.add_textbox(3250000, 4750000, 4000000, 200000)
            tf_f = formula_box.text_frame
            tf_f.paragraphs[0].alignment = 1 # Centered
            run_f = tf_f.paragraphs[0].add_run()
            run_f.text = "Calculation : sustainability Improvement (%) = (Current – Target) / Current x 100"
            run_f.font.size = Pt(5.5)
            run_f.font.bold = True
            run_f.font.color.rgb = RGBColor(0, 176, 240) # Cyan
        except Exception as e:
            log_message(LevelType.WARNING, f"Could not add formula text: {e}")

        
        # Manually color points for Carbon Chart
        try:
            series = chart2.series[0]
            points = series.points
            points[0].format.fill.solid()
            points[0].format.fill.fore_color.rgb = RGBColor(192, 0, 0)
            points[1].format.fill.solid()
            points[1].format.fill.fore_color.rgb = RGBColor(146, 208, 80)
            points[2].format.fill.solid()
            points[2].format.fill.fore_color.rgb = RGBColor(146, 208, 80)
            
            if hasattr(chart2, 'category_axis'):
                chart2.category_axis.reverse_order = True
            if hasattr(chart2, 'value_axis'):
                chart2.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        except Exception as e:
            log_message(LevelType.WARNING, f"Could not manually color Carbon chart points: {e}")
            
        
        log_message(LevelType.INFO, "Inserted separate Power and Carbon charts with Sustainability Scores")
        

        # 2. Update detailed text blocks
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.has_text_frame:
                text = shape.text
                
                # OPTIMAL DETAILS
                if "Cost:" in text and "lower cost" in text and shape.top < 4000000:
                    try:
                        tf = shape.text_frame
                        tf.text = ""  # Clear existing
                        tf.margin_top = 0  # Absolute top
                        tf.margin_left = 0  # Absolute left
                        tf.margin_right = 0
                        tf.margin_bottom = 0
                        tf.vertical_anchor = MSO_ANCHOR.TOP
                        # Details (5.1pt Normal, Indented 4 spaces)
                        def add_detail(paragraph, label, value_str, space_after=Pt(3.5)):
                            paragraph.space_after = space_after
                            paragraph.space_before = 0
                            paragraph.line_spacing = 1.0
                            run = paragraph.add_run()
                            run.text = f"    {label} {value_str}"
                            run.font.size = Pt(5.1)
                            run.font.bold = False
                            run.font.color.rgb = RGBColor(255, 255, 255)
                        
                        # Use existing first paragraph for the first line
                        p0 = tf.paragraphs[0]
                        add_detail(p0, "Cost:", f"${format_value(optimal_cost * 12)} ({opt_cost_save_pct:.2f}% lower cost)")
                        
                        p1 = tf.add_paragraph()
                        add_detail(p1, "Power:", f"{format_value(optimal_power)} kW ({opt_power_save_pct:.2f}% lower power)")
                        
                        # Carbon split into two lines
                        p2 = tf.add_paragraph()
                        add_detail(p2, "Carbon:", f"{format_value(optimal_carbon)} kgCO₂eq ({opt_carbon_save_pct:.2f}% lower", space_after=0)
                        
                        p3 = tf.add_paragraph()
                        add_detail(p3, "", "carbon emission)", space_after=Pt(3.5))

                    except Exception as e:
                        log_message(LevelType.ERROR, f"Error rebuilding Optimal details text: {e}", ErrorCode=-1)
                        
                # GOOD DETAILS
                elif "Cost:" in text and "lower cost" in text and shape.top > 4000000:
                    try:
                        tf = shape.text_frame
                        tf.text = ""  # Clear existing
                        tf.margin_top = 0  # Absolute top
                        tf.margin_left = 0  # Absolute left
                        tf.margin_right = 0
                        tf.margin_bottom = 0
                        tf.vertical_anchor = MSO_ANCHOR.TOP
                        # Details (5.1pt Normal, Indented 4 spaces)
                        def add_detail_good(paragraph, label, value_str, space_after=Pt(3.5)):
                            paragraph.space_after = space_after
                            paragraph.space_before = 0
                            paragraph.line_spacing = 1.0
                            run = paragraph.add_run()
                            run.text = f"    {label} {value_str}"
                            run.font.size = Pt(5.1)
                            run.font.bold = False
                            run.font.color.rgb = RGBColor(255, 255, 255)
                        
                        # Use existing first paragraph for the first line
                        p0 = tf.paragraphs[0]
                        add_detail_good(p0, "Cost:", f"${format_value(good_cost * 12)} ({good_cost_save_pct:.2f}% lower cost)")
                        
                        p1 = tf.add_paragraph()
                        add_detail_good(p1, "Power:", f"{format_value(good_power)} kW ({good_power_save_pct:.2f}% lower power)")
                        
                        # Carbon split into two lines
                        p2 = tf.add_paragraph()
                        add_detail_good(p2, "Carbon:", f"{format_value(good_carbon)} kgCO₂eq ({good_carbon_save_pct:.2f}% lower", space_after=0)
                        
                        p3 = tf.add_paragraph()
                        add_detail_good(p3, "", "carbon emission)", space_after=Pt(3.5))

                    except Exception as e:
                        log_message(LevelType.ERROR, f"Error rebuilding Good details text: {e}", ErrorCode=-1)
        
    except Exception as e:
        log_message(LevelType.ERROR, f"Error updating Power/Carbon: {e}", ErrorCode=-1)
        raise
        



def update_header_info(slide, data, user_email):
    """
    Update header information: Cloud, Customer, and Date.
    
    Args:
        slide: PowerPoint slide object
        data: Extracted Excel data
        user_email: User's email address to extract customer name
    """
    try:
        # Get Cloud provider from first instance
        first_instance = data['instances'][0] if data['instances'] else {}
        cloud_provider = first_instance.get('cloud', 'Unknown')
        
        # Extract customer name from email (domain part, before .com)
        # e.g., pasaikum@amd.com -> amd
        if user_email and '@' in user_email:
            domain = user_email.split('@')[1]
            customer_name = domain.split('.')[0].upper()
        else:
            customer_name = "CLIENT"
        # Calculate annualized values for Summary
        current_cost = data['grand_total'].get('current_cost', 0)
        optimal_cost = data['grand_total'].get('optimal_cost', 0)
        good_cost = data['grand_total'].get('good_cost', 0)
        
        current_annual = current_cost * 12
        optimal_annual = optimal_cost * 12
        good_annual = good_cost * 12
        
        # Calculate savings percentages
        opt_cost_save_pct = round(((current_cost - optimal_cost) / current_cost * 100), 2) if current_cost > 0 else 0
        good_cost_save_pct = round(((current_cost - good_cost) / current_cost * 100), 2) if current_cost > 0 else 0

        # Power/Carbon savings
        current_power = data['grand_total'].get('current_power', 0)
        optimal_power = data['grand_total'].get('optimal_power', 0)
        good_power = data['grand_total'].get('good_power', 0)
        
        current_carbon = data['grand_total'].get('current_carbon', 0)
        optimal_carbon = data['grand_total'].get('optimal_carbon', 0)
        good_carbon = data['grand_total'].get('good_carbon', 0)
        
        opt_power_save_pct = round(((current_power - optimal_power) / current_power * 100), 2) if current_power > 0 else 0
        # good_power_save_pct = max(0, ((current_power - good_power) / current_power * 100)) if current_power > 0 else 0 # Not used in summary text
        
        opt_carbon_save_pct = round(((current_carbon - optimal_carbon) / current_carbon * 100), 2) if current_carbon > 0 else 0
        
        # Get today's date
        today_date = datetime.today().strftime("%B %d, %Y")
        
        # Update text shapes
        for shape in slide.shapes:
            if not hasattr(shape, 'text'):
                continue
                
            text = shape.text.strip()
            
            # Update Cloud
            if "Cloud:" in text:
                # Replace "Cloud: AWS" or similar with "Cloud: [Provider]"
                # We use regex to be safe, or just replace the known template string
                if "Cloud: AWS" in text:
                    replace_text_preserving_formatting(shape, f"Cloud: {cloud_provider}")
                else:
                    # Fallback: try to replace just the value if possible, or the whole string
                    # For now, let's assume the template has "Cloud: AWS" as seen in screenshot
                    pass
            
            # Update Customer
            elif "Customer:" in text:
                if "Customer: INFOBELLIT" in text:
                    replace_text_preserving_formatting(shape, f"Customer: {customer_name}")
            
            # Update Date
            elif "Date:" in text:
                # The template has "Date: October 23, 2025"
                # We'll try to match "Date:" and replace the whole thing
                if "Date:" in text:
                    replace_text_preserving_formatting(shape, f"Date: {today_date}")
            
            # Update Executive Summary Text
            elif "Modernizing" in text and "infrastructure can reduce annual spend" in text:
                # Dynamic Summary Text
                summary_text = (
                    f"Modernizing AMD EPYC infrastructure can reduce annual spend from ${format_value(current_annual)} (Current) "
                    f"to ${format_value(optimal_annual)} (Optimal) with {opt_cost_save_pct}% savings per year "
                    f"or ${format_value(good_annual)} (Good) {good_cost_save_pct}% savings per year, "
                    f"delivering significant cost efficiency and sustainability benefits. All this while reducing "
                    f"power consumption by {opt_power_save_pct}% and reducing carbon emissions by {opt_carbon_save_pct}%"
                )
                replace_text_preserving_formatting(shape, summary_text)
                     
        log_message(LevelType.INFO, f"Updated Header: Cloud={cloud_provider}, Customer={customer_name}, Date={today_date}")

    except Exception as e:
        log_message(LevelType.ERROR, f"Error updating Header info: {e}", ErrorCode=-1)
        # Don't raise, just log error as this is non-critical


def generate_ppt_from_excel(excel_path, template_path, output_path, user_email, password=None):
    """
    Main function to generate PowerPoint presentation from Excel data.
    
    Args:
        excel_path: Path to the Excel file
        template_path: Path to the PowerPoint template
        output_path: Path where the output PPT will be saved
        user_email: User's email address for customer name extraction
    """
    try:
        log_message(LevelType.INFO, f"Starting PPT generation from {excel_path}")
        
        # Extract data from Excel
        data = extract_excel_data(excel_path)
        
        # Load PowerPoint template
        log_message(LevelType.INFO, f"Loading PPT template: {template_path}")
        prs = Presentation(template_path)
        
        # Get the first slide (assuming single slide template)
        slide = prs.slides[0]
        
        # Update all sections
        update_header_info(slide, data, user_email)
        update_cloud_footprint(slide, data)
        update_business_value(slide, data)
        update_path_to_savings(slide, data)
        update_dollar_spend_distribution(slide, data)
        update_power_carbon(slide, data)
        
        # Save the presentation
        prs.save(output_path)
        # --------- Apply password protection if password provided ----------
        if password:
            protect_pptx(output_path, password)
            log_message(LevelType.INFO, f"ppt protected for {output_path}", ErrorCode=1)
        log_message(LevelType.INFO, f"PowerPoint successfully generated at: {output_path}")
        
    except Exception as e:
        log_message(LevelType.ERROR, f"Error generating PowerPoint: {e}", ErrorCode=-1)
        raise
