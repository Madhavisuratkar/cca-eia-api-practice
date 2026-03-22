import hashlib
from io import BytesIO
import subprocess
import csv
from datetime import datetime, date, timezone
import json
import math
import os
import logging
from enum import Enum
import re
import shutil
import uuid
from bson import ObjectId
import numpy as np
from pymongo import ASCENDING, MongoClient
from dotenv import load_dotenv
import pandas as pd
import h5py
import tempfile
from typing import List, Optional, Iterable
from msoffcrypto.format.ooxml import OOXMLFile
from openpyxl import Workbook
from openpyxl.styles import *
from openpyxl.utils import get_column_letter
from openpyxl.cell.cell import MergedCell
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Pt
from msal import ConfidentialClientApplication
import base64
import zipfile
import requests
import boto3
from concurrent_log_handler import ConcurrentRotatingFileHandler
from cryptography.fernet import Fernet
from typing import Any, Dict

from sqlalchemy import Column, BigInteger, String, Text, DateTime, Numeric, UniqueConstraint
from sqlalchemy.sql import func, or_
from sqlalchemy import create_engine, MetaData, text
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker
import sys
from sqlalchemy.exc import SQLAlchemyError, DBAPIError
import re
from openpyxl import load_workbook
from datetime import datetime
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR

env_path="/etc/cca_eia_secrets.env"
load_dotenv(env_path)

class LevelType(Enum):
    INFO = "INFO"
    WARNING = "WARNING"
    ERROR = "ERROR"
    DEBUG = "DEBUG"

def safe_round(value, decimals=2):
    try:
        if value is None or str(value).strip() == "":
            return value
        val = float(value)
        return round(val, decimals)
    except (ValueError, TypeError):
        return value


def format_savings(value, decimals=2):
    try:
        if value is None or str(value).strip() == "":
            return value
        val = float(value)
        if val < 0:
            return "EIA Recommended"
        return round(val, decimals)
    except (ValueError, TypeError):
        return value

class CollectionNames:
    PORTFOLIOS = "portfolios"
    CURRENT_INSTANCES = "current_instance"
    RECOMMENDED_INSTANCES = "recommended_instances"
    ENDPOINTS = "endpoints"
    HEALTH_CHECK = "health_check"
    RECCOMENDATION_TRACKING = "recommendation_tracking"
    INPUT_REMARKS = "input_remarks"
    NOTIFICATIONS = "notifications"
    RECOMMENDATION_ANALYTICS = "recommendation_analytics"
    ANALYTICS_WITHOUT_RECOMMENDATION = "recommendation_unsupported_analytics"
    INPUT_REMARKS = "input_remarks"

class RecommendationStatus:
    TO_PROCESS = "TO_PROCESS"
    QUEUE = "QUEUE"
    IN_PROGRESS = "IN_PROGRESS"
    COMPLETED = "COMPLETED"
    FAILED = "FAILED"

class LoggerManager:
    def __init__(self, name: str = None, level=logging.INFO, log_file: str = "bulk_data.log"):
        LOG_FORMAT = "%(asctime)s - %(levelname)s - %(message)s"

        # Full path for log file
        log_folder = os.path.join(os.path.dirname(os.path.realpath(__file__)), 'app', 'Logs')
        self.log_file = os.path.join(log_folder, log_file)

        # Create logger
        self.logger = logging.getLogger(name or self.__class__.__name__)
        self.logger.setLevel(level)

        # Avoid duplicate handlers
        if not self.logger.handlers:
            # Console handler
            ch = logging.StreamHandler()
            ch.setLevel(level)
            ch.setFormatter(logging.Formatter(LOG_FORMAT))
            self.logger.addHandler(ch)

            # File handler with rotation
            fh = ConcurrentRotatingFileHandler(
                self.log_file,
                maxBytes=50 * 1024 * 1024,  # 50 MB per file
                backupCount=20
            )
            fh.setLevel(level)
            fh.setFormatter(logging.Formatter(LOG_FORMAT))
            self.logger.addHandler(fh)

    def log_message(self, level: LevelType, message: str, ErrorCode: int = 0):
        """Log messages with level and optional error code."""
        msg = f"{message} | ErrorCode: {ErrorCode}" if ErrorCode else message
        if level == LevelType.INFO:
            self.logger.info(msg)
        elif level == LevelType.WARNING:
            self.logger.warning(msg)
        elif level == LevelType.ERROR:
            self.logger.error(msg)
        elif level == LevelType.DEBUG:
            self.logger.debug(msg)
        else:
            self.logger.info(msg)  

class AWSSecretManager(LoggerManager):
    """
    Handles decryption of AWS credentials using Fernet and fetching secrets
    from AWS Secrets Manager. Inherits from LoggerManager for logging.
    """

    def __init__(self):
        # Initialize logger
        super().__init__(name="AWSSecretManager", level=logging.INFO, log_file="aws_secret_manager.log")
        self.log_message(LevelType.INFO, "Initializing AWSSecretManager...")

        try:
            # Encrypted environment variables
            self.fernet_access_key = os.getenv("FERNET_ACCESS_KEY")
            self.fernet_secret_key = os.getenv("FERNET_SECRET_KEY")
            self.fernet_region_name = os.getenv("FERNET_REGION_NAME")
            self.fernet_secret_name = os.getenv("FERNET_SECRET_NAME")

            # Password used for Fernet key derivation
            key = base64.urlsafe_b64encode(hashlib.sha256("zenitsuagatsuma".encode()).digest())
            self.fernet = Fernet(key)

            # Decrypt AWS credentials
            self.AWS_ACCESS_KEY = self._decrypt(self.fernet_access_key)
            self.AWS_SECRET_KEY = self._decrypt(self.fernet_secret_key)
            self.REGION_NAME = self._decrypt(self.fernet_region_name)
            self.SECRET_NAME = self._decrypt(self.fernet_secret_name)

            # Initialize boto3 secrets client
            self.secrets_client = boto3.client(
                "secretsmanager",
                region_name=self.REGION_NAME,
                aws_access_key_id=self.AWS_ACCESS_KEY,
                aws_secret_access_key=self.AWS_SECRET_KEY,
            )

            self.log_message(LevelType.INFO, f"AWSSecretManager initialized successfully for region {self.REGION_NAME}")

        except Exception as e:
            self.log_message(LevelType.ERROR, f"Failed to initialize AWSSecretManager: {str(e)}", ErrorCode=-1)
            sys.exit(1)

    def _decrypt(self, encrypted_value: str) -> str:
        """Decrypts a Fernet-encrypted string."""
        if not encrypted_value:
            msg = "Missing encrypted value for Fernet decryption"
            self.log_message(LevelType.ERROR, msg, ErrorCode=-1)
            sys.exit(1)

        try:
            decrypted_value = self.fernet.decrypt(encrypted_value.encode()).decode()
            return decrypted_value
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Fernet decryption failed: {str(e)}", ErrorCode=-1)
            sys.exit(1)

    def fetch_secret_value(self, secret_name: str = None) -> str:
        """Fetch raw secret value (SecretString) from AWS Secrets Manager."""
        secret_name = secret_name or self.SECRET_NAME
        try:
            response = self.secrets_client.get_secret_value(SecretId=secret_name)
            secret_str = response.get("SecretString")
            if not secret_str:
                self.log_message(LevelType.ERROR, "No Secret String found in AWS Secrets Manager response", ErrorCode=-1)
                sys.exit(1)
            self.log_message(LevelType.INFO, f"Fetched secret: {secret_name}")
            return secret_str
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Failed to fetch secret value: {str(e)}", ErrorCode=-1)
            sys.exit(1)

    def parse_main_secret(self) -> Dict[str, Any]:
        """Parse the main secret and return its nested JSON."""
        try:
            secret_str = self.fetch_secret_value(self.SECRET_NAME)
            main_secret = json.loads(secret_str)
            if "secrets" not in main_secret:
                self.log_message(LevelType.ERROR, "'secrets' key not found in main secret payload", ErrorCode=-1)
                sys.exit(1)

            parsed = json.loads(main_secret["secrets"])
            self.log_message(LevelType.INFO, "Main secret parsed successfully")
            return parsed
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Failed to parse main secret: {str(e)}", ErrorCode=-1)
            sys.exit(1)

    def extract_secret_section(self, parsed_secret: Dict[str, Any], section_key: str) -> Dict[str, Any]:
        """Extract a specific section (like 'cca_secrets') from parsed secret."""
        section = parsed_secret.get(section_key)
        if not section:
            self.log_message(LevelType.ERROR, "'secrets' key not found", ErrorCode=-1)
            sys.exit(1)
        return section

secret_manager = AWSSecretManager()
parsed_secret = secret_manager.parse_main_secret()
secret_data = secret_manager.extract_secret_section(parsed_secret, "cca_secrets")
cs_secret_data = secret_manager.extract_secret_section(parsed_secret, "cs_secrets")

ROOT_DIR = os.getenv("ROOT_DIR")
CLIENT_ID = secret_data.get("AZURE_CLIENT_ID")
CLIENT_SECRET = secret_data.get("AZURE_CLIENT_SECRET")
TENANT_ID = secret_data.get("AZURE_TENANT_ID")
SENDER = os.getenv("AZURE_SENDER_EMAIL")
rec_failed_to_emails=os.getenv("rec_failed_to_emails")
AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
SCOPE = ["https://graph.microsoft.com/.default"]
GRAPH_ENDPOINT = f"https://graph.microsoft.com/v1.0/users/{SENDER}/sendMail"
NEGATIVE_REPLACEMENT_TEXT = "EIA Recommended"

h5_file_path = os.path.join(ROOT_DIR, 'database', 'pricing_database.h5')
with h5py.File(h5_file_path, 'r') as hdf:
    providers = list(hdf.keys())
    regions_map = {provider: list(hdf[provider].keys()) for provider in providers}
    instances_map = {
        (provider, region): set(hdf[f"{provider}/{region}"]['Instance'][:].astype(str))
        for provider, regions in regions_map.items()
        for region in regions
    }

# CS db set up and Policy engine Model
DB_PREFIX = "postgresql+psycopg2"

# Database connection settings from secrets
CS_DB_HOST = cs_secret_data.get("DB_HOST")
CS_DB_PORT = cs_secret_data.get("DB_PORT", 5432)
CS_DB_USER = cs_secret_data.get("DB_USERNAME")
CS_DB_PASSWORD = cs_secret_data.get("DB_PASSWORD")
CS_DB_NAME = cs_secret_data.get("DB_DATABASE")

# Full connection URL
CS_DB_URL = f"{DB_PREFIX}://{CS_DB_USER}:{CS_DB_PASSWORD}@{CS_DB_HOST}:{CS_DB_PORT}/{CS_DB_NAME}"


# SQLAlchemy base
Base = declarative_base()

# SQLAlchemy engine
try:
    cs_engine = create_engine(
        CS_DB_URL,
        pool_pre_ping=True,
        pool_size=10,
        max_overflow=20,
        echo=False  # change to True for SQL debug logging
    )
except (SQLAlchemyError, DBAPIError) as e:
    # handle known SQLAlchemy-related errors
    print(f"Database engine initialization error: {e}")
    sys.exit(1)
except Exception as e:
    # catch all other exceptions
    print(f"Unexpected error initializing DB engine: {e}")
    sys.exit(1)

# SQLAlchemy session maker
CSSessionLocal = sessionmaker(bind=cs_engine, autoflush=False, autocommit=False)

def check_cs_db_connection():
    """
    Check if PostgreSQL CS DB connection works (for script use).
    Returns True if OK, False if failed.
    """
    try:
        session = CSSessionLocal()
        session.execute(text("SELECT 1"))  # <-- FIXED
        session.close()
        return True
    except Exception as e:
        print("❌ CS DB connection failed:", str(e))
        return False

class PolicyEngine(Base):
    """
    Policy Engine database model
    """
    __tablename__ = 'policy_engine'

    id = Column(BigInteger, primary_key=True, autoincrement=True)
    user_email = Column(Text, nullable=False)
    provider = Column(Text, nullable=False)
    instance_type = Column(Text, nullable=False)
    scalar_value = Column(Numeric(10, 4), nullable=False)
    policy_name = Column(Text, nullable=False)
    policy_type = Column(Text, nullable=False)

    # Timestamps from DB server (timestamptz)
    created_at = Column(
        DateTime(timezone=True),
        server_default=func.now(),
        nullable=False
    )
    updated_at = Column(
        DateTime(timezone=True),
        server_default=func.now(),
        onupdate=func.now(),
        nullable=False
    )

    def to_dict(self):
        return {c.name: getattr(self, c.name) for c in self.__table__.columns}

# CCA - EIA
class PortfolioReader(LoggerManager):
    def __init__(self):
        super().__init__()  # initialize LoggerManager
        self.log_message(LevelType.INFO, "PortfolioReader initialized")

        # ----------------- MongoDB Config from ENV -----------------
        self.mongo_uri = secret_data.get("MONGO_URI")
        self.db_name = secret_data.get("DATABASE_NAME")
        self.ALLOWED_PROVIDERS = ["AWS", "AZURE", "GCP"]
        self.validatedata_path = f"{ROOT_DIR}/Validatedata.xlsx"
        self.MONTHLY_UTILIZATION = "monthly utilization (hourly)"
        self.results_path = os.getenv('RESULTS_PATH')
        self.ANNUAL_COST = "Annual Cost"
        self.ANNUAL_SAVINGS_I = "Annual Savings I"
        self.ANNUAL_SAVINGS_II = "Annual Savings II"
        self.ANNUAL_SAVINGS_III = "Annual Savings III"
        self.PERF_ENHANCEMENT_I = 'Perf Enhancement I'
        self.PERF_ENHANCEMENT_II = 'Perf Enhancement II'
        self.PERF_ENHANCEMENT_III = 'Perf Enhancement III'
        self.PRICE_MODEL = 'pricingModel'
        self.INSTANCE_TYPE = 'instance type'
        self.HEADROOM = 'headroom%'
        self.MAX_CPU = 'max cpu%'
        self.NO_DATA_IN_FILE = "No data available in provided file"
        self.REGION_REQUIRED = "Region is required"
        self.PRICING_MODEL = ['ondemand', 'reserved', 'spot']
        self.CLOUD_PROVIDERS = ["AWS", "AZURE", "GCP", "OCI", "DATADOG", "CLOUDWATCH", "AZUREINSIGHTS", "GCPTELEMETRY", "PROMETHEUS"]
        self.UNSUPPORTED_PROVIDERS = ["OCI"]
        self.INSTANCE_ERROR = "Instance type is required"
        self.PIPE = "|"
        self.azure_instance_pattern = r"^(standard|basic)_[A-Za-z]+\d+[a-z]*(_v\d+)?$"
        self.aws_instance_pattern = fr'\b(?:r7iz{self.PIPE}g4dn{self.PIPE}c7a{self.PIPE}mac2{self.PIPE}r6id{self.PIPE}c8g{self.PIPE}c5{self.PIPE}c7i{self.PIPE}r5{self.PIPE}c3{self.PIPE}m6i{self.PIPE}gr6{self.PIPE}f1{self.PIPE}mac2-m2{self.PIPE}c6in{self.PIPE}x8g{self.PIPE}is4gen{self.PIPE}u-9tb1{self.PIPE}i3{self.PIPE}c6i{self.PIPE}t2{self.PIPE}d3{self.PIPE}i2{self.PIPE}r5ad{self.PIPE}u-18tb1{self.PIPE}c7gn{self.PIPE}r5dn{self.PIPE}c5n{self.PIPE}r3{self.PIPE}c6a{self.PIPE}m1{self.PIPE}r5a{self.PIPE}m5zn{self.PIPE}u-3tb1{self.PIPE}i7ie{self.PIPE}z1d{self.PIPE}m6g{self.PIPE}r5n{self.PIPE}r5d{self.PIPE}g6{self.PIPE}r6in{self.PIPE}mac2-m1ultra{self.PIPE}inf1{self.PIPE}dl1{self.PIPE}m6in{self.PIPE}hpc7g{self.PIPE}hpc6a{self.PIPE}hpc6id{self.PIPE}hpc7a{self.PIPE}t1{self.PIPE}m8g{self.PIPE}d3en{self.PIPE}m7a{self.PIPE}r6a{self.PIPE}r7gd{self.PIPE}c6gn{self.PIPE}i8g{self.PIPE}m5d{self.PIPE}r7i{self.PIPE}inf2{self.PIPE}c6id{self.PIPE}trn1{self.PIPE}m7i{self.PIPE}r7a{self.PIPE}c6g{self.PIPE}g5g{self.PIPE}r8g{self.PIPE}c7i-flex{self.PIPE}g5{self.PIPE}c5a{self.PIPE}u-6tb1{self.PIPE}t3a{self.PIPE}c7g{self.PIPE}r7an{self.PIPE}r6idn{self.PIPE}p4d{self.PIPE}r6g{self.PIPE}c4{self.PIPE}r7g{self.PIPE}x1{self.PIPE}d2{self.PIPE}i3en{self.PIPE}x2idn{self.PIPE}m6gd{self.PIPE}t3{self.PIPE}p5{self.PIPE}vt1{self.PIPE}t4g{self.PIPE}c7gd{self.PIPE}u7in-16tb{self.PIPE}c5ad{self.PIPE}m5ad{self.PIPE}x2iedn{self.PIPE}m7i-flex{self.PIPE}m6a{self.PIPE}i4i{self.PIPE}mac1{self.PIPE}mac2-m2pro{self.PIPE}m5n{self.PIPE}r5b{self.PIPE}trn1n{self.PIPE}m5dn{self.PIPE}h1{self.PIPE}p3dn{self.PIPE}a1{self.PIPE}c1{self.PIPE}r4{self.PIPE}r6gd{self.PIPE}u-24tb1{self.PIPE}c6gd{self.PIPE}g6e{self.PIPE}u7in-32tb{self.PIPE}m2{self.PIPE}x1e{self.PIPE}u-12tb1{self.PIPE}m5a{self.PIPE}u7in-24tb{self.PIPE}p2{self.PIPE}x2iezn{self.PIPE}c5d{self.PIPE}m7g{self.PIPE}m6id{self.PIPE}i4g{self.PIPE}m7gd{self.PIPE}m4{self.PIPE}u7i-12tb{self.PIPE}x2gd{self.PIPE}g4ad{self.PIPE}m6idn{self.PIPE}im4gn{self.PIPE}p3{self.PIPE}r6i{self.PIPE}m3{self.PIPE}m5)\.(?:large{self.PIPE}medium{self.PIPE}metal-48xl{self.PIPE}18xlarge{self.PIPE}9xlarge{self.PIPE}micro{self.PIPE}12xlarge{self.PIPE}96xlarge{self.PIPE}6xlarge{self.PIPE}metal{self.PIPE}nano{self.PIPE}10xlarge{self.PIPE}3xlarge{self.PIPE}metal-32xl{self.PIPE}metal-24xl{self.PIPE}xlarge{self.PIPE}112xlarge{self.PIPE}2xlarge{self.PIPE}224xlarge{self.PIPE}small{self.PIPE}56xlarge{self.PIPE}32xlarge{self.PIPE}4xlarge{self.PIPE}48xlarge{self.PIPE}16xlarge{self.PIPE}metal-16xl{self.PIPE}24xlarge{self.PIPE}8xlarge)\b'
        self.gcp_instance_pattern = r"^(c4a|n4|a3|n2|z3|a2|m2|n2d|f1|g1|g2|c3d|m1|h3|c4|t2a|n1|c2d|c3|e2|m3|t2d|c2|c4d|n4d)-(ultragpu|megamem|standard|megagpu|highgpu|ultramem|medium|highcpu|hypermem|small|micro|highmem)(-(60|48|22|224|1g|416|208|44|96|56|4|72|128|16|112|80|90|180|16g|24|176|2|8|2g|32|30|160|1|192|12|40|8g|360|4g|64|88|384|144|288))?(-metal|-lssd)?$"
        self.MAX_MEM_USED = "max mem used"
        self.MAX_NW_BW = "max network bw"
        self.MAX_DISK_BW = "max disk bw used"
        self.MAX_IOPS = "max iops"
        self.NUMBER_VALIDATION = "must be a positive number"
        self.UNSUPPORTED_PRICING_MODEL = [""]
        self.results_path = "/var/www/html/results"
        self.DEFAULT_PT = Pt(10)
        self.LARGE_PT = Pt(24)
        self.LARGE_KEYS = {"customer_name", "cloud_provider", "Date_Format"}
        self.CCA_HEADERS = [
            "UUID", "CSP", "Pricing Model", "Zone", "Current Instance", "Number of Instances", "vCPU", "Current Monthly Cost",
            "Annual Cost", "Recommendation I Instance", "vCPU I", "Monthly Cost I", "Annual Cost I (perf scaled)",
            "Annual Savings I", "Perf Enhancement I", "Recommendation II Instance", "vCPU II", "Monthly Cost II",
            "Annual Cost II (perf scaled)", "Annual Savings II", "Perf Enhancement II", "Recommendation III Instance",
            "vCPU III", "Monthly Cost III", "Annual Cost III (perf scaled)", "Annual Savings III", "Perf Enhancement III",
            "STATUS"
        ]
        self.EIA_HEADERS = [
            "UUID", "CSP", "Pricing Model", "Zone", "Current Instance", "vCPU", "Current Monthly Price",
            "Current Instance Energy Consumption (kwh)", "Current Instance Emission", "Recommendation I Instance",
            "vCPU I", "Monthly Price I", "Monthly Savings I", "Instance Energy Consumption I (kwh)", "Instance Emission I",
            "Perf Enhancement I", "Recommendation II Instance", "vCPU II", "Monthly Price II", "Monthly Savings II",
            "Instance Energy Consumption II (kwh)", "Instance Emission II", "Perf Enhancement II", "Input Headroom", "STATUS"
        ]
        self.AI_VALIDATION = os.getenv("AI_VALIDATIONS")
        self.default_regions = {
            'AWS': 'us-east-1',
            'AZURE': 'eastus',
            'GCP': 'us-east1'
        }
        self.PRICEMODEL = 'pricing model'
        self.valid_pricing_models = {'ondemand', 'spot', 'reserved'}
        self.AWS_ACCESS_KEY = secret_data.get('AWS_ACCESS_KEY')
        self.AWS_SECRET_KEY = secret_data.get("AWS_SECRET_KEY")
        self.AWS_REGION = os.getenv('AWS_REGION')
        self.BUCKET_NAME = os.getenv('BUCKET_NAME')
        self.MAIN_FOLDER = os.getenv('MAIN_FOLDER')
        self.CCA_APP_NAME = "EPYC Cloud Cost Advisor"
        self.EIA_APP_NAME = "EPYC Cloud Instance Advisor"
        self.CCA_UI = os.getenv('CCA_UI')
        self.EIA_UI = os.getenv('EIA_UI')
        self.support_links = {
            "CCA": f"{self.CCA_UI}/support",
            "EIA": f"{self.EIA_UI}/support"
        }
        self.EMAIL_TEMPLATES = {
            "CCA": """
            <html>
            <body>
                <p>Hello,</p>
                <p>Your recommendations for the portfolio <b>{portfolio_name}</b> have been successfully generated and are now available.</p>
                <p>The attached ZIP file includes:</p>
                <ul>
                    <li>Exported Excel report of the generated recommendations</li>
                    <li>Analysis PowerPoint (PPT)</li>
                </ul>
                <p>Thank you for using AMD EPYC Cloud Cost Advisor.</p>
                <br/>
                <p>Best regards,<br/>{SENDER}</p>
            </body>
            </html>
            """,

            "EIA": """
            <html>
            <body>
                <p>Hello,</p>
                <p>Your recommendations for the portfolio <b>{portfolio_name}</b> have been successfully generated and are now available.</p>
                <p>The attached ZIP file includes:</p>
                <ul>
                    <li>Exported Excel report of the generated recommendations</li>
                    <li>Defective rows list (if any)</li>
                </ul>
                <p>Thank you for using AMD EPYC Cloud Instance Advisor.</p>
                <br/>
                <p>Best regards,<br/>{SENDER}</p>
            </body>
            </html>
            """
        }
        self.WHITE_BORDER = Border(
            left=Side(style="thin", color="FFFFFF"),
            right=Side(style="thin", color="FFFFFF"),
            top=Side(style="thin", color="FFFFFF"),
            bottom=Side(style="thin", color="FFFFFF"),
        )
        self.BLACK_HEADER_FILL = PatternFill(start_color="000000", end_color="000000", fill_type="solid")
        self.GRAY_FILL = PatternFill(start_color="D9D9D9", end_color="D9D9D9", fill_type="solid")
        self.LIGHT_GREEN_FILL = PatternFill(start_color="00FF00", end_color="00FF00", fill_type="solid")
        self.WHITE_FILL = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")
        self.WHITE_BOLD_FONT = Font(bold=True, color="FFFFFF")
        self.NORMAL_FONT = Font(bold=False, color="000000")
        self.NOTE_FONT = Font(name='Cambria', size=11, color="0000FF", italic=False)
        self.BOLD_NORMAL_FONT = Font(bold=True, color="000000")
        self.DISCLAIMER_FONT = Font(name='Cambria', size=11, bold=False, color="000000")
        self.DISCLAIMER_TITLE_FONT = Font(name='Cambria', size=11, bold=True, color="000000")
        self.HEADER_ALIGNMENT = Alignment(wrap_text=True, horizontal="center", vertical="center")
        self.CENTER_ALIGNMENT = Alignment(horizontal="center", vertical="center")
        self.LEFT_ALIGNMENT = Alignment(horizontal="left", vertical="center")
        self.RIGHT_ALIGNMENT = Alignment(horizontal="right", vertical="center")
        self.WHITE_SIDE = Side(border_style="thin", color="FFFFFF")
        self.BLACK_SIDE = Side(border_style="thin", color="000000")
        self.THIN_BLACK_SIDE = Side(border_style="thin", color="000000")
        # Defined Borders
        self.R1_BOTTOM_BORDER = Border(bottom=self.WHITE_SIDE, left=self.WHITE_SIDE, right=self.WHITE_SIDE, top=self.WHITE_SIDE)
        self.R4_HEADER_BORDER = Border(left=self.WHITE_SIDE, right=self.WHITE_SIDE, top=self.WHITE_SIDE, bottom=None)
        self.R5_BOTTOM_BORDER = Border(left=self.WHITE_SIDE, right=self.WHITE_SIDE, top=None, bottom=self.WHITE_SIDE) 
        self.TITLE_BORDER = Border(left=self.WHITE_SIDE, right=self.WHITE_SIDE, top=self.WHITE_SIDE, bottom=self.WHITE_SIDE)
        self.THICK_BLACK_RIGHT_SIDE = Side(border_style="medium", color="000000")
        self.R3_SEPARATOR_BORDER = Border()
        self.R4_SEPARATOR_BORDER = Border()
        # --- New/Modified Global Constants for Column Widths ---
        self.INSTANCE_TYPE_COLS = [1, 10, 17] 
        self.SMALL_INSTANCE_TYPE_WIDTH = 12 
        self.MAX_COLUMN_WIDTH = 35 # Set the overall maximum width limit for all columns

        if not self.mongo_uri or not self.db_name:
            self.log_message(LevelType.ERROR, "Environment variables MONGO_URI and DATABASE_NAME must be set!", ErrorCode=1)
            return

        try:
            self.client = MongoClient(self.mongo_uri)
            self.db = self.client[self.db_name]
            self.recommendation_tracking_collection = self.db[CollectionNames.RECCOMENDATION_TRACKING]
            self.portfolio_collection = self.db[CollectionNames.PORTFOLIOS]
            self.current_instance_collection = self.db[CollectionNames.CURRENT_INSTANCES]
            self.recommended_instance_collection = self.db[CollectionNames.RECOMMENDED_INSTANCES]
            self.input_remarks_collection = self.db[CollectionNames.INPUT_REMARKS]
            self.notification_collection = self.db[CollectionNames.NOTIFICATIONS]
            self.recommended_analytics = self.db[CollectionNames.RECOMMENDATION_ANALYTICS]
            self.unsupported_recommendation_analytics = self.db[CollectionNames.ANALYTICS_WITHOUT_RECOMMENDATION]
            self.log_message(LevelType.INFO, f"Connected to MongoDB database '{self.db_name}' successfully.")
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error connecting to MongoDB: {e}", ErrorCode=2)
            self.client = None
            self.db = None
    def close_connection(self):
        if self.client:
            self.client.close()
            self.client = None
            self.db = None

    def fetch_and_sort_portfolioids(self):
        portfolio_ids = self.recommendation_tracking_collection.distinct(
            "portfolio_id",
            {"recommendation_status": RecommendationStatus.QUEUE}
        )

        if not portfolio_ids:
            self.log_message(LevelType.INFO, "No portfolios with queued batches found.", ErrorCode=1)
            return []

        portfolio_times = {}
        for pid in portfolio_ids:
            batch_doc = self.recommendation_tracking_collection.find_one(
                {"portfolio_id": pid, "recommendation_status": RecommendationStatus.QUEUE},
                sort=[("created_at", ASCENDING)]
            )
            if batch_doc:
                portfolio_times[pid] = batch_doc["created_at"]

        sorted_portfolios = sorted(portfolio_times.items(), key=lambda x: x[1])
        return [pid for pid, _ in sorted_portfolios]
    
    def get_total_batches_for_portfolio(self, portfolio_id: str) -> int:
        """
        Returns the total number of batches in QUEUE status for a given portfolio.
        """
        query = {
            "portfolio_id": portfolio_id
        }
        total_batches = self.recommendation_tracking_collection.count_documents(query)
        return total_batches
    
    def get_batch_count_on_status(self, portfolio_id):
        """
        Returns the total number of batches for a portfolio
        with status either COMPLETED or FAILED.
        """
        query = {
            "portfolio_id": portfolio_id,
            "recommendation_status": {"$in": [RecommendationStatus.COMPLETED, RecommendationStatus.FAILED]}
        }
        processed_batch = self.recommendation_tracking_collection.count_documents(query)
        return processed_batch

    def get_next_batch_for_portfolio(self, portfolio_id):
        try:
            return self.recommendation_tracking_collection.find_one(
                {"portfolio_id": portfolio_id, "recommendation_status": RecommendationStatus.QUEUE},
                sort=[("batch_id", ASCENDING)]
            )
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error getting next batch for portfolio {portfolio_id}: {str(e)}")
            return None

    def get_access_token(self):
        try:
            app = ConfidentialClientApplication(
                client_id=CLIENT_ID,
                authority=AUTHORITY,
                client_credential=CLIENT_SECRET
            )
            token_response = app.acquire_token_for_client(scopes=SCOPE)
            if "access_token" not in token_response:
                self.log_message(LevelType.ERROR, f"Failed to get token: {token_response.get('error_description')}")
                return None
            return token_response['access_token']
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Exception getting access token: {str(e)}")
            return None

    def post_send_mail(self, payload, token):
        try:
            headers = {
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/json"
            }
            return requests.post(GRAPH_ENDPOINT, headers=headers, data=json.dumps(payload))
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Exception sending mail post: {str(e)}")
            return None

    def send_email(self, subject, to_email, body, content_type="HTML"):
        try:
            token = self.get_access_token()
            if not token:
                return "Failed to acquire access token", False

            if isinstance(to_email, str):
                recipients = [email.strip() for email in to_email.split(",")]
            else:
                recipients = to_email

            if not to_email:
                self.log_message(LevelType.WARNING, "No user emails specified to send")
                return "No user emails specified to send", False

            message_payload = {
                "message": {
                    "subject": subject,
                    "body": {
                        "contentType": content_type,
                        "content": body
                    },
                    "toRecipients": [
                        {"emailAddress": {"address": email}} for email in recipients
                    ]
                },
                "saveToSentItems": "false"
            }

            headers = {
                'Authorization': f'Bearer {token}',
                'Content-Type': 'application/json'
            }

            response = requests.post(GRAPH_ENDPOINT, headers=headers, data=json.dumps(message_payload))

            if response.status_code == 202:
                return f"Email sent to {to_email}", True
            else:
                self.log_message(LevelType.ERROR, f"Failed to send email: {response.text}")
                return f"Failed to send email: {response.text}", False
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Exception during email send: {str(e)}")
            return f"Exception during email send: {str(e)}", False

    def update_portfolio_status_if_all_batches_completed(self, portfolio_id, rec_failed_to_emails, user_email, portfolio_name, app_name):
        try:
            failed_batches = list(self.recommendation_tracking_collection.find({
                "portfolio_id": portfolio_id,
                "recommendation_status": "FAILED"
            }, {"batch_id": 1, "_id": 0}))

            if app_name == "CCA":
                application_name = self.CCA_APP_NAME
            else:
                application_name = self.EIA_APP_NAME
 
            if failed_batches:
                self.portfolio_collection.update_one(
                    {"_id": ObjectId(portfolio_id)},
                    {"$set": {"recommendation_status": "FAILED"}}
                )
                self.log_message(LevelType.INFO, f"For application {application_name} Portfolio {portfolio_id} recommendation_status set to FAILED")
 
                batch_ids = [str(batch["batch_id"]) for batch in failed_batches]

                subject = f"{application_name} Recommendation Failed for Portfolio: {portfolio_name}"

                body = f"""
                <html>
                <body>
                    <p>Hi Team,</p>
                    <p>The following recommendation batch of {application_name} has failed:</p>
                    <p>
                        <b>Portfolio:</b> {portfolio_name}<br/>
                        <b>Portfolio ID:</b> {portfolio_id}<br/>
                        <b>Batch ID:</b> {', '.join(batch_ids)}
                    </p>
                    <p>Please review the failed batch logs and take the necessary corrective actions promptly.</p>
                    <br/>
                    <p>Regards,<br/>EPYC Advisory Services System Notification</p>
                </body>
                </html>
                """
                recipients = [email.strip() for email in rec_failed_to_emails.split(",") if email.strip()]
 
                self.send_email(subject, recipients, body)

                # Get link for current app (fallback to EIA if not found)
                support_url = self.support_links.get(app_name.upper())

                # Build subject and body
                user_subject = f"{application_name} Recommendation Failed for Portfolio: {portfolio_name}"

                user_body = f"""
                <html>
                <body>
                    <p>Hello,</p>
                    <p>We regret to inform you that the recommendation generation for your portfolio <b>{portfolio_name}</b> has failed.</p>
                    <p>Please contact the 
                        <a href="{support_url}" target="_blank" style="color:#0078d4; text-decoration:none;">
                            EPYC Advisory Services
                        </a> 
                        team for assistance or more information.
                    </p>
                    <br/>
                    <p>Best regards,<br/>{SENDER}</p>
                </body>
                </html>
                """
                self.send_email(user_subject, user_email, user_body)
                return False
            return True
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error in update_portfolio_status_if_all_batches_completed: {str(e)}")
            return False

    def sanitize_for_json(self, obj):
        """
        Recursively convert NaN/NaT/Inf to None inside dicts/lists/scalars.
        Also handles pandas DataFrames/Series by returning built-in Python types.
        """
        if isinstance(obj, pd.DataFrame):
            df = obj.replace([np.inf, -np.inf], np.nan)
            df = df.where(pd.notna(df), None)
            return [self.sanitize_for_json(r) for r in df.to_dict(orient="records")]
        if isinstance(obj, pd.Series):
            return self.sanitize_for_json(obj.to_dict())
        if isinstance(obj, dict):
            return {k: self.sanitize_for_json(v) for k, v in obj.items()}
        if isinstance(obj, (list, tuple)):
            return [self.sanitize_for_json(v) for v in obj]
        try:
            if isinstance(obj, float) and (math.isnan(obj) or math.isinf(obj)):
                return None
            return obj
        except Exception:
            return obj

    def transform_rows(self, rows):
        """Transform rows by extracting remarks and cleaning fields."""
        transformed = []
        for row in rows:
            remarks = []
            new_row = {}
            for key, value in row.items():
                if key.endswith("_error") and value is not None:  # has error
                    field_name = key.replace("_error", "").strip()
                    remarks.append({"Field": field_name, "Message": value})
                else:
                    new_row[key] = value
            new_row["Remarks"] = remarks
            transformed.append(new_row)
        return transformed

    def cleanup_files(self, *file_paths):
        for file_path in file_paths:
            if file_path and os.path.exists(file_path):
                try:
                    os.remove(file_path)
                except Exception as e:
                    self.log_message(LevelType.ERROR, f"Error removing file {file_path}: {e}")

    def create_instance_udf_files_from_json(self, json_data, headroom_value, udf_data, instance_file_path, udf_file_path):
        data = json_data
        udf = udf_data
        udf_file_created = False

        if data:
            skip_fields = {'pavg', 'uavg', 'p95', 'u95', 'instance name'}
            updated_data = [
                {
                    k.replace('pricingModel', 'pricing model') if k == 'pricingModel' else k: v
                    for k, v in item.items()
                    if k not in skip_fields
                } for item in data
            ]

            # Inject headroom% to each row
            for row in updated_data:
                row["headroom%"] = headroom_value
            data_df = pd.DataFrame(updated_data)
            data_df.to_csv(instance_file_path, index=False)

        if udf:
            udf_df = pd.DataFrame(udf)
            udf_df.to_csv(udf_file_path, index=False)
            udf_file_created = True

        return instance_file_path, udf_file_path if udf_file_created else None

    def costadvise_utils(self, total_data, input_folder_path, output_folder_path, csv_file_name, policy_data, policy_engine_file_path=None):
        policy_file_created = False
        summarized_list = [
            {
                "instance type": entry["instance type"],
                "region": entry["region"],
                "quantity": float(entry["quantity"]),
                "monthly utilization": float(entry[self.MONTHLY_UTILIZATION]),
                "cloud_csp": entry["cloud_csp"],
                "pricing model": entry["pricingModel"],
                "uuid": entry["instance_name"] if "instance_name" in entry and entry["instance_name"].strip() != "" else entry.get("uuid", '')
            } for entry in total_data
        ]

        df = pd.DataFrame(summarized_list)

        input_csv_file_path = os.path.join(input_folder_path, csv_file_name)
        output_csv_file_path = os.path.join(output_folder_path, csv_file_name)

        try:
            df.to_csv(input_csv_file_path, index=False)
            self.log_message(LevelType.INFO, f"Written CSV file to {input_csv_file_path}")
            if policy_data:
                policy_df = pd.DataFrame(policy_data)
                # Rename columns to remove underscore
                policy_df = policy_df.rename(columns={
                    "instance_type": "instance type",
                    "scalar_value": "scalar value"
                })
                policy_df.to_csv(policy_engine_file_path, index=False)
                policy_file_created = True
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Failed to write CSV file {input_csv_file_path}: {e}")

        return input_csv_file_path, output_csv_file_path, policy_engine_file_path if policy_file_created else None

    def run_command(self, command):
        try:
            result = subprocess.run(command, universal_newlines=True, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            return result.stdout, result.stderr
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error running command '{command}': {e}")
            return "", str(e)

    def is_command_successful(self, std_out, std_err):
        if ("Error" in std_out or "HDF5-DIAG" in std_out or "terminate called" in std_out or
            "Segmentation fault" in std_out or "terminated" in std_out):
            return False
        if std_err:
            return False
        return True

    def cca_process_data_perf(self, perf_enhancement_i, perf_enhancement_ii, perf_enhancement_iii, row, key, sums):
        try:
            value = float(row.get(key, 0)) if row.get(key, "").strip() else 0
            sums[key] += value

            if row.get(self.PERF_ENHANCEMENT_I, "").strip() not in ["", "-", "inf"]:
                perf_enhancement_i.append(float(row[self.PERF_ENHANCEMENT_I]))

            if row.get(self.PERF_ENHANCEMENT_II, "").strip() not in ["", "-", "inf"]:
                perf_enhancement_ii.append(float(row[self.PERF_ENHANCEMENT_II]))

            if row.get(self.PERF_ENHANCEMENT_III, "").strip() not in ["", "-", "inf"]:
                perf_enhancement_iii.append(float(row[self.PERF_ENHANCEMENT_III]))
            return perf_enhancement_i, perf_enhancement_ii, perf_enhancement_iii, True
        except ValueError:
            return perf_enhancement_i, perf_enhancement_ii, perf_enhancement_iii, False


    def cca_process_output_data(self, dict_data, sums):
        perf_enhancement_i, perf_enhancement_ii, perf_enhancement_iii = [], [], []
        for row in dict_data:
            for key in sums.keys():
                perf_enhancement_i, perf_enhancement_ii, perf_enhancement_iii, flag = self.cca_process_data_perf(perf_enhancement_i, perf_enhancement_ii, perf_enhancement_iii, row, key, sums)
                if not flag:
                    continue
        return perf_enhancement_i, perf_enhancement_ii, perf_enhancement_iii

    def cca_process_output_file(self, output_file, std_out, std_err, is_downloadable=True):
        xl_file = output_file.replace('.csv', '.xlsx')
        if is_downloadable:
            self.log_message(LevelType.INFO, f"Copying cca xl_file {xl_file} to results_path : {self.results_path}")
            if os.path.exists(output_file):
                try:
                    shutil.copy(xl_file, self.results_path)
                except Exception as e:
                    self.log_message(LevelType.ERROR, f"Copy failed: {str(e)}")
        try:
            with open(output_file, 'r') as f:
                reader = csv.reader(f)
                data = list(reader)
                if not data:
                    return std_out.strip() + "\n" + std_err.strip(), None, False

                header = [col.strip() for col in data[0]]
                dict_data = [
                    {header[i].strip(): row[i].strip() for i in range(min(len(header), len(row)))}
                    for row in data[1:] if row
                ]
                sums = {
                    'Number of Instances': 0,
                    'Current Monthly Cost': 0,
                    self.ANNUAL_COST: 0,
                    'Monthly Cost I': 0,
                    'Annual Cost I (perf scaled)': 0,
                    self.ANNUAL_SAVINGS_I: 0,
                    self.PERF_ENHANCEMENT_I: 0,
                    'Monthly Cost II': 0,
                    'Annual Cost II (perf scaled)': 0,
                    self.ANNUAL_SAVINGS_II: 0,
                    self.PERF_ENHANCEMENT_II: 0,
                    'Monthly Cost III': 0,
                    'Annual Cost III (perf scaled)': 0,
                    self.ANNUAL_SAVINGS_III: 0,
                    self.PERF_ENHANCEMENT_III: 0
                }
                perf_enhancement_i, perf_enhancement_ii, perf_enhancement_iii = self.cca_process_output_data(dict_data, sums)

                avg_perf_enhancement_i = round(sum(perf_enhancement_i) / len(perf_enhancement_i), 2) if perf_enhancement_i else 0
                avg_perf_enhancement_ii = round(sum(perf_enhancement_ii) / len(perf_enhancement_ii), 2) if perf_enhancement_ii else 0
                avg_perf_enhancement_iii = round(sum(perf_enhancement_iii) / len(perf_enhancement_iii), 2) if perf_enhancement_iii else 0

                sums[self.PERF_ENHANCEMENT_I] = avg_perf_enhancement_i
                sums[self.PERF_ENHANCEMENT_II] = avg_perf_enhancement_ii
                sums[self.PERF_ENHANCEMENT_III] = avg_perf_enhancement_iii

                sums = {key: round(value, 2) for key, value in sums.items()}
                return std_out.strip() + "\n" + std_err.strip(), {'data': dict_data, 'grandTotal': sums}, True
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error processing cca output: {str(e)}")
            return std_out.strip() + "\n" + std_err.strip(), None, False

    def eia_process_output_file(self, output_file, std_out, std_err, is_downloadable=True):
        xl_file = output_file.replace('.csv', '.xlsx')
        if is_downloadable:
            self.log_message(LevelType.INFO, f"Copying eia xl_file {xl_file} to results_path : {self.results_path}")
            if os.path.exists(output_file):
                try:
                    shutil.copy(xl_file, self.results_path)
                except Exception as e:
                    self.log_message(LevelType.ERROR, f"Copy failed: {str(e)}")

        try:
            with open(output_file, 'r') as f:
                reader = csv.reader(f)
                data = list(reader)
                if not data:
                    return std_out.strip() + "\n" + std_err.strip(), None, False

                header = [col.strip() for col in data[0]]
                dict_data = [
                    {header[i].strip(): row[i].strip() for i in range(min(len(header), len(row)))}
                    for row in data[1:] if row
                ]

                sums = {
                    'Current Monthly Price': 0.0,
                    'Current Instance Energy Consumption (kwh)': 0.0,
                    'Current Instance Emission': 0.0,
                    'Monthly Price I': 0.0,
                    'Instance Energy Consumption I (kwh)': 0.0,
                    'Instance Emission I': 0.0,
                    'Monthly Savings I': 0.0,
                    'Monthly Price II': 0.0,
                    'Instance Energy Consumption II (kwh)': 0.0,
                    'Instance Emission II': 0.0,
                    'Monthly Savings II': 0.0,
                }

                perf_enhancement_i = []
                perf_enhancement_ii = []
                untapped_capacity_i = []
                untapped_capacity_ii = []

                for row in dict_data:
                    for key in sums:
                        try:
                            value = row.get(key, "").strip()
                            if value not in ["", "-", None]:
                                sums[key] += float(value)
                        except ValueError:
                            pass

                    # Collect performance enhancement values
                    try:
                        val_i = row.get('Perf Enhancement I', "").strip()
                        if val_i not in ["", "-", None]:
                            perf_enhancement_i.append(float(val_i))
                    except ValueError:
                        pass

                    try:
                        val_ii = row.get('Perf Enhancement II', "").strip()
                        if val_ii not in ["", "-", None]:
                            perf_enhancement_ii.append(float(val_ii))
                    except ValueError:
                        pass

                    try:
                        val_uc_i = row.get('Untapped Capacity I', "").strip()
                        if val_uc_i not in ["", "-", None]:
                            untapped_capacity_i.append(float(val_uc_i))
                    except ValueError:
                        pass

                    try:
                        val_uc_ii = row.get('Untapped Capacity II', "").strip()
                        if val_uc_ii not in ["", "-", None]:
                            untapped_capacity_ii.append(float(val_uc_ii))
                    except ValueError:
                        pass

                # Calculate averages
                sums['Perf Enhancement I'] = round(sum(perf_enhancement_i) / len(perf_enhancement_i), 2) if perf_enhancement_i else 0.0
                sums['Perf Enhancement II'] = round(sum(perf_enhancement_ii) / len(perf_enhancement_ii), 2) if perf_enhancement_ii else 0.0
                sums['Untapped Capacity I'] = round(sum(untapped_capacity_i) / len(untapped_capacity_i), 2) if untapped_capacity_i else 0.0
                sums['Untapped Capacity II'] = round(sum(untapped_capacity_ii) / len(untapped_capacity_ii), 2) if untapped_capacity_ii else 0.0

                sums = {key: round(value, 2) for key, value in sums.items()}

                return std_out.strip() + "\n" + std_err.strip(), {'data': dict_data, 'grandTotal': sums}, True
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error processing eia output: {str(e)}")
            return std_out.strip() + "\n" + std_err.strip(), None, False

    def data_collection(self, input_file, output_file, udf_file, policy_engine_file_path, app_name, is_downloadable = True):
        msg = ""
        try:
            if udf_file:
                command = f"bash run.sh {input_file} {output_file} {app_name.upper()} {udf_file} no"
            elif policy_engine_file_path:
                command = f"bash run.sh {input_file} {output_file} {app_name.upper()} {policy_engine_file_path} no"
            else:
                command = f"bash run.sh {input_file} {output_file} {app_name.upper()} - no"
            out = "STDOUT: "
            err = "STDERR: "
            std_out, std_err = self.run_command(command)
            msg = out + std_out + err + std_err
            if not self.is_command_successful(std_out, std_err):
                return msg, None, False
            if os.path.exists(output_file):
                if app_name.upper() == 'CCA':
                    return self.cca_process_output_file(output_file, out + std_out, err + std_err, is_downloadable)
                else:
                    return self.eia_process_output_file(output_file, out + std_out, err + std_err, is_downloadable)
            else:
                return msg, None, False
        except subprocess.TimeoutExpired as err:
            self.log_message(LevelType.ERROR, str(err))
            return msg, None, False
        except Exception as err:
            self.log_message(LevelType.ERROR, f"error in summit : {str(err)}")
            return msg, None, False

    def saving_calculation(self, savings, cost):
        if cost != "-" and savings != '-':
            cost_val = float(cost)
            if cost_val == 0:
                return 0.0
            save_value = float(savings)
            if save_value <= 0:
                save_value = 0
            return round((save_value / cost_val) * 100, 2)
        else:
            return "-"

    def cca_transformed_data(self, entry, transformed_data, result):
        entry = {k.strip(): v.strip() if isinstance(v, str) else v for k, v in entry.items()}
        try:
            transformed_entry = {
                "id": entry["UUID"],
                "data": {
                    "currentPlatform": {
                        "zone": entry["Zone"],
                        "instanceType": entry["Current Instance"],
                        "numberOfInstances": entry["Number of Instances"],
                        "vCPU": entry["vCPU"],
                        "monthlyCost": entry["Current Monthly Cost"],
                        "annualCost": entry[self.ANNUAL_COST],
                        "cspProvider": entry["CSP"],
                        "pricingModel": entry["Pricing Model"],
                        "status": entry["STATUS"]
                    },
                    "recommendations": [
                        {
                            "zone": entry["Zone"],
                            "instanceType": entry["Recommendation I Instance"],
                            "vCPU": entry["vCPU I"],
                            "monthlyCost": entry["Monthly Cost I"],
                            "totalCost": entry["Annual Cost I (perf scaled)"],
                            "annualSavings": 0.0 if entry[self.ANNUAL_SAVINGS_I] == "-0.000000" else entry[self.ANNUAL_SAVINGS_I],
                            "savingsInPercentage": self.saving_calculation(entry[self.ANNUAL_SAVINGS_I], entry[self.ANNUAL_COST]),
                            "perf": entry["Perf Enhancement I"]
                        },
                        {
                            "zone": entry["Zone"],
                            "instanceType": entry["Recommendation II Instance"],
                            "vCPU": entry["vCPU II"],
                            "monthlyCost": entry["Monthly Cost II"],
                            "totalCost": entry["Annual Cost II (perf scaled)"],
                            "annualSavings": 0.0 if entry[self.ANNUAL_SAVINGS_II] == "-0.000000" else entry[self.ANNUAL_SAVINGS_II],
                            "savingsInPercentage": self.saving_calculation(entry[self.ANNUAL_SAVINGS_II], entry[self.ANNUAL_COST]),
                            "perf": entry["Perf Enhancement II"]
                        },
                        {
                            "zone": entry["Zone"],
                            "instanceType": entry["Recommendation III Instance"],
                            "vCPU": entry["vCPU III"],
                            "monthlyCost": entry["Monthly Cost III"],
                            "totalCost": entry["Annual Cost III (perf scaled)"],
                            "annualSavings": 0.0 if entry[self.ANNUAL_SAVINGS_III] == "-0.000000" else entry[self.ANNUAL_SAVINGS_III],
                            "savingsInPercentage": self.saving_calculation(entry[self.ANNUAL_SAVINGS_III], entry[self.ANNUAL_COST]),
                            "perf": entry["Perf Enhancement III"]
                        }
                    ]
                }
            }

            h_saving_percentage = round((float(result['grandTotal'][self.ANNUAL_SAVINGS_I]) / (
                max(1, float(result['grandTotal'][self.ANNUAL_COST]))) * 100), 2)
            m_saving_percentage = round((float(result['grandTotal'][self.ANNUAL_SAVINGS_II]) / (
                max(1, float(result['grandTotal'][self.ANNUAL_COST]))) * 100), 2)
            md_saving_percentage = round((float(result['grandTotal'][self.ANNUAL_SAVINGS_III]) / (
                max(1, float(result['grandTotal'][self.ANNUAL_COST]))) * 100), 2)
            result['grandTotal']['hSavingsInPercentage'] = h_saving_percentage
            result['grandTotal']['mSavingsInPercentage'] = m_saving_percentage
            result['grandTotal']['mdSavingsInPercentage'] = md_saving_percentage
            transformed_data.append(transformed_entry)
            return transformed_data
        except Exception as e:
            self.log_message(LevelType.ERROR, str(e))
            return None

    def get_transformed_rec_data(self, app_flag, result):
        transformed_data = []
        if app_flag.lower() == 'cca':
            for entry in result['data']:
                transformed_data = self.cca_transformed_data(entry, transformed_data, result)
        else:
            for entry in result['data']:
                entry = {k.strip(): v.strip() if isinstance(v, str) else v for k, v in entry.items()}
                try:
                    transformed_entry = {
                        "id": entry["UUID"],
                        "csp": entry["CSP"],
                        "data": {
                            "currentPlatform": {
                                "type": entry["Current Instance"],
                                "cost": entry["Current Monthly Price"],
                                "power": entry["Current Instance Energy Consumption (kwh)"],
                                "carbon": entry["Current Instance Emission"],
                                "status": entry["STATUS"],
                                "vCPU": entry["vCPU"],
                                "pricingModel": entry["Pricing Model"],
                                "region": entry["Zone"]
                            },
                            "recommendations": [
                                {
                                    "cost": entry["Monthly Price I"],
                                    "type": entry["Recommendation I Instance"],
                                    "power": entry["Instance Energy Consumption I (kwh)"],
                                    "carbon": entry["Instance Emission I"],
                                    "perf": entry["Perf Enhancement I"],
                                    "monthlySavings": entry["Monthly Savings I"],
                                    "vCPU": entry["vCPU I"],
                                    "untappedCapacity": entry["Untapped Capacity I"]
                                },
                                {
                                    "cost": entry["Monthly Price II"],
                                    "type": entry["Recommendation II Instance"],
                                    "power": entry["Instance Energy Consumption II (kwh)"],
                                    "carbon": entry["Instance Emission II"],
                                    "perf": entry["Perf Enhancement II"],
                                    "monthlySavings": entry["Monthly Savings II"],
                                    "vCPU": entry["vCPU II"],
                                    "untappedCapacity": entry["Untapped Capacity II"]
                                }
                            ]
                        }
                    }
                    transformed_data.append(transformed_entry)
                except Exception as e:
                    self.log_message(LevelType.ERROR, str(e))
                    return None
        return transformed_data

    def store_eia_recommendations(self,
        transformed_data,
        portfolio_id,is_delete=True,
        app_name="EIA"
    ):
        if isinstance(transformed_data, dict):
            transformed_data = [transformed_data]

        if is_delete:
            delete_result = self.recommended_instance_collection.delete_many({"portfolio_id": portfolio_id})
            self.log_message(LevelType.INFO, f"Deleted {delete_result.deleted_count} old recommendations where app: {app_name} for portfolio_id : {portfolio_id}")

        records_to_insert = []
        for rec in transformed_data:
            current = rec.get("data", {}).get("currentPlatform", {})
            recommendations = rec.get("data", {}).get("recommendations", [])

            rec1 = recommendations[0] if len(recommendations) > 0 else {}
            rec2 = recommendations[1] if len(recommendations) > 1 else {}

            records_to_insert.append({
                "portfolio_id": str(portfolio_id),
                "UUID": rec.get("id", ""),
                "CSP": rec.get("csp", "").upper(),
                "Pricing Model": current.get("pricingModel", ""),
                "Zone": current.get("region", ""),
                "Current Instance": current.get("type", ""),
                "vCPU": current.get("vCPU", ""),
                "Current Monthly Price": current.get("cost", ""),
                "Current Instance Energy Consumption (kwh)": current.get("power", ""),
                "Current Instance Emission": current.get("carbon", ""),
                "Recommendation I Instance": rec1.get("type", ""),
                "vCPU I": rec1.get("vCPU", ""),
                "Monthly Price I": rec1.get("cost", ""),
                "Monthly Savings I": rec1.get("monthlySavings", ""),
                "Instance Energy Consumption I (kwh)": rec1.get("power", ""),
                "Instance Emission I": rec1.get("carbon", ""),
                "Perf Enhancement I": rec1.get("perf", ""),
                "Untapped Capacity I": rec1.get("untappedCapacity", ""),
                "Recommendation II Instance": rec2.get("type", ""),
                "vCPU II": rec2.get("vCPU", ""),
                "Monthly Price II": rec2.get("cost", ""),
                "Monthly Savings II": rec2.get("monthlySavings", ""),
                "Instance Energy Consumption II (kwh)": rec2.get("power", ""),
                "Instance Emission II": rec2.get("carbon", ""),
                "Perf Enhancement II": rec2.get("perf", ""),
                "Untapped Capacity II": rec2.get("untappedCapacity", ""),
                "STATUS": current.get("status", ""),
                "comments": rec.get("comments", ""),
                "created_at": datetime.utcnow()
            })

        if records_to_insert:
            result = self.recommended_instance_collection.insert_many(records_to_insert)
            self.log_message(LevelType.INFO, f"Inserted {len(result.inserted_ids)} recommendations where app: {app_name} portfolio_id : {portfolio_id}")
        else:
            self.log_message(LevelType.INFO, f"No recommendations to insert for app: {app_name} portfolio_id : {portfolio_id}")

        self.log_message(LevelType.INFO, f"Finished storing {len(records_to_insert)} recommendations where app: {app_name} portfolio_id : {portfolio_id}")

        return records_to_insert

    def store_cca_recommendations_bulk(self, transformed_data, portfolio_id,is_delete=True, app_name="CCA"):
        if isinstance(transformed_data, dict):
            transformed_data = [transformed_data]
        if is_delete:
            delete_result = self.recommended_instance_collection.delete_many({"portfolio_id": portfolio_id})
            self.log_message(LevelType.INFO, f"Deleted {delete_result.deleted_count} old recommendations for portfolio_id={portfolio_id} where app : {app_name}")
        records_to_insert = []
        for idx, rec in enumerate(transformed_data, start=1):
            current = rec.get("data", {}).get("currentPlatform", {})
            recs = rec.get("data", {}).get("recommendations", [])

            rec1 = recs[0] if len(recs) > 0 else {}
            rec2 = recs[1] if len(recs) > 1 else {}
            rec3 = recs[2] if len(recs) > 2 else {}

            records_to_insert.append({
                "portfolio_id": portfolio_id,
                "UUID": rec.get("id", ""),
                "CSP": current.get("cspProvider", "").upper(),
                "Pricing Model": current.get("pricingModel", ""),
                "Zone": current.get("zone", ""),
                "Current Instance": current.get("instanceType", ""),
                "vCPU": current.get("vCPU", ""),
                "Current Monthly Price": current.get("monthlyCost", ""),
                "Annual Cost": current.get("annualCost", ""),
                "Number of Instances": current.get("numberOfInstances", ""),
                "STATUS": current.get("status", ""),

                "Recommendation I Instance": rec1.get("instanceType", ""),
                "vCPU I": rec1.get("vCPU", ""),
                "Monthly Price I": rec1.get("monthlyCost", ""),
                "Annual Cost I": rec1.get("totalCost", ""),
                "Annual Savings I": rec1.get("annualSavings", ""),
                "Savings % I": rec1.get("savingsInPercentage", ""),
                "Perf Enhancement I": rec1.get("perf", ""),
                "Zone I": rec1.get("zone", ""),

                "Recommendation II Instance": rec2.get("instanceType", ""),
                "vCPU II": rec2.get("vCPU", ""),
                "Monthly Price II": rec2.get("monthlyCost", ""),
                "Annual Cost II": rec2.get("totalCost", ""),
                "Annual Savings II": rec2.get("annualSavings", ""),
                "Savings % II": rec2.get("savingsInPercentage", ""),
                "Perf Enhancement II": rec2.get("perf", ""),
                "Zone II": rec2.get("zone", ""),

                "Recommendation III Instance": rec3.get("instanceType", ""),
                "vCPU III": rec3.get("vCPU", ""),
                "Monthly Price III": rec3.get("monthlyCost", ""),
                "Annual Cost III": rec3.get("totalCost", ""),
                "Annual Savings III": rec3.get("annualSavings", ""),
                "Savings % III": rec3.get("savingsInPercentage", ""),
                "Perf Enhancement III": rec3.get("perf", ""),
                "Zone III": rec3.get("zone", ""),

                "comments": rec.get("comments", ""),
                "created_at": datetime.utcnow()
            })

        if records_to_insert:
            result = self.recommended_instance_collection.insert_many(records_to_insert)
            self.log_message(LevelType.INFO, f"Inserted {len(result.inserted_ids)} recommendations where app : {app_name} portfolio_id={portfolio_id}")
        else:
            self.log_message(LevelType.INFO, f"No recommendations to insert for portfolio_id={portfolio_id} where app : {app_name}")

        self.log_message(LevelType.INFO, f"Finished storing {len(records_to_insert)} recommendations for portfolio_id={portfolio_id} where app : {app_name}")

        return records_to_insert
    
    def get_policy_data(self, portfolio_doc):
        db = CSSessionLocal()

        try:
            policy_data = []
            default_policy = "No Policy Engine (Default)"
            policy_engine_name = portfolio_doc.get("policy_engine")
            cloud_provider = portfolio_doc.get("cloud_provider")
            user_mail = portfolio_doc.get("user_email")

            if not policy_engine_name:
                return policy_data

            if policy_engine_name.lower() != default_policy.lower():
                rows = (
                    db.query(
                        PolicyEngine.instance_type,
                        PolicyEngine.scalar_value
                    )
                    .filter(
                        func.lower(PolicyEngine.provider) == cloud_provider.lower(),
                        func.lower(PolicyEngine.policy_name) == policy_engine_name.lower(),
                        or_(PolicyEngine.user_email == user_mail, PolicyEngine.user_email == "")
                    )
                    .order_by(PolicyEngine.id)
                    .all()
                )

                policy_data = [
                    {"instance_type": r.instance_type, "scalar_value": r.scalar_value}
                    for r in rows
                ]

            return policy_data

        finally:
            db.close()   # always close session

    def process_cost_advise(self, instances, headroom, portfolio_doc, app_name, batch_id, portfolio_id, udf_data=None):
        udf_file_path = None
        policy_engine_file_path = None
        COST_ADVICE_MSG = ""
        try:
            if app_name.strip().upper() not in ["EIA", "CCA"]:
                self.log_message(LevelType.ERROR, f"Invalid App Name. batch:{batch_id}")
                return [], False

            transformed_data = {}
            file_name = f"{portfolio_doc.get('name')}_{portfolio_id}"
            csv_file_name = f"{file_name}.csv"

            input_folder_path = os.path.join(ROOT_DIR, 'input')
            output_folder_path = os.path.join(ROOT_DIR, 'output')

            if app_name.upper() == 'EIA':
                udf_folder_path = os.path.join(ROOT_DIR, 'udf')
                udf_file_name = f"{portfolio_id}_udf.csv"
                udf_file_path = os.path.join(udf_folder_path, udf_file_name)
                input_file_path = os.path.join(input_folder_path, csv_file_name)
                output_csv_file_path = os.path.join(output_folder_path, csv_file_name)
                input_csv_file_path, udf_file_path = self.create_instance_udf_files_from_json(
                    instances, headroom, udf_data, input_file_path, udf_file_path)
            else:
                policy_data = self.get_policy_data(portfolio_doc)
                policy_engine_folder_path = os.path.join(ROOT_DIR, 'udf')
                policy_engine_file_name = f"{portfolio_id}_udf.{'csv'}"
                policy_engine_file_path = os.path.join(policy_engine_folder_path, policy_engine_file_name)
                input_csv_file_path, output_csv_file_path, policy_engine_file_path = self.costadvise_utils(
                    instances, input_folder_path, output_folder_path, csv_file_name, policy_data, policy_engine_file_path)

            COST_ADVICE_MSG, result, flag = self.data_collection(
                input_csv_file_path, output_csv_file_path, udf_file_path, policy_engine_file_path, app_name, is_downloadable=False)

            if not flag:
                self.cleanup_files(input_csv_file_path)
                self.log_message(LevelType.ERROR, f"Unable to find recommendation data.batch:{batch_id} COST_ADVICE_MSG : {COST_ADVICE_MSG}, portfolio_id : {portfolio_id}")
                return [], False

            transformed_data = self.get_transformed_rec_data(app_name, result)

            if not transformed_data:
                files_to_cleanup = [input_csv_file_path, output_csv_file_path]
                xlsx_file_name = output_csv_file_path.replace('.csv', '.xlsx') if output_csv_file_path else None
                if xlsx_file_name:
                    files_to_cleanup.append(xlsx_file_name)
                if udf_file_path:
                    files_to_cleanup.append(udf_file_path)
                self.cleanup_files(*files_to_cleanup)
                self.log_message(LevelType.ERROR, f"Unable to transform the data, Details: {COST_ADVICE_MSG} batch:{batch_id}, portfolio_id : {portfolio_id}")
                return [], False

            if isinstance(transformed_data, dict):
                transformed_data = [transformed_data]
            transformed_data.sort(key=lambda x: x.get("data", {}).get("currentPlatform", {}).get("instanceType", ""))
            self.log_message(LevelType.INFO, f"length of transformed_data : {len(transformed_data)} where app : {app_name}. batch:{batch_id}, portfolio_id : {portfolio_id}")
            if int(batch_id)==1:
                is_delete = True
                result = self.recommendation_tracking_collection.update_many(
                    {"portfolio_id": portfolio_id},
                    {"$set": {"recommendation_status": "QUEUE"}}
                )
                self.log_message(LevelType.INFO, f"Matched {result.matched_count} documents and modified {result.modified_count} documents. portfolio_id : {portfolio_id}")
            else:
                is_delete = False
            if app_name.upper() == 'EIA':
                recommendation_data = self.store_eia_recommendations(
                    transformed_data=transformed_data,
                    portfolio_id=portfolio_id,is_delete=is_delete
                )
                self.log_message(LevelType.INFO, f"Recommendation successfully saved for EIA.batch:{batch_id}, portfolio_id : {portfolio_id}")
                if udf_file_path:
                    self.cleanup_files(udf_file_path)
            else:
                recommendation_data = self.store_cca_recommendations_bulk(
                    transformed_data=transformed_data,
                    portfolio_id=portfolio_id,
                    is_delete=is_delete
                )
                self.log_message(LevelType.INFO, f"Recommendation successfully saved for CCA.batch:{batch_id} portfolio_id : {portfolio_id}")

            files_to_cleanup = [input_csv_file_path, output_csv_file_path]
            xlsx_file_name = output_csv_file_path.replace('.csv', '.xlsx') if output_csv_file_path else None
            if xlsx_file_name:
                files_to_cleanup.append(xlsx_file_name)
            self.cleanup_files(*files_to_cleanup)
            return recommendation_data , True

        except Exception as e:
            self.log_message(LevelType.ERROR, f"{str(e)}, Details: {COST_ADVICE_MSG}.batch:{batch_id} portfolio_id : {portfolio_id}")
            return [], False


    def cca_eia_headers_validation(self, data, required_headers):
        missing_headers = [header for header in required_headers if header not in data.columns]
        if missing_headers:
            self.log_message(LevelType.ERROR, f"Missing headers: [{', '.join(missing_headers)}]. Please download the template to get exact headers.")
            return False, f"Missing headers: [{', '.join(missing_headers)}]. Please download the template to get exact headers."
        return True, "Headers validated sucessfully."


    def auto_correct_validation_data_eia(self, entry, provider):
        if not entry['cloud_csp'] or not isinstance(entry['cloud_csp'],str) or entry['cloud_csp'].upper() not in self.CLOUD_PROVIDERS or entry['cloud_csp'].upper() != provider:
            entry['cloud_csp'] = provider
        if self.PRICE_MODEL not in entry or (entry[self.PRICE_MODEL].lower() not in self.PRICING_MODEL):
            entry[self.PRICE_MODEL] = 'ondemand'
        return entry


    def get_all_instances_for_provider(self, provider):
        all_instances = set()
        for region in regions_map.get(provider, []):
            all_instances.update(instances_map.get((provider, region), set()))

        instance_type_data = sorted([str(item) for item in all_instances])
        return instance_type_data


    def region_instance_validate(self, region, instance, remark_list):
        if not region:
            remark_list.append({"Field": "region", "Message": self.REGION_REQUIRED})
        if not instance:
            remark_list.append({"Field": self.INSTANCE_TYPE, "Message": "Size is required"})
        return remark_list


    def regex_validation(self, instance,instance_type_data, cloud_csp, remark_list):
        if not instance:
            remark_list.append({"Field": self.INSTANCE_TYPE, "Message": self.INSTANCE_ERROR})
        elif not isinstance(instance, str):
            remark_list.append({"Field": self.INSTANCE_TYPE, "Message": f"{instance} is invalid"})
        elif instance.lower() in instance_type_data:
            if not ((cloud_csp == "AZURE" and re.match(self.azure_instance_pattern, instance.lower())) or (
                    cloud_csp == "AWS" and re.match(self.aws_instance_pattern, instance.lower())) or
                    cloud_csp == "GCP" and re.match(self.gcp_instance_pattern, instance.lower())):
                remark_list.append({"Field": self.INSTANCE_TYPE, "Message": f"{instance} is unsupported"})
        else:
            remark_list.append({"Field": self.INSTANCE_TYPE, "Message": f"{instance} is invalid"})
        return remark_list


    def region_instance_cloud_data(self, region, cloud_csp, instance, remark_list, instance_type_data, flag=None):
        if not region:
            remark_list.append({"Field": "region", "Message": self.REGION_REQUIRED})
        elif not isinstance(region, str):
            remark_list.append({"Field": "region", "Message": f"{region} is invalid"})
        elif region.lower() not in regions_map[cloud_csp]:
            remark_list.append({"Field": "region", "Message": f"{region} is invalid"})
        if not flag:
            remark_list = self.regex_validation(instance,instance_type_data, cloud_csp, remark_list)
        else:
            region_instance_type_data = sorted([str(item) for item in instances_map.get((cloud_csp, region), set())])
            if not instance:
                remark_list.append({"Field": self.INSTANCE_TYPE, "Message": self.INSTANCE_ERROR})
            if not isinstance(instance, str):
                remark_list.append({"Field": self.INSTANCE_TYPE, "Message": f"{instance} is invalid"})
            elif instance.lower() not in region_instance_type_data and instance.lower() in instance_type_data:
                remark_list.append({"Field": self.INSTANCE_TYPE, "Message": f"{instance} is unsupported"})
            elif instance.lower() not in instance_type_data:
                remark_list.append({"Field": self.INSTANCE_TYPE, "Message": f"{instance} is invalid"})
        return remark_list

    def validate_cloud_data(self, provider, region, instance, remark_list, cloud_csp, instance_type_data, flag):
        if not provider:
            remark_list.append({"Field": "cloud_csp", "Message": "Cloud is required"})
            remark_list = self.region_instance_validate(region, instance, remark_list)
        elif not isinstance(provider, str):
            remark_list.append({"Field": "cloud_csp", "Message": f"{provider} is invalid "})
            remark_list = self.region_instance_validate(region, instance, remark_list)
        elif cloud_csp != provider.upper() and provider.upper() in self.CLOUD_PROVIDERS:
            remark_list.append({"Field": "cloud_csp", "Message": "Cloud input should be same as Cloud Service Provider"})
            remark_list = self.region_instance_validate(region, instance, remark_list)
        elif provider.upper() in self.CLOUD_PROVIDERS and provider.upper() in self.UNSUPPORTED_PROVIDERS:
            remark_list.append({"Field": "cloud_csp", "Message": f"{provider} is unsupported"})
            remark_list = self.region_instance_validate(region, instance, remark_list)
        elif provider.upper() not in self.CLOUD_PROVIDERS or provider.upper() != cloud_csp:
            remark_list.append({"Field": "cloud_csp", "Message": f"{provider} is invalid"})
            remark_list = self.region_instance_validate(region, instance, remark_list)
        else:
            remark_list = self.region_instance_cloud_data(region, cloud_csp, instance, remark_list, instance_type_data, flag)
        return remark_list

    def disk_fields_validation(self, disk_max, iops_max, remark_list):
        if not disk_max:
            remark_list.append({"Field": self.MAX_DISK_BW, "Message": "Max Disk BW used is required"})
        elif not isinstance(disk_max, (int, float)) or disk_max <= 0:
            remark_list.append({"Field": self.MAX_DISK_BW, "Message": "Max Disk BW used " + self.NUMBER_VALIDATION})
        if not iops_max:
            remark_list.append({"Field": self.MAX_IOPS, "Message": "Max IOPS is required"})
        elif not isinstance(iops_max, (int, float)) or iops_max <= 0:
            remark_list.append({"Field": self.MAX_IOPS, "Message": "Max IOPS " + self.NUMBER_VALIDATION})
        return remark_list

    def validate_fields_eia(self, uuid, cpu_max, mem_max, net_max, disk_max, iops_max, pricing_model):
        remark_list = []
        if not cpu_max:
            remark_list.append({"Field": "max cpu%", "Message": "Max CPU is required"})
        elif not isinstance(cpu_max, (int, float)) or not (0 < cpu_max <= 100):
            remark_list.append(
                {"Field": "max cpu%", "Message": "Max CPU must be a positive number and range between 1 to 100"})
        if not mem_max:
            remark_list.append({"Field": self.MAX_MEM_USED, "Message": "Max Mem used is required"})
        elif not isinstance(mem_max, (int, float)) or mem_max <= 0:
            remark_list.append({"Field": self.MAX_MEM_USED, "Message": "Max Mem used " + self.NUMBER_VALIDATION})
        if not net_max:
            remark_list.append({"Field": self.MAX_NW_BW, "Message": "Max Network BW is required"})
        elif not isinstance(net_max, (int, float)) or net_max <= 0:
            remark_list.append({"Field": self.MAX_NW_BW, "Message": "Max Network BW " + self.NUMBER_VALIDATION})
        remark_list = self.disk_fields_validation(disk_max, iops_max, remark_list)
        uuid = str(uuid).strip()
        if not uuid:
            remark_list.append({"Field": "uuid", "Message": "UUID is required"})
        if not pricing_model:
            remark_list.append({"Field": "pricingModel", "Message": "Pricing Model is required"})
        elif pricing_model.lower() in self.PRICING_MODEL and pricing_model.lower() in self.UNSUPPORTED_PRICING_MODEL:
            remark_list.append({"Field": "pricingModel", "Message": f"{pricing_model} pricing model is unsupported"})
        elif pricing_model.lower() not in self.PRICING_MODEL:
            remark_list.append({"Field": "pricingModel", "Message": f"{pricing_model} pricing model is invalid"})
        return remark_list

    def validate_utilization_metrics_fields(self,entry):
        """
        Validates pavg, uavg, p95, and u95 fields to ensure they are numbers (0-100).
        Returns a list of remarks.
        """
        remarks = []
        numeric_fields = ['pavg', 'uavg', 'p95', 'u95']
        for field in numeric_fields:
            value = entry.get(field, '')
            if value != '':
                if isinstance(value, (pd.Timestamp, datetime, date)):
                    remarks.append({
                        "Field": field,
                        "Message": f"Timestamps or dates are not allowed in {field}. Must be a number between 0 and 100."
                    })
                    continue
                try:
                    float_val = float(value)
                    if not (0 <= float_val <= 100):
                        remarks.append({
                            "Field": field,
                            "Message": f"Invalid value '{value}' in {field}. Must be a number between 0 and 100."
                        })
                except (ValueError, TypeError) as err:
                    self.log_message(LevelType.ERROR, f"Invalid value '{value}' in {field}. Must be a number between 0 and 100. : {str(err)}",ErrorCode=-1)
                    remarks.append({
                        "Field": field,
                        "Message": f"Invalid value '{value}' in {field}. Must be a number between 0 and 100."
                    })
        return remarks
    
    def auto_correct_eia(self, df, provider, regions_map):
        df['cloud_csp'] = df['cloud_csp'].apply(lambda x: provider.upper() if x != provider.upper() else x.upper())
        df['instance type'] = df['instance type'].str.lower()
        df['region'] = df.apply(
            lambda row: {r.lower(): r for r in regions_map.get(row['cloud_csp'].upper(), [])}.get(
                    str(row['region']).lower(),
                    self.default_regions.get(row['cloud_csp'].upper(), row['region'])
                ),
                axis=1
            )

        df['pricingModel'] = df['pricingModel'].apply(
            lambda x: x.lower() if str(x).lower() in self.valid_pricing_models else 'ondemand'
        )
        df["uuid"] = [str(uuid.uuid4()) for _ in range(len(df))]
        return df

    def remarks_input_eia(self, data, cloud_csp):
        remarks = []
        auto_corrected_data = []
        try:
            data = self.auto_correct_eia(data, cloud_csp, regions_map)
            for _, row in data.iterrows():
                corrected_entry = row.to_dict()
                uuid = corrected_entry['uuid']
                cpu_max = corrected_entry[self.MAX_CPU]
                mem_max = corrected_entry['max mem used']
                net_max = corrected_entry['max network bw']
                disk_max = corrected_entry['max disk bw used']
                iops_max = corrected_entry['max iops']
                remarks_list = self.validate_fields_eia(uuid, cpu_max, mem_max, net_max, disk_max, iops_max, corrected_entry[self.PRICE_MODEL])
                metric_remarks = self.validate_utilization_metrics_fields(corrected_entry)
                combined_remarks = remarks_list + metric_remarks
                remarks.append(combined_remarks)
                auto_corrected_data.append(corrected_entry)
            auto_corrected_data_df = pd.DataFrame(auto_corrected_data)
            return remarks, auto_corrected_data_df
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error for remarks_input_eia : {str(e)}")
            return "Failed", None

    def validate_file_input_data_eia(self, data_or_file, provider):
        try:
            required_headers = [
                'cloud_csp', 'region', 'instance type', 'max cpu%',
                'max mem used', 'max network bw', 'max disk bw used', 'max iops',
                'pricing model'
            ]

            if isinstance(data_or_file, pd.DataFrame):
                data = data_or_file.copy()
            elif hasattr(data_or_file, "read"):
                content = data_or_file.read()
                excel_data = BytesIO(content)
                data = pd.read_excel(excel_data)
            elif isinstance(data_or_file, list):
                if not data_or_file:
                    return {}, self.NO_DATA_IN_FILE, False
                data = pd.DataFrame(data_or_file)
            else:
                return {}, "Invalid input type", False

            if data.empty:
                return {}, self.NO_DATA_IN_FILE, False

            data.columns = data.columns.astype(str)
            data.columns = data.columns.str.strip().str.lower()

            # Keep a copy before any processing
            data_copy = data.copy()

            val, message = self.cca_eia_headers_validation(data, required_headers)
            if not val:
                return {}, message, False

            existing_cols = [col for col in required_headers if col in data.columns]
            data = data[existing_cols]

            if 'uuid' not in data.columns and 'uuid' in data_copy.columns:
                data['uuid'] = data_copy['uuid']
            if "uuid" not in data.columns:
                data['uuid'] = [str(uuid.uuid4()) for _ in range(len(data))]
            else:
                data['uuid'] = data['uuid'].apply(lambda x: str(uuid.uuid4()) if pd.isnull(x) or x == "" else x)

            if "pricing model" in data.columns:
                data = data.rename(columns={"pricing model": "pricingModel"})

            if "cloud_csp" in data.columns:
                data["cloud_csp"] = data["cloud_csp"].apply(lambda x: x.upper().strip() if isinstance(x, str) and x else "")

            remarks, data = self.remarks_input_eia(data, provider)
            if remarks == "Failed":
                return {}, "Unable to validate data", False
            data['Remarks'] = remarks

            output_dict = data.to_dict(orient='records')
            return output_dict, 'Uploaded file has some errors. Please click on VIEW DETAILS button to know more', True

        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error in validate_file_input_data_eia: {str(e)}")
            return {}, "Unable to process provided data", False

    def read_udf_file_data(self, input_data):
        try:
            if isinstance(input_data, list):
                data_dict = input_data
            elif isinstance(input_data, pd.DataFrame):
                data_dict = input_data.to_dict(orient='records')
            else:
                content = input_data.read()
                excel_data = BytesIO(content)
                data = pd.read_excel(excel_data, engine='openpyxl')
                data.columns = [col if 'Unnamed' not in col else '' for col in data.columns]
                data = data.fillna('')
                data_dict = data.to_dict(orient='records')
            return data_dict, 'Empty file uploaded for Self Perf Assessment'
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error in read_udf_file_data: {str(e)}")
            return {}, "Failed to validate data"


    def eia_data_read(self, data, udf_data, cloud_csp):
        data_out, udf_out, flag, message, udf_message = "", {}, False, "Invalid data provided", ""

        if (data is None or (isinstance(data, pd.DataFrame) and data.empty)) and \
        (udf_data is None or (isinstance(udf_data, pd.DataFrame) and udf_data.empty)):
            data_out, udf_out = [], []
            message = "Input data not provided"
            return data_out, udf_out, message, udf_message, flag

        if isinstance(data, pd.DataFrame) and not data.empty:
            data_out, _, flag = self.validate_file_input_data_eia(data, cloud_csp)
            message = "Input data validated sucessfully"
        else:
            data_out = []

        if isinstance(udf_data, pd.DataFrame) and not udf_data.empty:
            udf_out, udf_message = self.read_udf_file_data(udf_data)
            if not udf_out:
                message = f"Input data validated and Own metrics data: {udf_message}"
                self.log_message(LevelType.INFO, f"{message} for data : {udf_data}")
            else:
                message = "Input data and Own metrics data validated sucessfully"
                flag = True
        else:
            udf_out = []

        return data_out, udf_out, message, udf_message, flag
    def uuid_validation(self, data):

        if "uuid" not in data.columns:
            data["uuid"] = [str(uuid.uuid4()) for _ in range(len(data))]
            return data
        def fix_uuid(x):
            if x is None or pd.isna(x) or str(x).strip() == "":
                return str(uuid.uuid4())
            return str(x)

        data["uuid"] = data["uuid"].apply(fix_uuid)
        return data

    def auto_correct_cca(self, df, provider, regions_map):
        df['cloud_csp'] = df['cloud_csp'].apply(lambda x: provider.upper() if x != provider.upper() else x.upper())
        df['instance type'] = df['instance type'].str.lower()
        df['region'] = df.apply(
            lambda row: {r.lower(): r for r in regions_map.get(row['cloud_csp'].upper(), [])}.get(
                    str(row['region']).lower(),
                    self.default_regions.get(row['cloud_csp'].upper(), row['region'])
                ),
                axis=1
            )
        df['quantity'] = pd.to_numeric(df['quantity'], errors='coerce')
        df['quantity'] = df['quantity'].apply(lambda x: 1 if pd.isna(x) or x <= 0 else x)
        df['monthly utilization (hourly)'] = df.apply(
                    lambda row: row['monthly utilization (hourly)']
                        if isinstance(row['monthly utilization (hourly)'], (int, float)) and
                            row['monthly utilization (hourly)'] > 0 and
                            row['monthly utilization (hourly)'] <= row['quantity'] * 730
                        else row['quantity'] * 730,
                    axis=1
                )

        df['pricingModel'] = df['pricingModel'].apply(
            lambda x: x.lower() if str(x).lower() in self.valid_pricing_models else 'ondemand'
        )

        return df
    
    def validate_input_data_cca(self, data, provider):
        try:
            required_headers = ["cloud", "region", "size", "quantity", "total number of hours per month", self.PRICEMODEL]
            data_copy = data.copy()
            if not all(header in data.columns for header in required_headers):
                val, message = self.cca_eia_headers_validation(data, required_headers)
                if not val:
                    self.log_message(LevelType.ERROR, f"message : {message}",ErrorCode=-1)
                    return {}, message, False
            
            if "uuid" not in data.columns and "uuid" in data_copy.columns:
                data["uuid"] = data_copy["uuid"].fillna("").astype(str)

            data = self.uuid_validation(data)
                
            data = data.rename(columns={
                "cloud": "cloud_csp",
                "region": "region",
                "size": self.INSTANCE_TYPE,
                "quantity": "quantity",
                "total number of hours per month": "monthly utilization (hourly)",
                self.PRICEMODEL: "pricingModel"
            })

            if "cloud_csp" in data.columns:
                data["cloud_csp"] = data["cloud_csp"].apply(
                    lambda x: x.upper().strip() if isinstance(x, str) and x else ""  # leave blanks for null/empty
                )
            data = self.auto_correct_cca(data, provider, regions_map)
            filter_columns = ['cloud_csp', 'region', self.INSTANCE_TYPE, 'quantity', "monthly utilization (hourly)", 'pricingModel']
            data = data[~data[filter_columns].apply(lambda row: all(x in ["", []] for x in row), axis=1)]
            output_dict = data.to_dict(orient='records')
            return output_dict, "Data validated successfully", True
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error in validate_input_data_cca: {str(e)}",ErrorCode=-1)
            return {}, "Unable to process provided file", False

    def process_batch(self, portfolio_id, batch_id, cloud_csp, portfolio_doc, rec_progress_id, app_name):
        try:
            instances = list(self.current_instance_collection.find({
                "portfolio_id": portfolio_id,
                "batch_id": batch_id
            }))
            df_all = pd.DataFrame(instances)
            columns_to_drop = ['_id', 'created_at', 'uploaded_date', 'batch_id', "portfolio_id"]
            existing_columns_to_drop = [col for col in columns_to_drop if col in df_all.columns]
            if existing_columns_to_drop:
                df_all = df_all.drop(columns=existing_columns_to_drop)
                df_all = df_all.map(lambda x: x.strip() if isinstance(x, str) else x)
            if df_all.empty:
                self.log_message(LevelType.ERROR, f"No instance data found for portfolio {portfolio_id} batch {batch_id}")
                return False, None
            udf_data = []
            transformed_non_cloud_rows = []
            headroom = portfolio_doc.get('headroom', 20)
            if app_name.upper() == "CCA":
                final_df_again, _, flag = self.validate_input_data_cca(df_all, cloud_csp)
                if not flag:
                    self.recommendation_tracking_collection.update_one(
                        {"_id": rec_progress_id},
                        {"$set": {"recommendation_status": RecommendationStatus.FAILED}}
                    )
                    return False, None
                data_json = self.sanitize_for_json(final_df_again)
                cloud_rows = [row for row in data_json if str(row.get("cloud_csp", "")).lower() == cloud_csp.lower()]
                non_cloud_rows = [row for row in data_json if str(row.get("cloud_csp", "")).lower() != cloud_csp.lower()]

                data = self.transform_rows(cloud_rows)
                transformed_non_cloud_rows = self.transform_rows(non_cloud_rows)
                if transformed_non_cloud_rows:
                    for row in transformed_non_cloud_rows:
                        row['Remarks'] = [
                            {
                                "Field": "cloud_csp",
                                "Message": f"Invalid cloud provider for request. Expected {cloud_csp.upper()}"
                            }
                        ]
                if not data:
                    self.log_message(LevelType.WARNING, f"No matching records found for cloud provider {cloud_csp.upper()} for batch : {batch_id}.")
                message = "Data validated successfully."
            else:
                udf_raw_data = portfolio_doc.get("udf", [])
                if not udf_raw_data:
                    self.log_message(LevelType.INFO, f"No UDF data available for portfolio {portfolio_id}")
                    df_udf = pd.DataFrame()
                else:
                    df_udf = pd.DataFrame(udf_raw_data)
                data, udf_data, message, _, flag = self.eia_data_read(df_all, df_udf, cloud_csp)
                
                if not flag:
                    self.recommendation_tracking_collection.update_one(
                        {"_id": rec_progress_id},
                        {"$set": {"recommendation_status": RecommendationStatus.FAILED}}
                    )
                    return False, None

            self.log_message(LevelType.INFO, f"Validation for batch {batch_id}: {message} portfolio {portfolio_id}")
            error_records = []
            if data:
                error_records = [record for record in data if record.get('Remarks')]
                filtered_data = [record for record in data if not record.get('Remarks')]
                if filtered_data:
                    _, status = self.process_cost_advise(filtered_data, headroom, portfolio_doc, app_name, batch_id, portfolio_id, udf_data)
                else:
                    self.log_message(LevelType.INFO, f"No valid records found for portfolio {portfolio_id} for batch_id {batch_id}")
                    status = True
            else:
                self.log_message(LevelType.INFO, f"No data found for portfolio {portfolio_id} for batch_id {batch_id}")
                status = True
            if transformed_non_cloud_rows:
                error_records.extend(transformed_non_cloud_rows)

            self.log_message(LevelType.INFO, f"Cost advise status for portfolio {portfolio_id} for batch {batch_id}: {status}")

            new_status = "COMPLETED" if status else "FAILED"
            self.recommendation_tracking_collection.update_one(
                {"_id": rec_progress_id},
                {"$set": {"recommendation_status": new_status}}
            )
            self.log_message(LevelType.INFO, f"For portfolio id {portfolio_id} Batch {batch_id} recommendation_status updated to {new_status}")
            return status, error_records
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error processing batch {batch_id} for portfolio {portfolio_id}: {str(e)}")
            return False, None
    
    #Start of CCA Excel generation
    def style_cell(self, cell, fill=None, font=None, align=None, border=None):
        if fill:
            cell.fill = fill
        if font:
            cell.font = font
        if align:
            cell.alignment = align
        if border:
            cell.border = border


    def set_col_widths(self, ws):
        for col in ws.columns:
            max_len = 0
            col_letter = get_column_letter(col[0].column)

            for cell in col:
                try:
                    if cell and cell.value:
                        max_len = max(max_len, len(str(cell.value)))
                except:
                    pass

            ws.column_dimensions[col_letter].width = max_len + 2
        ws.column_dimensions["A"].width = 15

    def create_main_header(self, ws, total_columns, header_fill):
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=4)
        cell = ws["A1"]
        cell.value = "EPYC Cloud Cost Advisory Recommendations"

        self.style_cell(cell,
                fill=header_fill,
                font=Font(name="Calibri", size=11, bold=True, color="FFFFFF"),
                align=Alignment(horizontal="center", vertical="center"))

        for col in range(1, total_columns + 1):
            ws.cell(row=1, column=col).fill = header_fill
            ws.cell(row=1, column=col).border = self.WHITE_BORDER

        # Increase row 1 height
        ws.row_dimensions[1].height = 25

    def create_table_headers(self, ws, header_fill):
        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        header_align = Alignment(wrap_text=True, horizontal="center", vertical="center")

        main_cols = [
            "Region", "Current Instance", "Current Monthly Cost ($)", "Current Annual Cost ($)",
            "UUID/Instance Name", "Cloud", "Quantity", "Pricing Model",
            "Current vCPU(s)", "Remark"
        ]

        for idx, text in enumerate(main_cols, start=1):
            ws.merge_cells(start_row=2, start_column=idx, end_row=3, end_column=idx)
            cell = ws.cell(2, idx, value=text)
            self.style_cell(cell, header_fill, header_font, header_align, border=self.WHITE_BORDER)
            # Apply border to row 3 as well for merged cells
            ws.cell(3, idx).fill = header_fill
            ws.cell(3, idx).border = self.WHITE_BORDER

        rec_groups = [
            (11, "Hourly Cost Optimization"),
            (18, "Modernize"),
            (25, "Modernize & Downsize")
        ]

        for col, title in rec_groups:
            ws.merge_cells(start_row=2, start_column=col, end_row=2, end_column=col + 6)
            cell = ws.cell(2, col, value=title)
            self.style_cell(cell, header_fill, header_font, header_align, border=self.WHITE_BORDER)
            # Apply border to all cells in the merged range
            for c in range(col, col + 7):
                ws.cell(2, c).fill = header_fill
                ws.cell(2, c).border = self.WHITE_BORDER

        detailed_headers = [
            "Instance", "vCPU(s)", "Monthly Cost ($)", "Annual Cost ($)",
            "Annual Savings ($)", "Savings (%)", "Performance Improvement"
        ] * 3

        for col, text in enumerate(detailed_headers, start=11):
            cell = ws.cell(3, col, value=text)
            self.style_cell(cell, header_fill, header_font, header_align, border=self.WHITE_BORDER)
        
        # Adjust row heights for headers
        ws.row_dimensions[2].height = 20
        ws.row_dimensions[3].height = 20

    def _has_green_recommendation(self, item):
        """Returns True if any recommendation has positive annual savings (green cell)."""
        recs = item.get("data", {}).get("recommendations", [])
        for rec in recs:
            try:
                if float(rec.get("annualSavings", 0)) > 0:
                    return True
            except (ValueError, TypeError):
                pass
        return False

    def _sort_data_green_first(self, data):
        """Sort data so rows with green recommendations come first, white (no recommendations) last."""
        return sorted(data, key=lambda x: (0 if self._has_green_recommendation(x) else 1))

    def write_data_rows(self, ws, data, green_fill):
        for item in data:
            current = item["data"]["currentPlatform"]
            recs = item["data"]["recommendations"]

            row = [
                current["zone"], current["instanceType"], safe_round(current["monthlyCost"]),
                safe_round(current["annualCost"]), item["id"], current["cspProvider"],
                current["numberOfInstances"], current["pricingModel"],
                safe_round(current["vCPU"]), current["status"]
            ]

            for rec in recs:
                row.extend([
                    rec["instanceType"], safe_round(rec["vCPU"]), safe_round(rec["monthlyCost"]),
                    safe_round(rec["totalCost"]), format_savings(rec["annualSavings"]),
                    safe_round(rec["savingsInPercentage"]), safe_round(rec["perf"])
                ])

            ws.append(row)

            last_row = ws.max_row
            for col in range(1, len(row) + 1):
                ws.cell(row=last_row, column=col).alignment = Alignment(horizontal="left", vertical="center")
                
                # Apply thousand separator formatting to Cost, Annual Cost, Annual Savings cols
                if col in [3, 4] or (col >= 11 and (col - 11) % 7 in [2, 3, 4]):
                    ws.cell(row=last_row, column=col).number_format = '#,##0.00'

            for idx, rec in enumerate(recs):
                try:
                    if float(rec["annualSavings"]) > 0:
                        instance_col = 11 + (idx * 7)
                        ws.cell(last_row, instance_col).fill = green_fill
                except:
                    pass

    def write_grand_total(self, ws, grand_total):
        row = [
            "Grand Total", "", safe_round(grand_total["Current Monthly Cost"]),
            safe_round(grand_total["Annual Cost"]), "", "", grand_total["Number of Instances"],
            "", "", "", "", "", safe_round(grand_total["Monthly Cost I"]),
            safe_round(grand_total["Annual Cost I (perf scaled)"]),
            safe_round(grand_total["Annual Savings I"]), safe_round(grand_total["hSavingsInPercentage"]),
            safe_round(grand_total["Perf Enhancement I"]), "", "",
            safe_round(grand_total["Monthly Cost II"]), safe_round(grand_total["Annual Cost II (perf scaled)"]),
            safe_round(grand_total["Annual Savings II"]), safe_round(grand_total["mSavingsInPercentage"]),
            safe_round(grand_total["Perf Enhancement II"]), "", "",
            safe_round(grand_total["Monthly Cost III"]),
            safe_round(grand_total["Annual Cost III (perf scaled)"]),
            safe_round(grand_total["Annual Savings III"]), safe_round(grand_total["mdSavingsInPercentage"]),
            safe_round(grand_total["Perf Enhancement III"]),
        ]

        ws.append(row)
        last_row = ws.max_row

        for col in range(1, len(row) + 1):
            cell = ws.cell(row=last_row, column=col)
            cell.font = Font(name="Calibri", size=11)
            cell.alignment = Alignment(horizontal="left", vertical="center")
            cell.border = Border()
            
            # Apply thousand separator formatting to numeric total columns
            if col in [3, 4] or (col >= 13 and (col - 13) % 7 in [0, 1, 2]):
                cell.number_format = '#,##0.00'
    
    def add_note(self, ws):
        note_row = ws.max_row + 2
        ws.merge_cells(f"A{note_row}:F{note_row}")

        cell = ws[f"A{note_row}"]
        cell.value = "Note: Green color instances indicate positive savings."

        self.style_cell(cell,
                font=Font(name="Calibri", size=11, color="0000FF"),
                align=Alignment(horizontal="left", vertical="center"))
        
        # Add second note
        note_row2 = note_row + 1
        ws.merge_cells(f"A{note_row2}:F{note_row2}")
        
        cell2 = ws[f"A{note_row2}"]
        cell2.value = "Sizing instances - matching resources to actual demand (if Applicable)"
        
        self.style_cell(cell2,
                font=Font(name="Calibri", size=11, color="0000FF"),
                align=Alignment(horizontal="left", vertical="center"))

    def create_summary_sheet(self, wb, grand_total, header_fill):
        ws = wb.create_sheet("Total Annual Savings")

        white = Font(bold=True, color="FFFFFF")
        center = Alignment(horizontal="center", vertical="center")

        ws.merge_cells("A1:A2")
        self.style_cell(ws["A1"], header_fill, white, center, border=self.WHITE_BORDER)
        ws["A1"].value = "Current Cost"

        groups = {"B": "Hourly Cost Optimization", "D": "Modernize", "F": "Modernize & Downsize"}

        for col, title in groups.items():
            ws.merge_cells(f"{col}1:{chr(ord(col)+1)}1")
            cell = ws[f"{col}1"]
            cell.value = title
            self.style_cell(cell, header_fill, white, center, border=self.WHITE_BORDER)

        sub = ["Total Cost", "Total Savings"] * 3
        for idx, col in enumerate(["B","C","D","E","F","G"]):
            self.style_cell(ws[f"{col}2"], header_fill, white, center, border=self.WHITE_BORDER)
            ws[f"{col}2"].value = sub[idx]

        data = [
            grand_total["Annual Cost"],
            grand_total["Annual Cost I (perf scaled)"], grand_total["Annual Savings I"],
            grand_total["Annual Cost II (perf scaled)"], grand_total["Annual Savings II"],
            grand_total["Annual Cost III (perf scaled)"], grand_total["Annual Savings III"],
        ]

        for col, value in zip(["A","B","C","D","E","F","G"], data):
            cell = ws[f"{col}3"]
            cell.value = float(value)
            cell.number_format = '"$ "#,##0.00'

        self.set_col_widths(ws)

    def create_legal_sheet(self, wb):
        ws = wb.create_sheet("Legal Disclaimer")

        ws.merge_cells("A1:Q1")
        cell = ws["A1"]
        cell.value = (
            "Disclaimer: THE MATERIALS PROVIDED THROUGH THIS TOOL ARE PROVIDED 'AS IF', WITHOUT WARRANTY OF ANY KIND, EXPRESS OR IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY, FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM, OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE SOFTWARE."
        )
        cell.font = Font(name="Calibri", size=11)
        cell.alignment = Alignment(wrap_text=False)

        ws.row_dimensions[1].height = 25

        ws.merge_cells("A3:G3")
        ws["A3"] = "Cloud instance recommendation generated using AMD EPYC Cloud Cost Advisor"

        ws.merge_cells("A5:G5")
        current_year = datetime.now().year
        ws["A5"] = f"Copyright - {current_year} Advanced Microdevices Inc."

        ws.merge_cells("A7:H7")
        ws["A7"] = "For terms: https://www.amd.com/en/legal/copyright.html"

        for col in "ABCDEFGHIJKLMNOPQ":
            ws.column_dimensions[col].width = 20

    def generate_excel_from_json(self, input_json, output_path, app_name, user_email, excel_file, portfolio_id):
        try:
            # log_message(LevelType.INFO, f"Generating Excel from JSON: {input_json}")
            self.log_message(LevelType.INFO, f"Output path: {output_path}")
            
            data = input_json["data"]
            grand_total = input_json["grandTotal"]

            wb = Workbook()
            ws = wb.active
            ws.title = "Recommended-Instance"

            header_fill = PatternFill(start_color="000000", end_color="000000", fill_type="solid")
            green_fill = PatternFill(start_color="00FF00", end_color="00FF00", fill_type="solid")

            self.create_main_header(ws, 31, header_fill)
            self.create_table_headers(ws, header_fill)

            ws.freeze_panes = "E4"

            sorted_data = self._sort_data_green_first(data)
            self.write_data_rows(ws, sorted_data, green_fill)
            self.write_grand_total(ws, grand_total)
            self.add_note(ws)
            self.set_col_widths(ws)

            self.create_summary_sheet(wb, grand_total, header_fill)
            self.create_legal_sheet(wb)

            excel_buffer = BytesIO()
            wb.save(excel_buffer)
            excel_buffer.seek(0)
            wb.save(output_path)
            output_s3_path=self.upload_file_to_s3(
                file_bytes=excel_buffer.getvalue(),
                app_name=app_name,
                user_email=user_email,
                file_name=excel_file,
                sub_folder="output"
            )
            if output_s3_path:
                self.portfolio_collection.update_one(
                    {"_id": ObjectId(portfolio_id)},
                    {
                        "$set": {
                            "advice_s3_key": output_s3_path
                        }
                    }
                )
            self.log_message(LevelType.INFO, "Excel successfully generated!")
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error generating Excel: {e}", ErrorCode=-1)

    #End of CCA excel generation

    #Start of EIA Excel Generation
    def apply_cell_style(self, cell, fill=None, font=None, alignment=None, border=None):
        """Applies multiple styles to an openpyxl cell."""
        if fill:
            cell.fill = fill
        if font:
            cell.font = font
        if alignment:
            cell.alignment = alignment
        if border:
            cell.border = border


    def safe_float_eia(self, value):
        """Safely converts a value to float, defaulting to 0.0 on error."""
        try: return float(value)
        except (ValueError, TypeError): return 0.0


    def create_disclaimer_sheet(self, wb):
        ws = wb.create_sheet("Disclaimer")
        
        ws.column_dimensions['A'].width = 150
        
        ws.cell(1, 1, value="Disclaimer: THE MATERIALS PROVIDED THROUGH THIS TOOL ARE PROVIDED 'AS IF', WITHOUT WARRANTY OF ANY KIND, EXPRESS OR IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY, FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM, OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE SOFTWARE.")
        self.apply_cell_style(ws.cell(1, 1), font=self.DISCLAIMER_TITLE_FONT, alignment=self.LEFT_ALIGNMENT)
        
        ws.cell(3, 1, value="EPYC Cloud Instance Advisory recommendations generated using AMD EPYC Cloud Instance Advisor")
        self.apply_cell_style(ws.cell(3, 1), font=self.DISCLAIMER_FONT, alignment=self.LEFT_ALIGNMENT)

        current_year = datetime.now().year
        ws.cell(5, 1, value=f"Copyright - {current_year} Advanced Microdevices Inc.")
        self.apply_cell_style(ws.cell(5, 1), font=self.DISCLAIMER_FONT, alignment=self.LEFT_ALIGNMENT)
        
        ws.cell(7, 1, value="For Terms of Use / Copyrights: please refer https://www.amd.com/en/legal/copyright.html")
        self.apply_cell_style(ws.cell(7, 1), font=self.DISCLAIMER_FONT, alignment=self.LEFT_ALIGNMENT)
        
        return ws


    def create_total_savings_sheet(self, wb, grand_total):
        ws = wb.create_sheet("Total Annual Savings")
        ws.row_dimensions[1].height = 20
        ws.row_dimensions[2].height = 30
        ws.row_dimensions[3].height = 20
        
        COL_CURRENT_COST = 1
        COL_OPT_COST = 2
        COL_OPT_SAVINGS = 3
        COL_GOOD_COST = 4
        COL_GOOD_SAVINGS = 5
        TOTAL_COLS = 5

        
        ws.merge_cells(start_row=1, start_column=COL_CURRENT_COST, end_row=2, end_column=COL_CURRENT_COST)
        current_cost_cell = ws.cell(1, COL_CURRENT_COST, value="Current Cost")
        self.apply_cell_style(current_cost_cell, self.BLACK_HEADER_FILL, self.WHITE_BOLD_FONT, self.CENTER_ALIGNMENT, Border(left=self.WHITE_SIDE, right=self.WHITE_SIDE, top=self.WHITE_SIDE, bottom=self.WHITE_SIDE))
        
        ws.merge_cells(start_row=1, start_column=COL_OPT_COST, end_row=1, end_column=COL_OPT_SAVINGS)
        optimal_cell = ws.cell(1, COL_OPT_COST, value="OPTIMAL")
        self.apply_cell_style(optimal_cell, self.BLACK_HEADER_FILL, self.WHITE_BOLD_FONT, self.CENTER_ALIGNMENT, Border(left=self.WHITE_SIDE, right=self.WHITE_SIDE, top=self.WHITE_SIDE, bottom=self.WHITE_SIDE))
        
        ws.merge_cells(start_row=1, start_column=COL_GOOD_COST, end_row=1, end_column=COL_GOOD_SAVINGS)
        good_cell = ws.cell(1, COL_GOOD_COST, value="GOOD")
        self.apply_cell_style(good_cell, self.BLACK_HEADER_FILL, self.WHITE_BOLD_FONT, self.CENTER_ALIGNMENT, Border(left=self.WHITE_SIDE, right=self.WHITE_SIDE, top=self.WHITE_SIDE, bottom=self.WHITE_SIDE))

        header_row_2 = {
            COL_OPT_COST: "Total Cost", COL_OPT_SAVINGS: "Total Savings",
            COL_GOOD_COST: "Total Cost", COL_GOOD_SAVINGS: "Total Savings"
        }
        for col_idx, header_text in header_row_2.items():
            cell = ws.cell(2, col_idx, value=header_text)
            self.apply_cell_style(cell, self.BLACK_HEADER_FILL, self.WHITE_BOLD_FONT, self.CENTER_ALIGNMENT, Border(left=self.WHITE_SIDE, right=self.WHITE_SIDE, top=self.WHITE_SIDE, bottom=self.WHITE_SIDE))

        # Multiply by 12 to convert monthly to annual
        data_map = {
            COL_CURRENT_COST: safe_round(grand_total.get("Current Monthly Price", 0.0) * 12),
            COL_OPT_COST: safe_round(grand_total.get("Monthly Price I", 0.0) * 12),
            COL_OPT_SAVINGS: safe_round(grand_total.get("Monthly Savings I", 0.0) * 12), 
            COL_GOOD_COST: safe_round(grand_total.get("Monthly Price II", 0.0) * 12),
            COL_GOOD_SAVINGS: safe_round(grand_total.get("Monthly Savings II", 0.0) * 12),
        }

        data_row = 3
        for col_idx in range(1, TOTAL_COLS + 1):
            cell = ws.cell(data_row, col_idx, value=data_map.get(col_idx))
            
            # Changed from GRAY_FILL to WHITE_FILL and alignment from RIGHT to CENTER
            cell.fill = self.WHITE_FILL
            cell.number_format = '"$"#,##0.00'
            self.apply_cell_style(cell, alignment=self.CENTER_ALIGNMENT, font=self.BOLD_NORMAL_FONT)
            
            ws.column_dimensions[get_column_letter(col_idx)].width = 18

        return ws

    def create_main_sheet(self, wb, data, grand_total, safety_margin):
        """Creates and populates the main sheet with the requested 5-row structure."""
        
        if 'Sheet' in wb.sheetnames:
            del wb['Sheet']
            
        ws = wb.create_sheet(title="Recommended-Instance", index=0)

        TOTAL_COLUMNS = 26
        

        ROW_TITLE_START = 1 
        ROW_TITLE_END = 2 
        ROW_GROUP_HEADER = 3 
        ROW_DETAIL_HEADER = 4 
        ROW_DATA_START = 5 

        ws.row_dimensions[ROW_TITLE_START].height = 15
        ws.row_dimensions[ROW_TITLE_END].height = 10 
        ws.row_dimensions[ROW_GROUP_HEADER].height = 15
        ws.row_dimensions[ROW_DETAIL_HEADER].height = 35
        
        

        ws.merge_cells(start_row=ROW_TITLE_START, start_column=1, end_row=ROW_TITLE_END, end_column=TOTAL_COLUMNS) 
        cell_a1 = ws["A1"]
        cell_a1.value = f"EPYC Cloud Instance Advisory Recommendations - ( Operational Safety Margin : {safety_margin} )%"
        self.apply_cell_style(
            cell_a1, 
            self.BLACK_HEADER_FILL, 
            self.WHITE_BOLD_FONT, 
            Alignment(horizontal="left", vertical="center"), 
            self.TITLE_BORDER
        )


        for row_idx in [ROW_GROUP_HEADER, ROW_DETAIL_HEADER]:
            for col_idx in range(1, TOTAL_COLUMNS + 1):
                ws.cell(row_idx, col_idx).fill = self.BLACK_HEADER_FILL
                ws.cell(row_idx, col_idx).border = Border()
        
        
        ws.merge_cells(start_row=ROW_GROUP_HEADER, start_column=1, end_row=ROW_GROUP_HEADER, end_column=8)
        self.apply_cell_style(ws[f"A{ROW_GROUP_HEADER}"], self.BLACK_HEADER_FILL, self.WHITE_BOLD_FONT, self.CENTER_ALIGNMENT, self.R5_BOTTOM_BORDER) 


        ws.merge_cells(start_row=ROW_GROUP_HEADER, start_column=11, end_row=ROW_GROUP_HEADER, end_column=18)
        cell_k3 = ws[f"K{ROW_GROUP_HEADER}"]
        cell_k3.value = "Optimal" 


        self.apply_cell_style(cell_k3, self.BLACK_HEADER_FILL, self.WHITE_BOLD_FONT, self.CENTER_ALIGNMENT, self.R1_BOTTOM_BORDER) 


        ws.merge_cells(start_row=ROW_GROUP_HEADER, start_column=19, end_row=ROW_GROUP_HEADER, end_column=26)
        ws[f"S{ROW_GROUP_HEADER}"].value = "Good" 
        self.apply_cell_style(ws[f"S{ROW_GROUP_HEADER}"], self.BLACK_HEADER_FILL, self.WHITE_BOLD_FONT, self.CENTER_ALIGNMENT, self.R1_BOTTOM_BORDER)
        
        

        header_map_final = [
            # Current (1-8)
            (1, "Region", "region", "current"),
            (2, "Instance Type", "type", "current"), (3, "Cost ($)", "cost", "current"), (4, "Power (kW)", "power", "current"), 
            (5, "Carbon (kgCO2eq)", "carbon", "current"), (6, "UUID/Instance Name", "id", "top"), (7, "Cloud", "csp", "top"), 
            (8, "Pricing Model", "pricingModel", "current"), 
            
            (9, "vCPU(s)", "vCPU", "current_vcpus_for_remark"), 
            (10, "Remark", "status", "current"), 
            (11, "Instance Type", "type", 0), # COL 11 is Instance Type (Optimal)
            (12, "vCPU(s)", "vCPU", 0), (13, "Cost ($)", "cost", 0), 
            (14, "Power (kW)", "power", 0), (15, "Carbon (kgCO2eq)", "carbon", 0), (16, "Monthly Savings ($)", "monthlySavings", 0), 
            (17, "Performance Improvement", "perf", 0), (18, "Untapped Capacity", "untappedCapacity", 0),
            
            (19, "Instance Type", "type", 1), # COL 19 is Instance Type (Good)
            (20, "vCPU(s)", "vCPU", 1), (21, "Cost ($)", "cost", 1), 
            (22, "Power (kW)", "power", 1), (23, "Carbon (kgCO2eq)", "carbon", 1), (24, "Monthly Savings ($)", "monthlySavings", 1), 
            (25, "Performance Improvement", "perf", 1), (26, "Untapped Capacity", "untappedCapacity", 1),
        ]

        for col_idx, header_text, *rest in header_map_final:
            cell = ws.cell(ROW_DETAIL_HEADER, col_idx, value=header_text)
            

            self.apply_cell_style(cell, self.BLACK_HEADER_FILL, self.WHITE_BOLD_FONT, self.HEADER_ALIGNMENT, self.R4_HEADER_BORDER)
                
    

        
        row_num = ROW_DATA_START
        for item in data:
            current = item["data"]["currentPlatform"]
            recs = item["data"]["recommendations"]
            row_data = [""] * TOTAL_COLUMNS 
            
            rec1_savings = self.safe_float_eia(recs[0]["monthlySavings"]) if len(recs) > 0 and recs[0].get("monthlySavings") else 0.0
            rec2_savings = self.safe_float_eia(recs[1]["monthlySavings"]) if len(recs) > 1 and recs[1].get("monthlySavings") else 0.0

            # Populate Data Row
            for col_idx, _, key, source_type in header_map_final:
                value = "-"
                
                if source_type == "top":
                    value = str(item.get(key, "-"))
                elif source_type == "current":
                    val = current.get(key, "-")
                    if key in ["cost", "power", "carbon", "vCPU"]:
                        value = safe_round(val)
                    else:
                        value = str(val)
                elif source_type == "current_vcpus_for_remark":
                    value = safe_round(current.get("vCPU", "-"))
                elif isinstance(source_type, int) and len(recs) > source_type:
                    val = recs[source_type].get(key, "-")
                    if key in ["cost", "power", "carbon", "vCPU", "monthlySavings", "perf", "untappedCapacity"]:
                         value = safe_round(val)
                    else:
                         value = str(val)
                
                row_data[col_idx - 1] = value

            ws.append(row_data)
            
            data_row_idx = row_num 
            for col_idx in range(1, TOTAL_COLUMNS + 1):
                cell = ws.cell(data_row_idx, col_idx)
                
                is_green = False
                # Check for Optimal Instance Type (Col 11)
                if col_idx == 11 and rec1_savings > 0:
                    is_green = True
                # Check for Good Instance Type (Col 19)
                elif col_idx == 19 and rec2_savings > 0:
                    is_green = True
                
                if is_green:
                    cell.fill = self.LIGHT_GREEN_FILL
                
                if col_idx in [3, 4, 5, 13, 14, 15, 16, 21, 22, 23, 24]:
                    cell.number_format = '#,##0.00'
                    
                cell.alignment = self.LEFT_ALIGNMENT
                cell.border = Border()
                    
            row_num += 1

        gt_row = [""] * TOTAL_COLUMNS 
        
        gt_row[0] = "Grand Total"

        gt_map = {
            3: "Current Monthly Price", 4: "Current Instance Energy Consumption (kwh)", 5: "Current Instance Emission",
            13: "Monthly Price I", 14: "Instance Energy Consumption I (kwh)", 15: "Instance Emission I", 16: "Monthly Savings I", 17: "Perf Enhancement I", 18: "Untapped Capacity I",
            21: "Monthly Price II", 22: "Instance Energy Consumption II (kwh)", 23: "Instance Emission II", 24: "Monthly Savings II", 25: "Perf Enhancement II", 26: "Untapped Capacity II",
        }
        
        for col_idx, key in gt_map.items():
            if key in grand_total:
                gt_row[col_idx - 1] = safe_round(grand_total[key])

        ws.append(gt_row) 
        last_row = ws.max_row 
        
        for col_idx in range(1, TOTAL_COLUMNS + 1):
            cell = ws.cell(last_row, col_idx)
            cell.font = self.NORMAL_FONT
            cell.alignment = self.LEFT_ALIGNMENT

            if col_idx in gt_map.keys():
                try:
                    # Apply number formatting 
                    if col_idx in [3, 13, 16, 21, 24]:  # Cost and Savings columns
                        cell.number_format = '"$"#,##0.00'
                    else:
                        cell.number_format = '#,##0.00'
                except:
                    pass

            if col_idx == 18 or col_idx == 26:
                cell.border = Border()


        # --- Note/Legend ---
        note_text = "Note : Green color instances indicate positive savings."
        note_row = ws.max_row + 2
        ws.merge_cells(f'A{note_row}:F{note_row}')
        note_cell = ws[f'A{note_row}']
        note_cell.value = note_text
        self.apply_cell_style(note_cell, font=self.NOTE_FONT, alignment=self.LEFT_ALIGNMENT)
        
        # Add second note
        note_row2 = note_row + 1
        ws.merge_cells(f'A{note_row2}:F{note_row2}')
        note_cell2 = ws[f'A{note_row2}']
        note_cell2.value = "Sizing instances - matching resources to actual demand (if Applicable)"
        self.apply_cell_style(note_cell2, font=self.NOTE_FONT, alignment=self.LEFT_ALIGNMENT)
        
        ws.freeze_panes = "F5"

        return ws


    def adjust_column_widths(self, ws):
        """
        Adjusts column widths for the main sheet based on content, 
        applying a maximum width limit and skipping the 'Instance Type' columns.
        """
        for col in ws.columns:
            col_idx = col[0].column
            col_letter = get_column_letter(col_idx)

            # Skip Instance Type columns to keep their fixed small width
            if col_idx in [2, 11, 19]: 
                continue 

            max_len = 0
            
            # Check from the header row (Row 4)
            for cell in col:
                if cell.value and cell.row >= 4: 
                    # Use a factor of 1.2 to account for font size/padding, but use 1.0 for UUID/Instance Name (Col 6)
                    # as it often contains long non-readable strings that can be cut off.
                    multiplier = 1.0 if col_idx == 6 else 1.2
                    current_len = len(str(cell.value)) * multiplier
                    max_len = max(max_len, current_len)

            if max_len > 0:
                # Apply auto-adjustment, respecting the MAX_COLUMN_WIDTH
                adjusted_width = max(16, max_len)
                ws.column_dimensions[col_letter].width = min(self.MAX_COLUMN_WIDTH, adjusted_width)


    def generate_excel_report(self, data, grand_total, output_path, headroom, app_name, user_email, excel_file, portfolio_id):
        """Main function to generate the complete Excel report with three tabs."""
        try:
            if not data and not grand_total:
                self.log_message(LevelType.ERROR, "Aborting Excel generation due to data loading errors.", ErrorCode=-1)
                return

            wb = Workbook()
            
            ws_main = self.create_main_sheet(wb, data, grand_total, headroom)
            ws_totals = self.create_total_savings_sheet(wb, grand_total)
            ws_disclaimer = self.create_disclaimer_sheet(wb)
            
            wb.active = ws_main
            
            # Adjust widths only after all content is written
            self.adjust_column_widths(ws_main)

            # wb.save(output_path)
            excel_buffer = BytesIO()
            wb.save(excel_buffer)
            excel_buffer.seek(0)
            wb.save(output_path)
            output_s3_path=self.upload_file_to_s3(
                file_bytes=excel_buffer.getvalue(),
                app_name=app_name,
                user_email=user_email,
                file_name=excel_file,
                sub_folder="output"
            )
            if output_s3_path:
                self.portfolio_collection.update_one(
                    {"_id": ObjectId(portfolio_id)},
                    {
                        "$set": {
                            "advice_s3_key": output_s3_path
                        }
                    }
                )
            self.log_message(LevelType.INFO, f"Excel successfully generated at: {output_path}")
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error generating or saving Excel file: {e}", ErrorCode=-1)
    #End EIA excel generation
    #End of xlsx generation code

    #start of PPT generation code

    # --------- Currency Formatter ----------
    def format_currency(self, num):
        try:
            num = float(num)
        except Exception:
            return str(num)

        abs_num = abs(num)
        if abs_num >= 1_000_000_000:
            return f"${num / 1_000_000_000:.1f}B"
        elif abs_num >= 1_000_000:
            return f"${num / 1_000_000:.1f}M"
        elif abs_num >= 1_000:
            return f"${num / 1_000:.1f}K"
        else:
            return f"${num:.2f}"

    def extract_customer_from_email(self, email: str) -> str:
        try:
            if '@' in email and '.' in email.split('@')[-1]:
                return email.split('@')[1].split('.')[0].upper()
            return email.upper()
        except Exception as e:
            self.log_message(LevelType.ERROR, f"extract_customer_from_email: {str(e)}")
            return "UNKNOWN"

    def resolve_input_data(self, excel_or_df, sheet_name='Recommended-Instance', header_levels=[0,1,2]):
        try:
            if isinstance(excel_or_df, pd.DataFrame):
                data = excel_or_df
                customer_short = "REPORT"
                formatted_date = datetime.now().strftime("%B %d, %Y")
                return data, customer_short, formatted_date
            else:
                excel_file = excel_or_df
                creation_time = os.path.getctime(excel_file)
                dt_obj = datetime.fromtimestamp(creation_time)
                formatted_date = dt_obj.strftime("%B %d, %Y")
                customer_short = os.path.splitext(os.path.basename(excel_file))[0]
                data = pd.read_excel(excel_file, sheet_name=sheet_name, header=header_levels)
                return data, customer_short, formatted_date
        except Exception as e:
            self.log_message(LevelType.ERROR, f"resolve_input_data: {str(e)}")
            # Return default values on failure
            return pd.DataFrame(), "UNKNOWN", datetime.now().strftime("%B %d, %Y")

    def trim_drill_down(self, analysis_sku):
        try:
            columns_to_select = [
                ('EPYC Cloud Cost Advisory Recommendations', 'Current Instance', 'Unnamed: 1_level_2'),
                ('EPYC Cloud Cost Advisory Recommendations', 'Current Annual Cost ($)', 'Unnamed: 3_level_2')
            ]
            analysis_sku_item4 = analysis_sku[columns_to_select].copy()
            analysis_sku_item4.columns = ['_'.join([str(i) for i in col if i]).strip() for col in analysis_sku_item4.columns]
            analysis_sku_item4 = analysis_sku_item4.rename(columns={
                'EPYC Cloud Cost Advisory Recommendations_Current Instance_Unnamed: 1_level_2': 'Current Instance',
                'EPYC Cloud Cost Advisory Recommendations_Current Annual Cost ($)_Unnamed: 3_level_2': 'Current Annual Cost ($)'
            })
            analysis_sku_item4['Current Annual Cost ($)'] = pd.to_numeric(
                analysis_sku_item4['Current Annual Cost ($)'], errors='coerce'
            ).fillna(0)
            grouped_sum = analysis_sku_item4.groupby('Current Instance')['Current Annual Cost ($)'].sum().reset_index()
            grouped_sum = grouped_sum.rename(columns={'Current Annual Cost ($)': 'Total Annual Cost'})
            sorted_grouped = grouped_sum.sort_values(by='Total Annual Cost', ascending=False)
            return sorted_grouped
        except Exception as e:
            self.log_message(LevelType.ERROR, f"trim_drill_down: {str(e)}")
            return pd.DataFrame()

    def dollar_spend_eval(self, sorted_grouped, total_cost):
        try:
            sorted_grouped['Total Annual Cost'] = pd.to_numeric(
                sorted_grouped['Total Annual Cost'], errors='raise'
            )
            first_cost = sorted_grouped.iloc[0]['Total Annual Cost']
            max_spend = sorted_grouped.iloc[0]['Current Instance']
            item_4a = round((first_cost/total_cost)*100, 1) if total_cost != 0 else 0
            next_10_sum = sorted_grouped.iloc[1:11]['Total Annual Cost'].sum()
            item_4b = round((next_10_sum/total_cost)*100, 1) if total_cost != 0 else 0
            remaining_sum = sorted_grouped.iloc[11:]['Total Annual Cost'].sum()
            item_4c = round((remaining_sum/total_cost)*100, 1) if total_cost != 0 else 0
            return max_spend, item_4a, item_4b, item_4c
        except Exception as e:
            self.log_message(LevelType.ERROR, f"dollar_spend_eval: {str(e)}")
            return "N/A", 0.0, 0.0, 0.0

    def process_text_shape(self, shape, mapping):
        try:
            if not getattr(shape, "has_text_frame", False):
                return
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original = run.text
                    for key, val in mapping.items():
                        if key in original:
                            new_text = original.replace(key, str(val))
                            run.text = new_text
                            run.font.size = self.LARGE_PT if key in self.LARGE_KEYS else self.DEFAULT_PT
                            run.font.bold = key.startswith("<")
                            original = new_text
        except Exception as e:
            self.log_message(LevelType.ERROR, f"process_text_shape: {str(e)}")

    def protect_pptx(self, file_path: str, password: str):
        try:
            with open(file_path, "rb") as f:
                file_data = BytesIO(f.read())
            file = OOXMLFile(file_data)
            encrypted_data = BytesIO()
            file.encrypt(password, encrypted_data)
            with open(file_path, "wb") as f:
                f.write(encrypted_data.getvalue())
        except Exception as e:
            self.log_message(LevelType.ERROR, f"protect_pptx: {str(e)}")

    def generate_ppt(self, excel_or_df, email, pptx_template, results_path, app_name, portfolio_id):
        try:
            customer = self.extract_customer_from_email(email)
            data, customer_short, formatted = self.resolve_input_data(excel_or_df)
            item_2a = data.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Cloud', 'Unnamed: 5_level_2')]
            analysis_total = data[data[('EPYC Cloud Cost Advisory Recommendations', 'Region', 'Unnamed: 0_level_2')] == "Grand Total"].reset_index(drop=True)
            analysis_sku = data[
                (data[('EPYC Cloud Cost Advisory Recommendations', 'Remark', 'Unnamed: 9_level_2')].isna()) &
                (data[('EPYC Cloud Cost Advisory Recommendations', 'Cloud', 'Unnamed: 5_level_2')] == item_2a)
            ].reset_index(drop=True)
            item_3a = int(data.loc[data[('EPYC Cloud Cost Advisory Recommendations', 'Remark', 'Unnamed: 9_level_2')].isna() & (data[('EPYC Cloud Cost Advisory Recommendations', 'Region', 'Unnamed: 0_level_2')] != "Grand Total"), ('EPYC Cloud Cost Advisory Recommendations', 'Quantity', 'Unnamed: 6_level_2')].sum())
            item_3b = analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Current Annual Cost ($)', 'Unnamed: 3_level_2')]
            item_3bF = self.format_currency(item_3b)
            item_3c = len(analysis_sku[('EPYC Cloud Cost Advisory Recommendations', 'Region', 'Unnamed: 0_level_2')].unique())
            sorted_group = self.trim_drill_down(analysis_sku)
            if not sorted_group.empty:
                max_spend, item_4a, item_4b, item_4c = self.dollar_spend_eval(sorted_group, item_3b)
            else:
                max_spend, item_4a, item_4b, item_4c = "N/A", 0.0, 0.0, 0.0
            item_5a = self.format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Hourly Cost Optimization', 'Annual Cost ($)')])
            item_5b = self.format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Modernize', 'Annual Cost ($)')])
            item_5c = self.format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Modernize & Downsize', 'Annual Cost ($)')])
            item_6a = self.format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Hourly Cost Optimization', 'Annual Savings ($)')])
            item_6b = str(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Hourly Cost Optimization', 'Performance Improvement')]) + "x"
            item_7a = self.format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Modernize', 'Annual Savings ($)')])
            item_7b = str(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Modernize', 'Performance Improvement')]) + "x"
            item_8a = self.format_currency(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Modernize & Downsize', 'Annual Savings ($)')])
            item_8b = str(analysis_total.loc[0, ('EPYC Cloud Cost Advisory Recommendations', 'Modernize & Downsize', 'Performance Improvement')]) + "x"
            ph1 = str(analysis_sku.iloc[0][('EPYC Cloud Cost Advisory Recommendations', 'Current Instance', 'Unnamed: 1_level_2')]) if not analysis_sku.empty else ""
            ph4 = str(analysis_sku.iloc[0][('EPYC Cloud Cost Advisory Recommendations', 'Hourly Cost Optimization', 'Instance')]) if not analysis_sku.empty else ""
            ph3 = str(analysis_sku.iloc[0][('EPYC Cloud Cost Advisory Recommendations', 'Modernize', 'Instance')]) if not analysis_sku.empty else ""
            ph2 = str(analysis_sku.iloc[0][('EPYC Cloud Cost Advisory Recommendations', 'Modernize & Downsize', 'Instance')]) if not analysis_sku.empty else ""
            prs = Presentation(pptx_template)
            slide = prs.slides[0]
            chart_shape = None
            for shape in slide.shapes:
                if hasattr(shape, "has_chart") and shape.has_chart:
                    chart_shape = shape
                    break
            if chart_shape:
                left, top, width, height = chart_shape.left, chart_shape.top, chart_shape.width, chart_shape.height
                sp = chart_shape._element
                sp.getparent().remove(sp)
                chart_data = CategoryChartData()
                chart_data.categories = [max_spend, 'Next 10', 'Rest']
                chart_data.add_series('Dollar Spend', (float(item_4a), float(item_4b), float(item_4c)))
                chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
                chart.has_title = False
                chart.value_axis.has_major_gridlines = False
                chart.category_axis.tick_labels.font.size = Pt(7)
                chart.value_axis.tick_labels.font.size = Pt(7)
                for series in chart.series:
                    series.has_data_labels = True
                    for i, point in enumerate(series.points):
                        dl = point.data_label
                        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                        val = series.values[i]
                        dl.text_frame.text = f"{val:.2f}%"
                        # Force run-level styling to avoid reset
                        for paragraph in dl.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(7)
                                run.font.color.rgb = RGBColor(255, 255, 255)
            tag_map = {
                "customer_name": customer,
                "cloud_provider": item_2a,
                "Date_Format": formatted,
                "<3a>": item_3a,
                "<3b>": item_3bF,
                "<3bF>": item_3bF,
                "<3c>": item_3c,
                "<4a>": item_4a,
                "<4b>": item_4b,
                "<4c>": item_4c,
                "<5a>": item_5a,
                "<5b>": item_5b,
                "<5c>": item_5c,
                "<6a>": item_6a,
                "<6b>": item_6b,
                "<7a>": item_7a,
                "<7b>": item_7b,
                "<8a>": item_8a,
                "<8b>": item_8b,
                "ph1": ph1,
                "ph2": ph2,
                "ph3": ph3,
                "ph4": ph4,
            }
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 6:
                        for sub_shape in shape.shapes:
                            self.process_text_shape(sub_shape, tag_map)
                    else:
                        self.process_text_shape(shape, tag_map)
            output_name = os.path.join(results_path, f"{customer_short}.pptx")
            if os.path.exists(output_name):
                os.remove(output_name)

            prs.save(output_name)  # Reset to start
            return f"{customer_short}.pptx"
        except Exception as err:
            self.log_message(LevelType.ERROR, f"generate_ppt error: {str(err)}")
            return None

    #EIA ppt generation code
    def safe_float(self, value, default=0.0):
        """Safely convert a value to float, returning default if conversion fails."""
        try:
            if value is None:
                return default
            return float(value)
        except (ValueError, TypeError):
            return default

    def format_value(self, value):
        """
        Format numeric values with K (thousands) or M (millions) suffix.
        Examples:
            1234 -> 1.23K
            1234567 -> 1.23M
            123 -> 123
        """
        try:
            num = float(value)
            if abs(num) >= 1_000_000_000:
                return f"{num / 1_000_000_000:.2f}B"
            elif abs(num) >= 1_000_000:
                return f"{num / 1_000_000:.2f}M"
            elif abs(num) >= 1_000:
                return f"{num / 1_000:.2f}k"
            else:
                return f"{num:.2f}"
        except (ValueError, TypeError):
            return str(value)

    def extract_excel_data(self, excel_path):
        """
        Extract all required data from the Excel file.
        
        Args:
            excel_path: Path to the Excel file
            
        Returns:
            dict: Dictionary containing all extracted data
        """
        try:
            self.log_message(LevelType.INFO, f"Reading Excel file: {excel_path}")
            
            wb = load_workbook(excel_path)
            ws = wb['Recommended-Instance']
            
            data = {
                'instances': [],
                'grand_total': {},
                'unique_regions': set()
            }
            
            # Find data rows (start from row 5, skip headers in rows 1-4)
            data_row_start = 5
            
            for row_idx in range(data_row_start, ws.max_row + 1):
                # Check if this is the Grand Total row
                if ws.cell(row_idx, 1).value == "Grand Total":
                    data['grand_total'] = {
                        'current_cost': self.safe_float(ws.cell(row_idx, 3).value),
                        'current_power': self.safe_float(ws.cell(row_idx, 4).value),
                        'current_carbon': self.safe_float(ws.cell(row_idx, 5).value),
                        'optimal_cost': self.safe_float(ws.cell(row_idx, 13).value),
                        'optimal_power': self.safe_float(ws.cell(row_idx, 14).value),
                        'optimal_carbon': self.safe_float(ws.cell(row_idx, 15).value),
                        'optimal_savings': self.safe_float(ws.cell(row_idx, 16).value),
                        'optimal_perf': ws.cell(row_idx, 17).value or 0,
                        'good_cost': self.safe_float(ws.cell(row_idx, 21).value),
                        'good_power': self.safe_float(ws.cell(row_idx, 22).value),
                        'good_carbon': self.safe_float(ws.cell(row_idx, 23).value),
                        'good_savings': self.safe_float(ws.cell(row_idx, 24).value),
                        'good_perf': ws.cell(row_idx, 25).value or 0,
                    }
                    break
                
                # Skip empty rows or note rows
                region = ws.cell(row_idx, 1).value
                if not region or region.startswith("Note"):
                    continue
                    
                instance_data = {
                    'region': region,
                    'current_instance': ws.cell(row_idx, 2).value,
                    'current_cost': self.safe_float(ws.cell(row_idx, 3).value),
                    'current_power': self.safe_float(ws.cell(row_idx, 4).value),
                    'current_carbon': self.safe_float(ws.cell(row_idx, 5).value),
                    'uuid': ws.cell(row_idx, 6).value,
                    'cloud': ws.cell(row_idx, 7).value,
                    'pricing_model': ws.cell(row_idx, 8).value,
                    'vcpu': self.safe_float(ws.cell(row_idx, 9).value),
                    'remark': ws.cell(row_idx, 10).value,
                    'optimal_instance': ws.cell(row_idx, 11).value,
                    'optimal_vcpu': self.safe_float(ws.cell(row_idx, 12).value),
                    'optimal_cost': self.safe_float(ws.cell(row_idx, 13).value),
                    'optimal_power': self.safe_float(ws.cell(row_idx, 14).value),
                    'optimal_carbon': self.safe_float(ws.cell(row_idx, 15).value),
                    'optimal_savings': self.safe_float(ws.cell(row_idx, 16).value),
                    'optimal_perf': ws.cell(row_idx, 17).value or 0,
                    'good_instance': ws.cell(row_idx, 19).value,
                    'good_vcpu': self.safe_float(ws.cell(row_idx, 20).value),
                    'good_cost': self.safe_float(ws.cell(row_idx, 21).value),
                    'good_power': self.safe_float(ws.cell(row_idx, 22).value),
                    'good_carbon': self.safe_float(ws.cell(row_idx, 23).value),
                    'good_savings': self.safe_float(ws.cell(row_idx, 24).value),
                    'good_perf': ws.cell(row_idx, 25).value or 0,
                }
                
                data['instances'].append(instance_data)
                data['unique_regions'].add(region)
            
            self.log_message(LevelType.INFO, f"Extracted {len(data['instances'])} instances from Excel")
            return data
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error extracting Excel data: {e}", ErrorCode=-1)
            raise

    def replace_text_preserving_formatting(self, shape, new_text):
        """
        Replace text in a shape while preserving the formatting of the first run.
        
        Args:
            shape: The PowerPoint shape object
            new_text: The new text to insert
        """
        if not shape.has_text_frame:
            return

        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            return

        # Target the first paragraph
        p = text_frame.paragraphs[0]
        
        if not p.runs:
            p.add_run().text = new_text
            return

        # Set the text of the first run to the new text
        p.runs[0].text = new_text

        # Clear text from all subsequent runs in this paragraph
        for i in range(1, len(p.runs)):
            p.runs[i].text = ""
            
        # Clear all subsequent paragraphs completely
        # We loop through remaining paragraphs and clear their text
        # Note: We cannot easily 'remove' paragraphs in python-pptx without low-level XML manipulation
        # so we just empty their text.
        for i in range(1, len(text_frame.paragraphs)):
            for run in text_frame.paragraphs[i].runs:
                run.text = ""

    def update_cloud_footprint(self, slide, data):
        """
        Update Cloud Footprint section with instance count, monthly spend, and regions.
        
        Args:
            slide: PowerPoint slide object
            data: Extracted Excel data
        """
        try:
            # Count instances with recommendations 
            instances_analyzed = sum(1 for inst in data['instances'] if not inst.get('remark'))
            
            # Monthly spend (directly from current cost)
            monthly_spend = data['grand_total'].get('current_cost', 0)
            # Annualize if needed (User requested Annual Spend in Cloud Footprint)
            annual_spend = monthly_spend * 12
            
            # Unique regions count
            region_count = len(data['unique_regions'])
            
            # Format monthly spend using format_value function
            # Using Annual Spend as per request
            annual_spend_str = f"${self.format_value(annual_spend)}"
                
            # Update text shapes by targeting specific placeholders
            for shape in slide.shapes:
                if not hasattr(shape, 'text'):
                    continue
                    
                text = shape.text.strip()
                
                # Target "02" -> Instances Analyzed count
                if text == "<3a>":
                    self.replace_text_preserving_formatting(shape, str(instances_analyzed).zfill(2))
                    
                # Target "$2.5 K" or "$ 2.5 K" -> Annual Spend
                # The placeholder <3bF> is now for Annual Spend
                elif text == "<3bF>":
                    self.replace_text_preserving_formatting(shape, annual_spend_str)
                    
                # Target Summary text with region count
                elif "Summary : Current infrastructure across" in text:
                    # Use regex to replace [X] with actual count
                    new_text = re.sub(r'\[\d+\]', f'[{region_count}]', text)
                    self.replace_text_preserving_formatting(shape, new_text)
            
            self.log_message(LevelType.INFO, f"Updated Cloud Footprint: {instances_analyzed} instances, {annual_spend_str} annual spend, {region_count} regions")
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error updating Cloud Footprint: {e}", ErrorCode=-1)
            raise

    def style_chart_for_dark_theme(self, chart, bar_color_rgb=(0, 176, 240), is_percentage=False, font_size=6):
        """
        Apply dark theme styling to chart: white fonts, white lines, no gridlines, custom bar colors, no legend, data labels.
        
        Args:
            chart: Chart object to style
            bar_color_rgb: RGB tuple for bar color (default: cyan/turquoise)
            is_percentage: If True, format Y-axis as percentage (0-100), else use K format
            font_size: Font size for labels (default: 6)
        """
        try:
            # Remove chart title if it exists
            chart.has_title = False
            
            # Remove legend
            chart.has_legend = False
            
            # Style value axis (Y-axis)
            if hasattr(chart, 'value_axis'):
                value_axis = chart.value_axis
                value_axis.tick_labels.font.color.rgb = RGBColor(255, 255, 255)
                value_axis.tick_labels.font.size = Pt(font_size)
                # Remove gridlines
                value_axis.has_major_gridlines = False
                value_axis.has_minor_gridlines = False
                value_axis.format.line.color.rgb = RGBColor(255, 255, 255)
                # Format numbers based on chart type
                if is_percentage:
                    # For percentage charts (0-100), use simple number format
                    value_axis.tick_labels.number_format = '0'
                else:
                    # Conditional format: M for millions, K for thousands
                    # [>=1000000]#,##0.00,,"M";[>=1000]#,##0.00,"K";0.00
                    value_axis.tick_labels.number_format = '[>=1000000]#,##0.00,,"M";[>=1000]#,##0.00,"K";0.00'
                
                # Ensure labels are at the bottom (Low)
                value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
                
            # Style category axis (X-axis)
            if hasattr(chart, 'category_axis'):
                category_axis = chart.category_axis
                category_axis.tick_labels.font.color.rgb = RGBColor(255, 255, 255)
                category_axis.tick_labels.font.size = Pt(font_size)
                category_axis.format.line.color.rgb = RGBColor(255, 255, 255)
            
            # Set bar colors and add data labels
            for series in chart.series:
                # Set bar colors
                for point in series.points:
                    fill = point.format.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(*bar_color_rgb)
                
                # Add data labels manually per point for B/M/K support
                try:
                    # Activate at Plot level to ensure visibility
                    plot = chart.plots[0]
                    plot.has_data_labels = True
                    
                    for series in chart.series:
                        series.has_data_labels = True
                        for i, point in enumerate(series.points):
                            data_label = point.data_label
                            data_label.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                            
                            # Get the actual value
                            val = series.values[i]
                            
                            # Manual formatting
                            if is_percentage:
                                label_text = f"{val:.2f}%"
                            else:
                                if val >= 1_000_000_000:
                                    label_text = f"{val/1_000_000_000:.2f}B"
                                elif val >= 1_000_000:
                                    label_text = f"{val/1_000_000:.2f}M"
                                elif val >= 1_000:
                                    label_text = f"{val/1_000:.2f}K"
                                else:
                                    label_text = f"{val:.2f}"
                            
                            # Set the text directly
                            data_label.text_frame.text = label_text
                            
                            # Force run-level styling to avoid reset to black/default
                            for paragraph in data_label.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(7)
                                    run.font.color.rgb = RGBColor(255, 255, 255)
                except Exception as e:
                     self.log_message(LevelType.WARNING, f"Could not manually set data labels: {e}")
                
        except Exception as e:
            self.log_message(LevelType.WARNING, f"Could not fully style chart: {e}")

    def update_business_value(self, slide, data):
        """
        Update Business Value section with dynamic chart showing Current/Optimal/Good costs.
        Removes old static shapes and inserts a new chart.
        
        Args:
            slide: PowerPoint slide object
            data: Extracted Excel data
        """
        try:
            # Annualize costs for the chart
            current_cost = data['grand_total'].get('current_cost', 0) * 12
            optimal_cost = data['grand_total'].get('optimal_cost', 0) * 12
            good_cost = data['grand_total'].get('good_cost', 0) * 12
            
            # Remove old shapes in the Business Value chart area - EXTREMELY AGGRESSIVE
            # Need to remove Y-axis text far on the left side
            # Based on inspection: title at left=3205163, top=609600
            # Expand FAR left to catch all Y-axis labels: left 2800000-6000000, top 600000-2600000
            shapes_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, 'left') and hasattr(shape, 'top'):
                    # Check if shape is in the Business Value chart area (including far left Y-axis)
                    if (2800000 <= shape.left <= 6000000 and 
                        900000 <= shape.top <= 2900000):
                        # PRESERVE TITLES AND BACKGROUND BOXES
                        if hasattr(shape, 'text'):
                            # Keep the title
                            if 'Business Value' in shape.text:
                                continue  
                        
                        # Keep background boxes (AutoShapes without text or empty text)
                        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                            if not shape.has_text_frame or not shape.text.strip():
                                # Check dimensions to distinguish background from chart elements
                                # RELAXED FILTER: Keep if reasonably large (Width > 1,000,000)
                                if shape.width > 1000000:
                                    continue

                        # Remove everything else (Chart content, labels, small shapes)
                        shapes_to_remove.append(shape)
            
            # Remove shapes (must be done separately to avoid iterator issues)
            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)
            
            # Define chart data
            chart_data = CategoryChartData()
            chart_data.categories = ['Current', 'Optimal', 'Good']
            chart_data.add_series('Cost ($)', (current_cost, optimal_cost, good_cost))
            
            # Insert new chart - moved down more
            # Insert new chart - moved down MORE
            # Y position increased to 1350000
            x, y, cx, cy = 3250000, 1350000, 2700000, 1350000
            graphic_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            chart = graphic_frame.chart
            
            # Apply dark theme styling with cyan color and no legend
            self.style_chart_for_dark_theme(chart, bar_color_rgb=(0, 176, 240))
            
            self.log_message(LevelType.INFO, f"Inserted Business Value chart: Current=${current_cost}, Optimal=${optimal_cost}, Good=${good_cost}")
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error updating Business Value: {e}", ErrorCode=-1)
            raise

    def update_path_to_savings(self, slide, data):
        """
        Update Path to Savings section with instance types and savings from first row + Grand Total.
        
        Args:
            slide: PowerPoint slide object
            data: Extracted Excel data
        """
        try:
            # Get first row data for instance types
            first_instance = data['instances'][0] if data['instances'] else {}
            
            current_instance = first_instance.get('current_instance', '')
            optimal_instance = first_instance.get('optimal_instance', '')
            good_instance = first_instance.get('good_instance', '')
            
            # Get savings and perf from Grand Total (Monthly values)
            optimal_savings_monthly = data['grand_total'].get('optimal_savings', 0)
            optimal_perf = data['grand_total'].get('optimal_perf', 0)
            good_savings_monthly = data['grand_total'].get('good_savings', 0)
            good_perf = data['grand_total'].get('good_perf', 0)
            
            # Annualize Savings
            optimal_savings_annual = optimal_savings_monthly * 12
            good_savings_annual = good_savings_monthly * 12
            
            self.log_message(LevelType.INFO, f"Path to Savings - Optimal Perf: {optimal_perf}, Good Perf: {good_perf}")
            
            # Format values with K/M suffixes
            optimal_savings_str = f"$ {self.format_value(optimal_savings_annual)}"
            good_savings_str = f"$ {self.format_value(good_savings_annual)}"
            
            optimal_perf_str = f"{self.format_value(optimal_perf)} X"
            good_perf_str = f"{self.format_value(good_perf)} X"
            
            # Update text shapes by targeting specific placeholders
            for shape in slide.shapes:
                if not hasattr(shape, 'text'):
                    continue
                    
                text = shape.text.strip()
                
                # OPTIMAL SECTION PLACEHOLDERS
                if text == "ph1":  # Current Instance placeholder (Optimal section)
                    # Note: This placeholder appears twice (Optimal and Good sections).
                    # We need to distinguish based on position (Top vs Bottom)
                    # if shape.top < 5000000:  # Top half (Optimal)
                    #     self.replace_text_preserving_formatting(shape, current_instance)
                    # else:
                    self.replace_text_preserving_formatting(shape, current_instance)
                    
                elif text == "ph2":  # Optimal Instance placeholder
                    self.replace_text_preserving_formatting(shape, optimal_instance)
                    
                elif text == "<3os>":  # Optimal Savings placeholder
                    self.replace_text_preserving_formatting(shape, optimal_savings_str)
                    
                elif text == "<3op>":  # Optimal Perf placeholder
                    self.replace_text_preserving_formatting(shape, optimal_perf_str)
                    # Reduce font size to fit
                    if shape.has_text_frame and shape.text_frame.paragraphs:
                        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(5.3)
                    
                # GOOD SECTION PLACEHOLDERS
                # elif text == "ph1":  # Current Instance placeholder (Good section)
                #     self.replace_text_preserving_formatting(shape, current_instance)
                    
                elif text == "ph3":  # Good Instance placeholder
                    self.replace_text_preserving_formatting(shape, good_instance)
                    
                elif text == "<3gs>":  # Good Savings placeholder
                    self.replace_text_preserving_formatting(shape, good_savings_str)
                    
                elif text == "<3gp>":  # Good Perf placeholder
                    self.replace_text_preserving_formatting(shape, good_perf_str)
                    # Reduce font size to fit
                    if shape.has_text_frame and shape.text_frame.paragraphs:
                        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(5.3)
            
            self.log_message(LevelType.INFO, "Updated Path to Savings section")
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error updating Path to Savings: {e}", ErrorCode=-1)
            raise

    def update_dollar_spend_distribution(self, slide, data):
        """
        Update Dollar Spend Distribution section with: 
        1st highest cost instance, Next 10 (sum), Rest (sum).
        Always shows 3 bars even if values are 0.
        Removes old static shapes and inserts a new chart.
        
        Args:
            slide: PowerPoint slide object
            data: Extracted Excel data
        """
        try:
            # Group instances by type and sum their costs
            instance_costs = {}
            for inst in data['instances']:
                name = inst.get('current_instance', 'Unknown')
                cost = inst.get('current_cost', 0)
                instance_costs[name] = instance_costs.get(name, 0) + cost
            
            # Create list of dicts for sorting
            grouped_instances = [{'current_instance': k, 'current_cost': v} for k, v in instance_costs.items()]
            
            # Sort instances by current cost descending
            sorted_instances = sorted(grouped_instances, key=lambda x: x['current_cost'], reverse=True)
            
            total_cost = data['grand_total'].get('current_cost', 0)
            
            # ALWAYS create 3 categories: Top 1, Next 10, Rest
            # First bar: Top 1 (largest instance group)
            if sorted_instances:
                top_instance = sorted_instances[0]
                top_percentage = round((top_instance['current_cost'] / total_cost * 100), 2) if total_cost > 0 else 0
                top_name = top_instance['current_instance']
            else:
                top_percentage = 0
                top_name = "N/A"
            
            # Second bar: Next 10 (groups 2-11)
            if len(sorted_instances) > 1:
                next_10_instances = sorted_instances[1:11]
                next_10_cost = sum(inst['current_cost'] for inst in next_10_instances)
                next_10_percentage = round((next_10_cost / total_cost * 100), 2) if total_cost > 0 else 0
            else:
                next_10_percentage = 0
            
            # Third bar: Rest (groups 12+)
            if len(sorted_instances) > 11:
                rest_instances = sorted_instances[11:]
                rest_cost = sum(inst['current_cost'] for inst in rest_instances)
                rest_percentage = round((rest_cost / total_cost * 100), 2) if total_cost > 0 else 0
            else:
                rest_percentage = 0
            
            # Remove old shapes in the Dollar Spend chart area - EXTREMELY AGGRESSIVE
            # Expand area to catch all Y-axis labels on the FAR left (100, 75, 50, 25, 0)
            # Based on inspection: title at left=185738, top=2805113
            # Chart area: left -100000 to 3200000, top 2800000-4900000
            shapes_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, 'left') and hasattr(shape, 'top'):
                    # Catch shapes from far left (including Y-axis labels) to right
                    if (shape.left <= 3200000 and 
                        3100000 <= shape.top <= 4900000):
                        # PRESERVE TITLES AND BACKGROUND BOXES
                        if hasattr(shape, 'text'):
                            # Keep the title
                            if 'Dollar Spend' in shape.text:
                                continue  
                        
                        # Keep background boxes (AutoShapes without text or empty text)
                        # FILTER BY SIZE: Keep only large backgrounds
                        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                            if not shape.has_text_frame or not shape.text.strip():
                                if shape.width > 1000000 and shape.height > 1000000:
                                    continue
                                
                                # Preserve header background
                                if shape.top < 2900000: # Section starts at 2800000
                                    continue

                        # Remove everything else
                        shapes_to_remove.append(shape)
            
            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)
            
            # Define chart data - ALWAYS 3 categories
            chart_data = CategoryChartData()
            chart_data.categories = [top_name, 'Next 10', 'Rest']
            chart_data.add_series('Percentage', (top_percentage, next_10_percentage, rest_percentage))
            
            # Insert new chart - MOVED DOWN MORE within the section
            # Position adjusted: y increased to move down, staying within bounds
            x, y, cx, cy = 250000, 3500000, 2700000, 1350000
            graphic_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            chart = graphic_frame.chart
            
            # Apply dark theme styling with cyan color - PERCENTAGE FORMAT for Y-axis
            self.style_chart_for_dark_theme(chart, bar_color_rgb=(0, 176, 240), is_percentage=True)
            
            # Set Y-axis to percentage format (0-100)
            if hasattr(chart, 'value_axis'):
                chart.value_axis.maximum_scale = 100
                chart.value_axis.minimum_scale = 0
            
            self.log_message(LevelType.INFO, "Inserted Dollar Spend Distribution chart with 3 bars")
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error updating Dollar Spend Distribution: {e}", ErrorCode=-1)
            raise


    def update_power_carbon(self, slide, data):
        """
        Update Power/Carbon section with dynamic charts for Current/Optimal/Good.
        Also updates the detailed text blocks with calculated savings.
        
        Args:
            slide: PowerPoint slide object
            data: Extracted Excel data
        """
        try:
            current_power = data['grand_total'].get('current_power', 0)
            optimal_power = data['grand_total'].get('optimal_power', 0)
            good_power = data['grand_total'].get('good_power', 0)
            
            current_carbon = data['grand_total'].get('current_carbon', 0)
            optimal_carbon = data['grand_total'].get('optimal_carbon', 0)
            good_carbon = data['grand_total'].get('good_carbon', 0)
            
            current_cost = data['grand_total'].get('current_cost', 0)
            optimal_cost = data['grand_total'].get('optimal_cost', 0)
            good_cost = data['grand_total'].get('good_cost', 0)
            
            # Calculate savings percentages (clamp negative values to 0)
            opt_cost_save_pct = round(((current_cost - optimal_cost) / current_cost * 100), 2) if current_cost > 0 else 0
            opt_power_save_pct = round(((current_power - optimal_power) / current_power * 100), 2) if current_power > 0 else 0
            opt_carbon_save_pct = round(((current_carbon - optimal_carbon) / current_carbon * 100), 2) if current_carbon > 0 else 0
            
            good_cost_save_pct = round(((current_cost - good_cost) / current_cost * 100), 2) if current_cost > 0 else 0
            good_power_save_pct = round(((current_power - good_power) / current_power * 100), 2) if current_power > 0 else 0
            good_carbon_save_pct = round(((current_carbon - good_carbon) / current_carbon * 100), 2) if current_carbon > 0 else 0
            
            # Remove old shapes in the Power/Carbon chart area
            # Based on inspection: title at left=3205163, top=2819400
            # Chart area extended to bottom to catch all extra text
            # Chart area: left 3000000-9000000 (Expanded left and right), top 3000000-5100000
            shapes_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, 'left') and hasattr(shape, 'top'):
                    if (3000000 <= shape.left <= 9000000 and 
                        3100000 <= shape.top <= 5100000):
                        # PRESERVE TITLES AND BACKGROUND BOXES
                        if hasattr(shape, 'text'):
                            # Keep the section title - check robustly
                            if 'Power' in shape.text and 'Carbon' in shape.text:
                                continue
                            # Keep the detail text boxes (contain Cost: and lower cost)
                            if 'Cost:' in shape.text and 'lower cost' in shape.text:
                                continue
                            # Keep Optimal and Good boxes (headers) - ONLY in the detail area (Right side)
                            if ('Optimal' in shape.text or 'Good' in shape.text) and shape.left > 7000000:
                                continue
                        
                        # Keep background boxes (AutoShapes without text or empty text)
                        # FILTER BY SIZE: Keep only large backgrounds
                        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                            if not shape.has_text_frame or not shape.text.strip():
                                if shape.width > 1000000 and shape.height > 1000000:
                                    continue
                                
                                # Preserve header background
                                if shape.top < 3100000: # Section starts at 3000000
                                    continue
                                
                                # Preserve Background Boxes for Optimal/Good Details (Right side)
                                # These might be smaller than the 1M x 1M filter above
                                # Detail area is roughly Left > 7.0M
                                if shape.left > 7000000:
                                    continue

                        # Remove everything else (including old chart elements)
                        shapes_to_remove.append(shape)
            
            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)
            
            # --- CHART 1: POWER (Left) ---
            power_chart_data = CategoryChartData()
            # For Bar Chart (Horizontal), categories are bottom-to-top.
            # Template visual: Top=Current, Middle=Optimal, Bottom=Good.
            power_chart_data.categories = ['Current', 'Optimal', 'Good']
            power_chart_data.add_series('Power (kW)', (current_power, optimal_power, good_power))
            
            # Position: Left side of the container
            # Increased height slightly for better visibility
            x1, y1, cx1, cy1 = 3200000, 3300000, 2000000, 1150000
            graphic_frame1 = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, x1, y1, cx1, cy1, power_chart_data
            )
            chart1 = graphic_frame1.chart
            self.style_chart_for_dark_theme(chart1, bar_color_rgb=(255, 153, 51), font_size=5) # Orange

            # Add Manual Title for Power
            # Moved up again to avoid overlap with boxes
            tx_box_1 = slide.shapes.add_textbox(3200000, 4320000, 2200000, 250000)
            tf_1 = tx_box_1.text_frame
            tf_1.text = "Power (kW)"
            tf_1.paragraphs[0].font.size = Pt(8)
            tf_1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf_1.paragraphs[0].alignment = 2 # Centered

            # Move up to Y = 4.52M
            try:
                score_box_y = 4520000
                score_box_cx = 1950000
                score_box_cy = 180000
                shape1 = slide.shapes.add_shape(
                    1, 3250000, score_box_y, score_box_cx, score_box_cy
                )
                shape1.fill.solid()
                shape1.fill.fore_color.rgb = RGBColor(25, 25, 25) # Slightly darker for "transparent dark" look
                shape1.line.color.rgb = RGBColor(89, 89, 89) # Subtler border
                shape1.line.width = Pt(0.5)
                
                tf_score1 = shape1.text_frame
                tf_score1.word_wrap = False
                tf_score1.margin_top = 0
                tf_score1.margin_bottom = 0
                tf_score1.margin_left = 0
                tf_score1.margin_right = 0
                p1 = tf_score1.paragraphs[0]
                p1.alignment = 1 # Center
                
                # Add info icon (ⓘ)
                run_i = p1.add_run()
                run_i.text = " ⓘ "
                run_i.font.size = Pt(6)
                run_i.font.color.rgb = RGBColor(180, 180, 180)
                
                # Add score text (Multi-run for conditional coloring)
                run_lbl = p1.add_run()
                # 2 Decimal Places Max as requested
                run_lbl.text = f"Sustainability Score : Optimal {opt_power_save_pct:.2f}% "
                run_lbl.font.size = Pt(5)
                run_lbl.font.bold = True
                run_lbl.font.color.rgb = RGBColor(255, 255, 255)

                # Optimal Arrow
                run_a1 = p1.add_run()
                run_a1.text = "↑" if opt_power_save_pct >= 0 else "↓"
                run_a1.font.size = Pt(5)
                run_a1.font.bold = True
                run_a1.font.color.rgb = RGBColor(0, 176, 80) if opt_power_save_pct >= 0 else RGBColor(255, 0, 0)

                run_lbl2 = p1.add_run()
                run_lbl2.text = f" , Good {good_power_save_pct:.2f}% "
                run_lbl2.font.size = Pt(5)
                run_lbl2.font.bold = True
                run_lbl2.font.color.rgb = RGBColor(255, 255, 255)

                # Good Arrow
                run_a2 = p1.add_run()
                run_a2.text = "↑" if good_power_save_pct >= 0 else "↓"
                run_a2.font.size = Pt(5)
                run_a2.font.bold = True
                run_a2.font.color.rgb = RGBColor(0, 176, 80) if good_power_save_pct >= 0 else RGBColor(255, 0, 0)
            except Exception as e:
                self.log_message(LevelType.WARNING, f"Could not add Power sustainability score box: {e}")

            # Manually color points for Power Chart
            try:
                series = chart1.series[0]
                points = series.points
                points[0].format.fill.solid()
                points[0].format.fill.fore_color.rgb = RGBColor(192, 0, 0) # Current -> Dark Red
                points[1].format.fill.solid()
                points[1].format.fill.fore_color.rgb = RGBColor(255, 192, 0) # Optimal -> Orange
                points[2].format.fill.solid()
                points[2].format.fill.fore_color.rgb = RGBColor(255, 192, 0) # Good -> Orange
                
                if hasattr(chart1, 'category_axis'):
                    chart1.category_axis.reverse_order = True
                if hasattr(chart1, 'value_axis'):
                    chart1.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
            except Exception as e:
                self.log_message(LevelType.WARNING, f"Could not manually color Power chart points: {e}")

            # --- CHART 2: CARBON (Right) ---
            carbon_chart_data = CategoryChartData()
            carbon_chart_data.categories = ['Current', 'Optimal', 'Good']
            carbon_chart_data.add_series('Carbon (kgCO2eq)', (current_carbon, optimal_carbon, good_carbon))
            
            # --- VERTICAL DIVIDER LINE ---
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, 5250000, 3150000, 5250000, 4750000
            )
            connector.line.color.rgb = RGBColor(128, 128, 128) # Grey
            connector.line.width = Pt(1)

            # Position: Right side of the container
            # Increased height slightly
            x2, y2, cx2, cy2 = 5300000, 3300000, 2000000, 1150000
            graphic_frame2 = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, x2, y2, cx2, cy2, carbon_chart_data
            )
            chart2 = graphic_frame2.chart
            self.style_chart_for_dark_theme(chart2, bar_color_rgb=(146, 208, 80), font_size=5) # Green

            # Add Manual Title for Carbon
            # Moved up again
            tx_box_2 = slide.shapes.add_textbox(5300000, 4320000, 2200000, 250000)
            tf_2 = tx_box_2.text_frame
            tf_2.text = "Carbon (kgCO2eq)"
            tf_2.paragraphs[0].font.size = Pt(8)
            tf_2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf_2.paragraphs[0].alignment = 2 # Centered

            # Sustainability Score Box for Carbon
            try:
                shape2 = slide.shapes.add_shape(
                    1, 5325000, score_box_y, score_box_cx, score_box_cy
                )
                shape2.fill.solid()
                shape2.fill.fore_color.rgb = RGBColor(25, 25, 25)
                shape2.line.color.rgb = RGBColor(89, 89, 89)
                shape2.line.width = Pt(0.5)
                
                tf_score2 = shape2.text_frame
                tf_score2.word_wrap = False
                tf_score2.margin_top = 0
                tf_score2.margin_bottom = 0
                tf_score2.margin_left = 0
                tf_score2.margin_right = 0
                p2 = tf_score2.paragraphs[0]
                p2.alignment = 1
                
                # Add info icon (ⓘ)
                run_i2 = p2.add_run()
                run_i2.text = " ⓘ "
                run_i2.font.size = Pt(6)
                run_i2.font.color.rgb = RGBColor(180, 180, 180)
                
                # Add score text (Multi-run for conditional coloring)
                run_lbl_c = p2.add_run()
                run_lbl_c.text = f"Sustainability Score : Optimal {opt_carbon_save_pct:.2f}% "
                run_lbl_c.font.size = Pt(5)
                run_lbl_c.font.bold = True
                run_lbl_c.font.color.rgb = RGBColor(255, 255, 255)

                # Optimal Arrow
                run_ca1 = p2.add_run()
                run_ca1.text = "↑" if opt_carbon_save_pct >= 0 else "↓"
                run_ca1.font.size = Pt(5)
                run_ca1.font.bold = True
                run_ca1.font.color.rgb = RGBColor(0, 176, 80) if opt_carbon_save_pct >= 0 else RGBColor(255, 0, 0)

                run_lbl_c2 = p2.add_run()
                run_lbl_c2.text = f" , Good {good_carbon_save_pct:.2f}% "
                run_lbl_c2.font.size = Pt(5)
                run_lbl_c2.font.bold = True
                run_lbl_c2.font.color.rgb = RGBColor(255, 255, 255)

                # Good Arrow
                run_ca2 = p2.add_run()
                run_ca2.text = "↑" if good_carbon_save_pct >= 0 else "↓"
                run_ca2.font.size = Pt(5)
                run_ca2.font.bold = True
                run_ca2.font.color.rgb = RGBColor(0, 176, 80) if good_carbon_save_pct >= 0 else RGBColor(255, 0, 0)
            except Exception as e:
                self.log_message(LevelType.WARNING, f"Could not add Carbon sustainability score box: {e}")

            # --- Calculation Formula at bottom ---
            # Moved up to Y = 4.75M
            try:
                formula_box = slide.shapes.add_textbox(3250000, 4750000, 4000000, 200000)
                tf_f = formula_box.text_frame
                tf_f.paragraphs[0].alignment = 1 # Centered
                run_f = tf_f.paragraphs[0].add_run()
                run_f.text = "Calculation : sustainability Improvement (%) = (Current – Target) / Current x 100"
                run_f.font.size = Pt(5.5)
                run_f.font.bold = True
                run_f.font.color.rgb = RGBColor(0, 176, 240) # Cyan
            except Exception as e:
                self.log_message(LevelType.WARNING, f"Could not add formula text: {e}")

            # Manually color points for Carbon Chart
            try:
                series = chart2.series[0]
                points = series.points
                points[0].format.fill.solid()
                points[0].format.fill.fore_color.rgb = RGBColor(192, 0, 0)
                points[1].format.fill.solid()
                points[1].format.fill.fore_color.rgb = RGBColor(146, 208, 80)
                points[2].format.fill.solid()
                points[2].format.fill.fore_color.rgb = RGBColor(146, 208, 80)
                
                if hasattr(chart2, 'category_axis'):
                    chart2.category_axis.reverse_order = True
                if hasattr(chart2, 'value_axis'):
                    chart2.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
            except Exception as e:
                self.log_message(LevelType.WARNING, f"Could not manually color Carbon chart points: {e}")
                
            self.log_message(LevelType.INFO, "Inserted separate Power and Carbon charts with Sustainability Scores")
            

            # 2. Update detailed text blocks
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.has_text_frame:
                    text = shape.text
                    
                    # OPTIMAL DETAILS
                    if "Cost:" in text and "lower cost" in text and shape.top < 4000000:
                        try:
                            tf = shape.text_frame
                            tf.text = ""  # Clear existing
                            tf.margin_top = 0  # Absolute top
                            tf.margin_left = 0  # Absolute left
                            tf.margin_right = 0
                            tf.margin_bottom = 0
                            tf.vertical_anchor = MSO_ANCHOR.TOP
                            
                            # Details (5.1pt Normal, Indented 4 spaces)
                            def add_detail(paragraph, label, value_str, space_after=Pt(3.5)):
                                paragraph.space_after = space_after
                                paragraph.space_before = 0
                                paragraph.line_spacing = 1.0
                                run = paragraph.add_run()
                                run.text = f"    {label} {value_str}"
                                run.font.size = Pt(5.1)
                                run.font.bold = False
                                run.font.color.rgb = RGBColor(255, 255, 255)
                            
                            # Use existing first paragraph for the first line
                            p0 = tf.paragraphs[0]
                            add_detail(p0, "Cost:", f"${self.format_value(optimal_cost * 12)} ({opt_cost_save_pct:.2f}% lower cost)")
                            
                            p1 = tf.add_paragraph()
                            add_detail(p1, "Power:", f"{self.format_value(optimal_power)} kW ({opt_power_save_pct:.2f}% lower power)")
                            
                            # Carbon split into two lines
                            p2 = tf.add_paragraph()
                            add_detail(p2, "Carbon:", f"{self.format_value(optimal_carbon)} kgCO₂eq ({opt_carbon_save_pct:.2f}% lower", space_after=0)
                            
                            p3 = tf.add_paragraph()
                            add_detail(p3, "", "carbon emission)", space_after=Pt(3.5))

                        except Exception as e:
                            self.log_message(LevelType.ERROR, f"Error rebuilding Optimal details text: {e}", ErrorCode=-1)
                            
                    # GOOD DETAILS
                    elif "Cost:" in text and "lower cost" in text and shape.top > 4000000:
                        try:
                            tf = shape.text_frame
                            tf.text = ""  # Clear existing
                            tf.margin_top = 0  # Absolute top
                            tf.margin_left = 0  # Absolute left
                            tf.margin_right = 0
                            tf.margin_bottom = 0
                            tf.vertical_anchor = MSO_ANCHOR.TOP
                            
                            # Details (5.1pt Normal, Indented 4 spaces)
                            def add_detail_good(paragraph, label, value_str, space_after=Pt(3.5)):
                                paragraph.space_after = space_after
                                paragraph.space_before = 0
                                paragraph.line_spacing = 1.0
                                run = paragraph.add_run()
                                run.text = f"    {label} {value_str}"
                                run.font.size = Pt(5.1)
                                run.font.bold = False
                                run.font.color.rgb = RGBColor(255, 255, 255)
                            
                            # Use existing first paragraph for the first line
                            p0 = tf.paragraphs[0]
                            add_detail_good(p0, "Cost:", f"${self.format_value(good_cost * 12)} ({good_cost_save_pct:.2f}% lower cost)")
                            
                            p1 = tf.add_paragraph()
                            add_detail_good(p1, "Power:", f"{self.format_value(good_power)} kW ({good_power_save_pct:.2f}% lower power)")
                            
                            # Carbon split into two lines
                            p2 = tf.add_paragraph()
                            add_detail_good(p2, "Carbon:", f"{self.format_value(good_carbon)} kgCO₂eq ({good_carbon_save_pct:.2f}% lower", space_after=0)
                            
                            p3 = tf.add_paragraph()
                            add_detail_good(p3, "", "carbon emission)", space_after=Pt(3.5))

                        except Exception as e:
                            self.log_message(LevelType.ERROR, f"Error rebuilding Good details text: {e}", ErrorCode=-1)
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error updating Power/Carbon: {e}", ErrorCode=-1)
            raise

    def update_header_info(self, slide, data, user_email):
        """
        Update header information: Cloud, Customer, and Date.
        
        Args:
            slide: PowerPoint slide object
            data: Extracted Excel data
            user_email: User's email address to extract customer name
        """
        try:
            # Get Cloud provider from first instance
            first_instance = data['instances'][0] if data['instances'] else {}
            cloud_provider = first_instance.get('cloud', 'Unknown')
            
            # Extract customer name from email (domain part, before .com)
            # e.g., pasaikum@amd.com -> amd
            if user_email and '@' in user_email:
                domain = user_email.split('@')[1]
                customer_name = domain.split('.')[0].upper()
            else:
                customer_name = "CLIENT"
            # Calculate annualized values for Summary
            current_cost = data['grand_total'].get('current_cost', 0)
            optimal_cost = data['grand_total'].get('optimal_cost', 0)
            good_cost = data['grand_total'].get('good_cost', 0)
            
            current_annual = current_cost * 12
            optimal_annual = optimal_cost * 12
            good_annual = good_cost * 12
            
            # Calculate savings percentages
            opt_cost_save_pct = round(((current_cost - optimal_cost) / current_cost * 100), 2) if current_cost > 0 else 0
            good_cost_save_pct = round(((current_cost - good_cost) / current_cost * 100), 2) if current_cost > 0 else 0

            # Power/Carbon savings
            current_power = data['grand_total'].get('current_power', 0)
            optimal_power = data['grand_total'].get('optimal_power', 0)
            good_power = data['grand_total'].get('good_power', 0)
            
            current_carbon = data['grand_total'].get('current_carbon', 0)
            optimal_carbon = data['grand_total'].get('optimal_carbon', 0)
            good_carbon = data['grand_total'].get('good_carbon', 0)
            
            opt_power_save_pct = round(((current_power - optimal_power) / current_power * 100), 2) if current_power > 0 else 0
            # good_power_save_pct = max(0, ((current_power - good_power) / current_power * 100)) if current_power > 0 else 0 # Not used in summary text
            
            opt_carbon_save_pct = round(((current_carbon - optimal_carbon) / current_carbon * 100), 2) if current_carbon > 0 else 0
            
            # Get today's date
            today_date = datetime.today().strftime("%B %d, %Y")
            
            # Update text shapes
            for shape in slide.shapes:
                if not hasattr(shape, 'text'):
                    continue
                    
                text = shape.text.strip()
                
                # Update Cloud
                if "Cloud:" in text:
                    # Replace "Cloud: AWS" or similar with "Cloud: [Provider]"
                    # We use regex to be safe, or just replace the known template string
                    if "Cloud: AWS" in text:
                        self.replace_text_preserving_formatting(shape, f"Cloud: {cloud_provider}")
                    else:
                        # Fallback: try to replace just the value if possible, or the whole string
                        # For now, let's assume the template has "Cloud: AWS" as seen in screenshot
                        pass
                
                # Update Customer
                elif "Customer:" in text:
                    if "Customer: INFOBELLIT" in text:
                        self.replace_text_preserving_formatting(shape, f"Customer: {customer_name}")
                
                # Update Date
                elif "Date:" in text:
                    # The template has "Date: October 23, 2025"
                    # We'll try to match "Date:" and replace the whole thing
                    if "Date:" in text:
                        self.replace_text_preserving_formatting(shape, f"Date: {today_date}")
                
                # Update Executive Summary Text
                elif "Modernizing" in text and "infrastructure can reduce annual spend" in text:
                    # Dynamic Summary Text
                    summary_text = (
                        f"Modernizing AMD EPYC infrastructure can reduce annual spend from ${self.format_value(current_annual)} (Current) "
                        f"to ${self.format_value(optimal_annual)} (Optimal) with {int(opt_cost_save_pct)}% savings per year "
                        f"for ${self.format_value(good_annual)} (Good) {int(good_cost_save_pct)}% savings per year, "
                        f"delivering significant cost efficiency and sustainability benefits. All this while reducing "
                        f"power consumption by {int(opt_power_save_pct)}% and reducing carbon emissions by {int(opt_carbon_save_pct)}%"
                    )
                    self.replace_text_preserving_formatting(shape, summary_text)
                        
            self.log_message(LevelType.INFO, f"Updated Header: Cloud={cloud_provider}, Customer={customer_name}, Date={today_date}")

        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error updating Header info: {e}", ErrorCode=-1)
            # Don't raise, just log error as this is non-critical


    def generate_ppt_from_excel(self, excel_path, template_path, output_path, user_email):
        """
        Main function to generate PowerPoint presentation from Excel data.
        
        Args:
            excel_path: Path to the Excel file
            template_path: Path to the PowerPoint template
            output_path: Path where the output PPT will be saved
            user_email: User's email address for customer name extraction
        """
        try:
            self.log_message(LevelType.INFO, f"Starting PPT generation from {excel_path}")

            # Extract data from Excel
            data = self.extract_excel_data(excel_path)
            
            # Load PowerPoint template
            self.log_message(LevelType.INFO, f"Loading PPT template: {template_path}")
            prs = Presentation(template_path)
            
            # Get the first slide (assuming single slide template)
            slide = prs.slides[0]
            
            # Update all sections
            self.update_header_info(slide, data, user_email)
            self.update_cloud_footprint(slide, data)
            self.update_business_value(slide, data)
            self.update_path_to_savings(slide, data)
            self.update_dollar_spend_distribution(slide, data)
            self.update_power_carbon(slide, data)
            
            # Save the presentation
            prs.save(output_path)
            self.log_message(LevelType.INFO, f"PowerPoint successfully generated at: {output_path}")
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Error generating PowerPoint: {e}", ErrorCode=-1)
            raise

    #email sending code
    def normalize_recipients(self, to_email):
        if not to_email:
            return []
        if isinstance(to_email, str):
            return [addr.strip() for addr in to_email.split(",") if addr.strip()]
        return list({str(addr).strip() for addr in to_email if str(addr).strip()})

    def get_access_token(self):
        try:
            app = ConfidentialClientApplication(
                client_id=CLIENT_ID,
                authority=AUTHORITY,
                client_credential=CLIENT_SECRET
            )
            token_response = app.acquire_token_for_client(scopes=SCOPE)
            if "access_token" not in token_response:
                msg = f"Failed to get token: {token_response.get('error_description')}"
                self.log_message(LevelType.ERROR, msg)
                return None
            return token_response['access_token']
        except Exception as e:
            self.log_message(LevelType.ERROR, f"get_access_token error: {str(e)}")
            return None

    def post_send_mail(self, payload, token):
        try:
            headers = {
                "Authorization": f"Bearer {token}",
                "Content-Type": "application/json"
            }
            response = requests.post(GRAPH_ENDPOINT, headers=headers, data=json.dumps(payload))
            return response
        except Exception as e:
            self.log_message(LevelType.ERROR, f"post_send_mail error: {str(e)}")
            return None

    def send_recommendations_completed_email(self,
                                            to_email,
                                            app_name,
                                            portfolio_id,
                                            portfolio_name=None,
                                            excel_path=None,
                                            pptx_path=None,
                                            invalid_file=None,
                                            excel_filename="report.xlsx",
                                            pptx_filename="slides.pptx",
                                            zip_filename="recommendations_bundle.zip",
                                            invalid_file_name="invalid_records.xlsx"
                                            ):
        try:
            recipients = self.normalize_recipients(to_email)
            if not recipients:
                return "No user emails specified to send completion email", False

            if not excel_path or not os.path.exists(excel_path):
                return f"Excel file not found at {excel_path}", False

            ppt_available = bool(pptx_path and os.path.exists(pptx_path))
            invalid_available = bool(invalid_file and os.path.exists(invalid_file))
            if invalid_available:
                with open(invalid_file, "rb") as f:
                    file_bytes = f.read()

                invalid_record_s3_key = self.upload_file_to_s3(
                    file_bytes=file_bytes,
                    app_name=app_name,
                    user_email=to_email,
                    file_name=invalid_file_name,
                    sub_folder="input_remarks"
                )
                if invalid_record_s3_key:
                    self.portfolio_collection.update_one(
                        {"_id": ObjectId(portfolio_id)},
                        {
                            "$set": {
                                "invalid_record_s3_key": invalid_record_s3_key
                            }
                        }
                    )
            rid_fragment = f" {portfolio_name}" if portfolio_name else ""

            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
                zf.write(excel_path, excel_filename)
                if ppt_available:
                    zf.write(pptx_path, pptx_filename)
                if invalid_available:
                    zf.write(invalid_file, invalid_file_name)

            zip_bytes = zip_buffer.getvalue()

            included_parts = ["the Excel report"]
            if ppt_available:
                included_parts.append("the PowerPoint presentation")

            if app_name.upper() == "CCA":
                application_name = self.CCA_APP_NAME
            else:
                application_name = self.EIA_APP_NAME

            subject = f"Your {application_name} Recommendations Are Ready for Portfolio: {rid_fragment}"
            

            template = self.EMAIL_TEMPLATES.get(app_name.upper())
            body_html = template.format(portfolio_name=rid_fragment, SENDER=SENDER)

            attachment = {
                "@odata.type": "#microsoft.graph.fileAttachment",
                "name": zip_filename,
                "contentType": "application/zip",
                "contentBytes": base64.b64encode(zip_bytes).decode("ascii"),
            }

            message_payload = {
                "message": {
                    "subject": subject,
                    "body": {"contentType": "HTML", "content": body_html},
                    "toRecipients": [{"emailAddress": {"address": addr}} for addr in recipients],
                    "attachments": [attachment],
                }
            }

            token = self.get_access_token()
            if not token:
                return "Failed to obtain access token for sending email", False
            response = self.post_send_mail(message_payload, token)
            if response and response.status_code == 202:
                return f"Recommendations email sent to {', '.join(recipients)}", True
            return f"Failed to send recommendations email: {response.text if response else 'No response'}", False

        except Exception as e:
            self.log_message(LevelType.ERROR, f"Exception during recommendations email send: {e}")
            return f"Exception during recommendations email send: {e}", False

    def safe_unlink(self, path):
        try:
            if path and os.path.exists(path) and os.path.isfile(path):
                os.remove(path)
                return True
        except Exception as e:
            self.log_message(LevelType.WARNING, f"safe_unlink error: {e}")
        return False

    def safe_rmtree(self, path):
        try:
            if path and os.path.exists(path) and os.path.isdir(path):
                shutil.rmtree(path, ignore_errors=True)
                return True
        except Exception as e:
            self.log_message(LevelType.WARNING, f"safe_rmtree error: {e}")
        return False

    def cleanup_artifacts(self, csv_paths=(),excel_path=None, invalid_file=None, ppt_path=None):
        try:
            for p in csv_paths:
                self.safe_unlink(p)
            if excel_path:
                self.safe_unlink(excel_path)
                excel_dir = os.path.dirname(excel_path)
                if excel_dir and os.path.isdir(excel_dir) and not os.listdir(excel_dir):
                    self.safe_rmtree(excel_dir)
            if ppt_path:
                self.safe_unlink(ppt_path)
                ppt_dir = os.path.dirname(ppt_path)
                if ppt_dir and os.path.isdir(ppt_dir) and not os.listdir(ppt_dir):
                    self.safe_rmtree(ppt_dir)
            if invalid_file:
                self.safe_unlink(invalid_file)
                invalid_dir = os.path.dirname(invalid_file)
                if invalid_dir and os.path.isdir(invalid_dir) and not os.listdir(invalid_dir):
                    self.safe_rmtree(invalid_dir)
        except Exception as e:
            self.log_message(LevelType.WARNING, f"cleanup_artifacts error: {e}")

    def fetch_portfolio_recommendations(self, portfolio_id):
        try:
            collection = self.recommended_instance_collection
            query = {"portfolio_id": portfolio_id}
            pipeline = [
                {"$match": query},
                {"$addFields": {
                    "_status_sort": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$STATUS", ""]},
                                {"$not": ["$STATUS"]}
                            ]},
                            0,
                            1
                        ]
                    }
                }},
                {"$sort": {"_status_sort": 1}}
            ]
            cursor = collection.aggregate(pipeline)
            recs = cursor.to_list(length=None)
            self.log_message(LevelType.INFO, f"Found {len(recs)} recommendations for portfolio_id={portfolio_id}")
            return recs
        except Exception as e:
            self.log_message(LevelType.ERROR, f"fetch_portfolio_recommendations error: {e}")
            return []

    def fetch_invalid_records(self, portfolio_id: str, app_name: str) -> str:
        try:
            collection = self.input_remarks_collection
            query = {"portfolio_id": portfolio_id, "app_name": app_name}
            results = list(collection.find(query))
            if app_name.upper() == "CCA":
                col_map = {
                    'uuid': 'UUID',
                    'cloud_csp': 'Cloud',
                    'region': 'Region',
                    'instance type': 'Size',
                    'quantity': 'Quantity',
                    'monthly utilization (hourly)': 'Total number of hours per month',
                    'pricingModel': 'Pricing Model'
                }
            else:
                col_map ={
                    'uuid':'uuid',
                    'cloud_csp':'cloud_csp',
                    'region':'region',
                    'instance type':'instance type',
                    'max cpu%':'max cpu%',
                    'max mem used':'max mem used',
                    'max network bw':'max network bw',
                    'max disk bw used':'max disk bw used',
                    'max iops':'max iops',
                    'pricingModel':'Pricing Model', 
                    'uavg':'UAVG', 
                    'u95':'U95'
                }
            cleaned_data = []
            for doc in results:
                row = {new_col: doc.get(old_col, "") for old_col, new_col in col_map.items()}
                remarks_field = doc.get('Remarks', [])
                if isinstance(remarks_field, list):
                    messages = [item.get('Message', "") for item in remarks_field if isinstance(item, dict) and 'Message' in item]
                    row['Remarks'] = ", ".join([msg for msg in messages if msg])
                else:
                    row['Remarks'] = ""
                cleaned_data.append(row)
            if cleaned_data:
                df = pd.DataFrame(cleaned_data)

                with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
                    temp_path = tmp.name

                df.to_excel(temp_path, index=False)
                return temp_path
            return ""
        except Exception as e:
            self.log_message(LevelType.ERROR, f"Fetch_invalid_records error: {str(e)}")
            return ""

    def normalize_cca(self, doc: dict) -> dict:
        try:
            synonyms = {
                "Current Monthly Price": "Current Monthly Cost",
                "Monthly Price I": "Monthly Cost I",
                "Monthly Price II": "Monthly Cost II",
                "Monthly Price III": "Monthly Cost III",
                "Annual Cost I": "Annual Cost I (perf scaled)",
                "Annual Cost II": "Annual Cost II (perf scaled)",
                "Annual Cost III": "Annual Cost III (perf scaled)",
                "Number of Instances": "Number of Instances",
            }
            normalized = dict(doc)
            for src, dest in synonyms.items():
                if src in doc:
                    normalized[dest] = doc.get(src, "")
            for h in self.CCA_HEADERS:
                normalized.setdefault(h, "")
            if not normalized.get("Number of Instances"):
                normalized["Number of Instances"] = "1"
            return normalized
        except Exception as e:
            self.log_message(LevelType.ERROR, f"normalize_cca error: {str(e)}")
            return doc

    def build_cca_csv(self, portfolio_id: str, app_name: str) -> (bytes, str):
        try:
            docs = self.fetch_portfolio_recommendations(portfolio_id)
            invalid_file = self.fetch_invalid_records(portfolio_id, app_name)
            if not docs:
                return self.to_csv_bytes([], self.CCA_HEADERS), invalid_file
            rows = []
            field_map = {h: h for h in self.CCA_HEADERS}
            for d in docs:
                nd = self.normalize_cca(d)
                rows.append(self.row_from_doc(nd, self.CCA_HEADERS, field_map))
            return self.to_csv_bytes(rows, self.CCA_HEADERS), invalid_file
        except Exception as e:
            self.log_message(LevelType.ERROR, f"build_cca_csv error: {str(e)}")
            return b"", ""

    def get_portfolio_headroom(self, portfolio_id: str, app_name: str) -> str:
        try:
            portfolio_coll = self.portfolio_collection
            portfolio = portfolio_coll.find_one(
                {"_id": ObjectId(portfolio_id), "app_name": app_name},
                {"_id": 0, "headroom": 1},
            )
            if not portfolio:
                return ""
            hr = portfolio.get("headroom")
            return "" if hr is None else str(hr)
        except Exception as e:
            self.log_message(LevelType.ERROR, f"get_portfolio_headroom error: {str(e)}")
            return ""

    def normalize_eia(self, doc: dict, input_headroom: str) -> dict:
        try:
            normalized = dict(doc)
            if input_headroom is not None:
                normalized["Input Headroom"] = input_headroom
            else:
                normalized.setdefault("Input Headroom", "")
            for h in self.EIA_HEADERS:
                normalized.setdefault(h, "")
            return normalized
        except Exception as e:
            self.log_message(LevelType.ERROR, f"normalize_eia error: {str(e)}")
            return doc

    def to_csv_bytes(self, rows: list, headers: list) -> bytes:
        try:
            from io import StringIO
            sio = StringIO()
            writer = csv.DictWriter(sio, fieldnames=headers, extrasaction="ignore")
            writer.writeheader()
            writer.writerows(rows)
            return sio.getvalue().encode("utf-8")
        except Exception as e:
            self.log_message(LevelType.ERROR, f"to_csv_bytes error: {str(e)}")
            return b""

    def row_from_doc(self, doc: dict, headers: list, field_map: dict) -> dict:
        try:
            row = {h: doc.get(field_map.get(h, h), "") for h in headers}
            return row
        except Exception as e:
            self.log_message(LevelType.ERROR, f"row_from_doc error: {str(e)}")
            return dict.fromkeys(headers, "")

    def build_eia_csv(self, portfolio_id: str, app_name: str) -> (bytes, str):
        try:
            docs = self.fetch_portfolio_recommendations(portfolio_id)
            headroom = self.get_portfolio_headroom(portfolio_id, app_name)
            invalid_file = self.fetch_invalid_records(portfolio_id, app_name)
            if not docs:
                return self.to_csv_bytes([], self.EIA_HEADERS), invalid_file
            rows = []
            field_map = {h: h for h in self.EIA_HEADERS}
            for d in docs:
                nd = self.normalize_eia(d, headroom)
                rows.append(self.row_from_doc(nd, self.EIA_HEADERS, field_map))
            return self.to_csv_bytes(rows, self.EIA_HEADERS), invalid_file
        except Exception as e:
            self.log_message(LevelType.ERROR, f"build_eia_csv error: {str(e)}")
            return b"", ""

    def transform_cca_data(self, docs):
        """
        Transform CCA database documents to JSON format and calculate grandTotal.
        
        Args:
            docs: List of documents from recommended_instance_collection
            
        Returns:
            tuple: (transformed_data, grand_total)
        """
        try:
            transformed_data = []
            
            # Initialize sums for grandTotal
            sums = {
                'Number of Instances': 0,
                'Current Monthly Cost': 0,
                'Annual Cost': 0,
                'Monthly Cost I': 0,
                'Annual Cost I (perf scaled)': 0,
                'Annual Savings I': 0,
                'Perf Enhancement I': 0,
                'Monthly Cost II': 0,
                'Annual Cost II (perf scaled)': 0,
                'Annual Savings II': 0,
                'Perf Enhancement II': 0,
                'Monthly Cost III': 0,
                'Annual Cost III (perf scaled)': 0,
                'Annual Savings III': 0,
                'Perf Enhancement III': 0
            }
            
            perf_enhancement_i = []
            perf_enhancement_ii = []
            perf_enhancement_iii = []
            
            for doc in docs:
                # Transform to JSON format
                item = {
                    "id": doc.get("UUID", ""),
                    "data": {
                        "currentPlatform": {
                            "zone": doc.get("Zone", ""),
                            "instanceType": doc.get("Current Instance", ""),
                            "monthlyCost": safe_round(doc.get("Current Monthly Price", "")),
                            "annualCost": safe_round(doc.get("Annual Cost", "")),
                            "cspProvider": doc.get("CSP", ""),
                            "numberOfInstances": doc.get("Number of Instances", ""),
                            "pricingModel": doc.get("Pricing Model", ""),
                            "vCPU": doc.get("vCPU", ""),
                            "status": doc.get("STATUS", "")
                        },
                        "recommendations": [
                            {
                                "instanceType": doc.get("Recommendation I Instance", ""),
                                "vCPU": doc.get("vCPU I", ""),
                                "monthlyCost": safe_round(doc.get("Monthly Price I", "")),
                                "totalCost": safe_round(doc.get("Annual Cost I", "")),
                                "annualSavings": safe_round(doc.get("Annual Savings I", "")),
                                "savingsInPercentage": safe_round(doc.get("Savings % I", "")),
                                "perf": safe_round(doc.get("Perf Enhancement I", ""))
                            },
                            {
                                "instanceType": doc.get("Recommendation II Instance", ""),
                                "vCPU": doc.get("vCPU II", ""),
                                "monthlyCost": safe_round(doc.get("Monthly Price II", "")),
                                "totalCost": safe_round(doc.get("Annual Cost II", "")),
                                "annualSavings": safe_round(doc.get("Annual Savings II", "")),
                                "savingsInPercentage": safe_round(doc.get("Savings % II", "")),
                                "perf": safe_round(doc.get("Perf Enhancement II", ""))
                            },
                            {
                                "instanceType": doc.get("Recommendation III Instance", ""),
                                "vCPU": doc.get("vCPU III", ""),
                                "monthlyCost": safe_round(doc.get("Monthly Price III", "")),
                                "totalCost": safe_round(doc.get("Annual Cost III", "")),
                                "annualSavings": safe_round(doc.get("Annual Savings III", "")),
                                "savingsInPercentage": safe_round(doc.get("Savings % III", "")),
                                "perf": safe_round(doc.get("Perf Enhancement III", ""))
                            }
                        ]
                    },
                    "comments": doc.get("comments", "")
                }
                transformed_data.append(item)
                
                # Calculate sums for grandTotal
                # Map database field names to grandTotal keys
                field_mapping = {
                    'Number of Instances': 'Number of Instances',
                    'Current Monthly Cost': 'Current Monthly Price',  # DB field name
                    'Annual Cost': 'Annual Cost',
                    'Monthly Cost I': 'Monthly Price I',  # DB field name
                    'Annual Cost I (perf scaled)': 'Annual Cost I',  # DB field name
                    'Annual Savings I': 'Annual Savings I',
                    'Perf Enhancement I': 'Perf Enhancement I',
                    'Monthly Cost II': 'Monthly Price II',  # DB field name
                    'Annual Cost II (perf scaled)': 'Annual Cost II',  # DB field name
                    'Annual Savings II': 'Annual Savings II',
                    'Perf Enhancement II': 'Perf Enhancement II',
                    'Monthly Cost III': 'Monthly Price III',  # DB field name
                    'Annual Cost III (perf scaled)': 'Annual Cost III',  # DB field name
                    'Annual Savings III': 'Annual Savings III',
                    'Perf Enhancement III': 'Perf Enhancement III'
                }
                
                for sum_key, db_field in field_mapping.items():
                    if sum_key == 'Perf Enhancement I':
                        try:
                            val = doc.get(db_field, "")
                            if val and str(val).strip() not in ["", "-", "inf"]:
                                perf_enhancement_i.append(float(val))
                        except (ValueError, TypeError):
                            pass
                    elif sum_key == 'Perf Enhancement II':
                        try:
                            val = doc.get(db_field, "")
                            if val and str(val).strip() not in ["", "-", "inf"]:
                                perf_enhancement_ii.append(float(val))
                        except (ValueError, TypeError):
                            pass
                    elif sum_key == 'Perf Enhancement III':
                        try:
                            val = doc.get(db_field, "")
                            if val and str(val).strip() not in ["", "-", "inf"]:
                                perf_enhancement_iii.append(float(val))
                        except (ValueError, TypeError):
                            pass
                    else:
                        try:
                            val = doc.get(db_field, "")
                            if val and str(val).strip():
                                sums[sum_key] += float(val)
                        except (ValueError, TypeError):
                            pass
            
            # Calculate average performance enhancements
            sums['Perf Enhancement I'] = round(sum(perf_enhancement_i) / len(perf_enhancement_i), 2) if perf_enhancement_i else 0
            sums['Perf Enhancement II'] = round(sum(perf_enhancement_ii) / len(perf_enhancement_ii), 2) if perf_enhancement_ii else 0
            sums['Perf Enhancement III'] = round(sum(perf_enhancement_iii) / len(perf_enhancement_iii), 2) if perf_enhancement_iii else 0
            
            # Calculate savings percentages
            annual_cost = sums.get('Annual Cost', 0)
            if annual_cost > 0:
                sums['hSavingsInPercentage'] = round((sums.get('Annual Savings I', 0) / annual_cost) * 100, 2)
                sums['mSavingsInPercentage'] = round((sums.get('Annual Savings II', 0) / annual_cost) * 100, 2)
                sums['mdSavingsInPercentage'] = round((sums.get('Annual Savings III', 0) / annual_cost) * 100, 2)
            else:
                sums['hSavingsInPercentage'] = 0
                sums['mSavingsInPercentage'] = 0
                sums['mdSavingsInPercentage'] = 0
            
            # Round all values
            grand_total = {key: round(value, 2) for key, value in sums.items()}
            
            self.log_message(LevelType.INFO, f"Transformed {len(transformed_data)} CCA records")
            return transformed_data, grand_total
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"transform_cca_data error: {str(e)}")
            return [], {}

    def transform_eia_data(self, docs, portfolio_id):
        """
        Transform EIA database documents to JSON format and calculate grandTotal.
        
        Args:
            docs: List of documents from recommended_instance_collection
            portfolio_id: Portfolio ID for fetching headroom
            
        Returns:
            tuple: (transformed_data, grand_total)
        """
        try:
            transformed_data = []
            
            # Initialize sums for grandTotal
            sums = {
                'Current Monthly Price': 0.0,
                'Current Instance Energy Consumption (kwh)': 0.0,
                'Current Instance Emission': 0.0,
                'Monthly Price I': 0.0,
                'Instance Energy Consumption I (kwh)': 0.0,
                'Instance Emission I': 0.0,
                'Monthly Savings I': 0.0,
                'Monthly Price II': 0.0,
                'Instance Energy Consumption II (kwh)': 0.0,
                'Instance Emission II': 0.0,
                'Monthly Savings II': 0.0,
                'Perf Enhancement I': 0.0,
                'Perf Enhancement II': 0.0,
                'Untapped Capacity I': 0.0,
                'Untapped Capacity II': 0.0
            }
            
            perf_enhancement_i = []
            perf_enhancement_ii = []
            untapped_capacity_i = []
            untapped_capacity_ii = []
            
            for doc in docs:
                # Transform to JSON format
                item = {
                    "id": doc.get("UUID", ""),
                    "csp": doc.get("CSP", ""),
                    "data": {
                        "currentPlatform": {
                            "type": doc.get("Current Instance", ""),
                            "vCPU": doc.get("vCPU", ""),
                            "cost": safe_round(doc.get("Current Monthly Price", "")),
                            "power": safe_round(doc.get("Current Instance Energy Consumption (kwh)", "")),
                            "carbon": safe_round(doc.get("Current Instance Emission", "")),
                            "region": doc.get("Zone", ""),
                            "pricingModel": doc.get("Pricing Model", ""),
                            "status": doc.get("STATUS", "")
                        },
                        "recommendations": [
                            {
                                "type": doc.get("Recommendation I Instance", ""),
                                "vCPU": doc.get("vCPU I", ""),
                                "cost": safe_round(doc.get("Monthly Price I", "")),
                                "monthlySavings": safe_round(doc.get("Monthly Savings I", "")),
                                "power": safe_round(doc.get("Instance Energy Consumption I (kwh)", "")),
                                "carbon": safe_round(doc.get("Instance Emission I", "")),
                                "perf": safe_round(doc.get("Perf Enhancement I", "")),
                                "untappedCapacity": safe_round(doc.get("Untapped Capacity I", ""))
                            },
                            {
                                "type": doc.get("Recommendation II Instance", ""),
                                "vCPU": doc.get("vCPU II", ""),
                                "cost": safe_round(doc.get("Monthly Price II", "")),
                                "monthlySavings": safe_round(doc.get("Monthly Savings II", "")),
                                "power": safe_round(doc.get("Instance Energy Consumption II (kwh)", "")),
                                "carbon": safe_round(doc.get("Instance Emission II", "")),
                                "perf": safe_round(doc.get("Perf Enhancement II", "")),
                                "untappedCapacity": safe_round(doc.get("Untapped Capacity II", ""))
                            }
                        ]
                    },
                    "comments": doc.get("comments", "")
                }
                transformed_data.append(item)
                
                # Calculate sums for grandTotal
                for key in sums.keys():
                    if key == 'Perf Enhancement I':
                        try:
                            val = doc.get(key, "")
                            if val and str(val).strip() not in ["", "-", "inf"]:
                                perf_enhancement_i.append(float(val))
                        except (ValueError, TypeError):
                            pass
                    elif key == 'Perf Enhancement II':
                        try:
                            val = doc.get(key, "")
                            if val and str(val).strip() not in ["", "-", "inf"]:
                                perf_enhancement_ii.append(float(val))
                        except (ValueError, TypeError):
                            pass
                    elif key == 'Untapped Capacity I':
                        try:
                            val = doc.get(key, "")
                            if val and str(val).strip() not in ["", "-", "inf"]:
                                untapped_capacity_i.append(float(val))
                        except (ValueError, TypeError):
                            pass
                    elif key == 'Untapped Capacity II':
                        try:
                            val = doc.get(key, "")
                            if val and str(val).strip() not in ["", "-", "inf"]:
                                untapped_capacity_ii.append(float(val))
                        except (ValueError, TypeError):
                            pass
                    else:
                        try:
                            val = doc.get(key, "")
                            if val and str(val).strip():
                                sums[key] += float(val)
                        except (ValueError, TypeError):
                            pass
            
            # Calculate average performance enhancements
            sums['Perf Enhancement I'] = round(sum(perf_enhancement_i) / len(perf_enhancement_i), 2) if perf_enhancement_i else 0.0
            sums['Perf Enhancement II'] = round(sum(perf_enhancement_ii) / len(perf_enhancement_ii), 2) if perf_enhancement_ii else 0.0
            sums['Untapped Capacity I'] = round(sum(untapped_capacity_i) / len(untapped_capacity_i), 2) if untapped_capacity_i else 0.0
            sums['Untapped Capacity II'] = round(sum(untapped_capacity_ii) / len(untapped_capacity_ii), 2) if untapped_capacity_ii else 0.0
            
            # Round all values
            grand_total = {key: round(value, 2) for key, value in sums.items()}
            
            self.log_message(LevelType.INFO, f"Transformed {len(transformed_data)} EIA records")
            return transformed_data, grand_total
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"transform_eia_data error: {str(e)}")
            return [], {}

    def download_csv(self, portfolio_id: str, app_name: str, portfolio_name: str, user_email: str):
        try:
            self.log_message(LevelType.INFO, f"Starting download_csv for portfolio_id={portfolio_id}, app_name={app_name}")
            
            # Fetch recommendations from database
            docs = self.fetch_portfolio_recommendations(portfolio_id)
            if not docs:
                self.log_message(LevelType.WARNING, f"No recommendations found for portfolio_id={portfolio_id}")
                return
            
            # Fetch invalid records
            invalid_file = self.fetch_invalid_records(portfolio_id, app_name)
            
            # Transform data and calculate grandTotal based on app_name
            if app_name.upper() == "CCA":
                transformed_data, grand_total = self.transform_cca_data(docs)
            elif app_name.upper() == "EIA":
                transformed_data, grand_total = self.transform_eia_data(docs, portfolio_id)
            else:
                self.log_message(LevelType.ERROR, f"Invalid app_name: {app_name}")
                return
            
            # Generate file paths
            file_name = f"{portfolio_name}_{portfolio_id}"
            excel_file = f"{file_name}.xlsx"
            excel_path = os.path.join(self.results_path, excel_file)
            # Generate Excel file
            if app_name.upper() == "CCA":
                self.generate_excel_from_json(
                    {"data": transformed_data, "grandTotal": grand_total},
                    excel_path,
                    app_name,
                    user_email,
                    excel_file,
                    portfolio_id
                )
            elif app_name.upper() == "EIA":
                headroom = self.get_portfolio_headroom(portfolio_id, app_name)
                self.generate_excel_report(
                    transformed_data,
                    grand_total,
                    excel_path,
                    headroom,
                    app_name,
                    user_email,
                    excel_file,
                    portfolio_id
                )
            
            # Generate PPT file
            ppt_name = ""
            ppt_path = ""
            ppt_s3_key = ""
            
            if app_name.upper() == "CCA":
                ppt_template = os.path.join(ROOT_DIR, "AMD_EPYC_PPT_TEMPLATE.pptx")
                if os.path.exists(excel_path) and os.path.exists(ppt_template):
                    ppt_name = self.generate_ppt(
                        excel_path,
                        user_email,
                        ppt_template,
                        self.results_path,
                        app_name,
                        portfolio_id
                    )
                    if ppt_name:
                        ppt_path = os.path.join(self.results_path, ppt_name)
            elif app_name.upper() == "EIA":
                ppt_template = os.path.join(ROOT_DIR, "EIA_Presentation_Screen.pptx")
                if os.path.exists(excel_path) and os.path.exists(ppt_template):
                    ppt_name = f"{file_name}_EIA.pptx"
                    ppt_path = os.path.join(self.results_path, ppt_name)
                    try:
                        self.generate_ppt_from_excel(excel_path, ppt_template, ppt_path, user_email)
                        self.log_message(LevelType.INFO, f"EIA PPT generated: {ppt_path}")
                    except Exception as e:
                        self.log_message(LevelType.ERROR, f"Failed to generate EIA PPT: {str(e)}")
                        ppt_name = ""
                        ppt_path = ""
            
            # Upload PPT to S3 if generated
            if ppt_path and os.path.exists(ppt_path):
                with open(ppt_path, 'rb') as f:
                    ppt_bytes = f.read()
                
                ppt_s3_key = self.upload_file_to_s3(
                    file_bytes=ppt_bytes,
                    app_name=app_name,
                    user_email=user_email,
                    file_name=ppt_name,
                    sub_folder="output"
                )
                
                if ppt_s3_key:
                    self.portfolio_collection.update_one(
                        {"_id": ObjectId(portfolio_id)},
                        {"$set": {"ppt_s3_key": ppt_s3_key}}
                    )
                    self.log_message(LevelType.INFO, f"PPT S3 path stored: {ppt_s3_key}")
            
            # Send email notification
            zip_filename = f"Recommendations_{portfolio_name}.zip"
            invalid_file_name = f"{portfolio_name}_invalid_records.xlsx"
            
            msg, _ = self.send_recommendations_completed_email(
                to_email=user_email,
                app_name=app_name,
                portfolio_id=portfolio_id,
                portfolio_name=portfolio_name,
                excel_path=excel_path,
                pptx_path=ppt_path,
                invalid_file=invalid_file,
                excel_filename=excel_file,
                pptx_filename=ppt_name,
                zip_filename=zip_filename,
                invalid_file_name=invalid_file_name
            )
            
            self.log_message(LevelType.INFO, f"Email notification: {msg} for portfolio_id: {portfolio_id}")
            
            # Cleanup artifacts
            self.cleanup_artifacts(excel_path=excel_path, invalid_file=invalid_file, ppt_path=ppt_path)
            
            # Delete invalid records from collection
            collection = self.input_remarks_collection
            query = {"portfolio_id": portfolio_id, "app_name": app_name}
            collection.delete_many(query)
            
            self.log_message(LevelType.INFO, f"Completed download_csv for portfolio_id={portfolio_id}")
            
        except Exception as e:
            self.log_message(LevelType.ERROR, f"download_csv error: {str(e)}")

    def add_notification(self, user_email: str, app_name: str, portfolio_id: str, portfolio_name: str, purpose: str = None, message: str = None, title: str = None, is_seen: bool = False) -> str:
        try:
            doc = {
                "user_email": user_email,
                "app_name": app_name,
                "portfolio_id": portfolio_id,
                "portfolio_name": portfolio_name,
                "is_seen": is_seen,
                "created_at": datetime.utcnow(),
                "updated_at": datetime.utcnow()
            }
            if purpose:
                doc["purpose"] = purpose
            if message:
                doc["message"] = message
            if title:
                doc["title"] = title
            insert_result = self.notification_collection.insert_one(doc)
            return str(insert_result.inserted_id)
        except Exception as e:
            self.log_message(LevelType.ERROR, f"add_notification error: {str(e)}")
            return ""

    def generate_user_name_from_email(self, user_email: str) -> str:
        """
        Generate a user_name from user_email.
        
        Rules:
        - Split email by '@'
        - Replace '.' with '_' in local part
        - Take first part of domain before '.' 
        - Combine with '_' → local_part_domain_part
        """
        try:
            local_part, domain_part = user_email.split("@")
            local_part = local_part.replace(".", "_")
            domain_part = domain_part.split(".")[0]
            return f"{local_part}_{domain_part}"
        except Exception:
            # Fallback: replace @ and dots with underscores if unexpected format
            return user_email.replace("@", "_").replace(".", "_")

    def build_s3_key(self, app_name: str, user_name: str, file_name: str, sub_folder : str) -> str:
        """
        Build S3 object key in format:
        main_folder/app_name/user_name/input/<file_name>_input
        """
        return f"{self.MAIN_FOLDER}/{app_name}/{user_name}/{sub_folder}/{file_name}"

    def upload_file_to_s3(self, file_bytes: bytes, app_name: str, user_email: str, file_name: str, sub_folder: str) -> str:
        """
        Upload file bytes to S3 synchronously.
        Returns the S3 key.
        """
        s3_client = boto3.client(
            "s3",
            aws_access_key_id=self.AWS_ACCESS_KEY,
            aws_secret_access_key=self.AWS_SECRET_KEY,
            region_name=self.AWS_REGION,
        )
        user_name = self.generate_user_name_from_email(user_email)
        s3_key = f"{self.MAIN_FOLDER}/{app_name}/{user_name}/{sub_folder}/{file_name}"
        # Determine the correct ContentType based on file extension
        content_type = "application/octet-stream"  # default
        if file_name.lower().endswith('.xlsx'):
            content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif file_name.lower().endswith('.pptx'):
            content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif file_name.lower().endswith('.csv'):
            content_type = "text/csv"
        try:
            s3_client.put_object(
                Bucket=self.BUCKET_NAME,
                Key=s3_key,
                Body=file_bytes,
                ContentType=content_type
            )
            return s3_key
        except Exception as e:
            print(f"Error uploading to S3: {str(e)}")
            return ""
    @staticmethod
    def extract_org_and_user_from_email(email: str):
        if '@' in email:
            user, org_part = email.split('@', 1)
            org = org_part.split('.')[0].lower() if '.' in org_part else org_part.lower()
            return user.lower(), org
        else:
            lowered = email.lower()
            return lowered, lowered

    @staticmethod
    def convert_to_utc(dt):
        if dt is None:
            return None
        if dt.tzinfo is None:
            return dt.replace(tzinfo=timezone.utc)
        return dt.astimezone(timezone.utc)
    @staticmethod
    def get_cca_pipeline(portfolio_id):
        return [
                {"$match": {"portfolio_id": portfolio_id}},
                {"$group": {
                    "_id": None,
                    "Current_Vcpus": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$vCPU", "-"]},
                                {"$eq": ["$vCPU", None]},
                                {"$eq": ["$vCPU", ""]}
                            ]},
                            0,
                            {"$toDouble": "$vCPU"}
                        ]
                    }},
                    "HC_VCPUs": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$vCPU I", "-"]},
                                {"$eq": ["$vCPU I", None]},
                                {"$eq": ["$vCPU I", ""]}
                            ]},
                            0,
                            {"$toDouble": "$vCPU I"}
                        ]
                    }},
                    "M_VCPUs": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$vCPU II", "-"]},
                                {"$eq": ["$vCPU II", None]},
                                {"$eq": ["$vCPU II", ""]}
                            ]},
                            0,
                            {"$toDouble": "$vCPU II"}
                        ]
                    }},
                    "MD_VCPUs": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$vCPU III", "-"]},
                                {"$eq": ["$vCPU III", None]},
                                {"$eq": ["$vCPU III", ""]}
                            ]},
                            0,
                            {"$toDouble": "$vCPU III"}
                        ]
                    }},
                    "Sum_H_Saving": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Annual Savings I", "-"]},
                                {"$eq": ["$Annual Savings I", None]},
                                {"$eq": ["$Annual Savings I", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Annual Savings I"}
                        ]
                    }},
                    "Sum_M_Saving": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Annual Savings II", "-"]},
                                {"$eq": ["$Annual Savings II", None]},
                                {"$eq": ["$Annual Savings II", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Annual Savings II"}
                        ]
                    }},
                    "Sum_MD_Saving": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Annual Savings III", "-"]},
                                {"$eq": ["$Annual Savings III", None]},
                                {"$eq": ["$Annual Savings III", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Annual Savings III"}
                        ]
                    }},
                    "count_perfI": {"$sum": {"$cond": [{"$in": ["$Perf Enhancement I", ["-", None, ""]]}, 0, 1]}},
                    "count_perfII": {"$sum": {"$cond": [{"$in": ["$Perf Enhancement II", ["-", None, ""]]}, 0, 1]}},
                    "count_perfIII": {"$sum": {"$cond": [{"$in": ["$Perf Enhancement III", ["-", None, ""]]}, 0, 1]}},
                    "sum_perfI": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Perf Enhancement I", "-"]},
                                {"$eq": ["$Perf Enhancement I", None]},
                                {"$eq": ["$Perf Enhancement I", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Perf Enhancement I"}
                        ]
                    }},
                    "sum_perfII": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Perf Enhancement II", "-"]},
                                {"$eq": ["$Perf Enhancement II", None]},
                                {"$eq": ["$Perf Enhancement II", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Perf Enhancement II"}
                        ]
                    }},
                    "sum_perfIII": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Perf Enhancement III", "-"]},
                                {"$eq": ["$Perf Enhancement III", None]},
                                {"$eq": ["$Perf Enhancement III", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Perf Enhancement III"}
                        ]
                    }},
                    
                    # --- Annual Costs ---
                    "current_cost": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Annual Cost", "-"]},
                                {"$eq": ["$Annual Cost", None]},
                                {"$eq": ["$Annual Cost", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Annual Cost"}
                        ]
                    }},
                    "HC_cost": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Annual Cost I", "-"]},
                                {"$eq": ["$Annual Cost I", None]},
                                {"$eq": ["$Annual Cost I", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Annual Cost I"}
                        ]
                    }},
                    "M_cost": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Annual Cost II", "-"]},
                                {"$eq": ["$Annual Cost II", None]},
                                {"$eq": ["$Annual Cost II", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Annual Cost II"}
                        ]
                    }},
                    "M_D_cost": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Annual Cost III", "-"]},
                                {"$eq": ["$Annual Cost III", None]},
                                {"$eq": ["$Annual Cost III", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Annual Cost III"}
                        ]
                    }},
                    
                    # --- Recommendation Counts ---
                    "recommendation_instances_count": {
                        "$sum": {
                            "$cond": [
                                {"$or": [
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price I", "-"]}, {"$eq": ["$Monthly Price I", None]}, {"$eq": ["$Monthly Price I", ""]}]},
                                            0, {"$toDouble": "$Monthly Price I"}
                                        ]
                                    }, 0]},
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price II", "-"]}, {"$eq": ["$Monthly Price II", None]}, {"$eq": ["$Monthly Price II", ""]}]},
                                            0, {"$toDouble": "$Monthly Price II"}
                                        ]
                                    }, 0]},
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price III", "-"]}, {"$eq": ["$Monthly Price III", None]}, {"$eq": ["$Monthly Price III", ""]}]},
                                            0, {"$toDouble": "$Monthly Price III"}
                                        ]
                                    }, 0]}
                                ]},
                                1, 0
                            ]
                        }
                    },
                    "unique_zones": {
                        "$addToSet": {
                            "$cond": [
                                {"$or": [
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price I", "-"]}, {"$eq": ["$Monthly Price I", None]}, {"$eq": ["$Monthly Price I", ""]}]},
                                            0, {"$toDouble": "$Monthly Price I"}
                                        ]
                                    }, 0]},
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price II", "-"]}, {"$eq": ["$Monthly Price II", None]}, {"$eq": ["$Monthly Price II", ""]}]},
                                            0, {"$toDouble": "$Monthly Price II"}
                                        ]
                                    }, 0]},
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price III", "-"]}, {"$eq": ["$Monthly Price III", None]}, {"$eq": ["$Monthly Price III", ""]}]},
                                            0, {"$toDouble": "$Monthly Price III"}
                                        ]
                                    }, 0]}
                                ]},
                                "$Zone I",
                                None
                            ]
                        }
                    },
                    "recommendation_date": {"$min": "$created_at"},
                }},
                {"$project": {
                    "_id": 0,
                    "Current_Vcpus": 1,
                    "HC_VCPUs": 1,
                    "M_VCPUs": 1,
                    "MD_VCPUs": 1,
                    "recommendation_date": 1,
                    "H_Perf": {
                        "$cond": [
                            {"$eq": ["$count_perfI", 0]},
                            0,
                            {"$divide": ["$sum_perfI", "$count_perfI"]}
                        ]
                    },
                    "M_Perf": {
                        "$cond": [
                            {"$eq": ["$count_perfII", 0]},
                            0,
                            {"$divide": ["$sum_perfII", "$count_perfII"]}
                        ]
                    },
                    "MD_Perf": {
                        "$cond": [
                            {"$eq": ["$count_perfIII", 0]},
                            0,
                            {"$divide": ["$sum_perfIII", "$count_perfIII"]}
                        ]
                    },
                    "Total_Perf": {
                        "$let": {
                            "vars": {
                                "sumPerf": {"$add": ["$sum_perfI", "$sum_perfII", "$sum_perfIII"]},
                                "countPerf": {"$add": ["$count_perfI", "$count_perfII", "$count_perfIII"]}
                            },
                            "in": {
                                "$cond": [
                                    {"$eq": ["$countPerf", 0]},
                                    0,
                                    {"$divide": ["$sumPerf", "$countPerf"]}
                                ]
                            }
                        }
                    },
                    "H_Saving": "$Sum_H_Saving",
                    "M_Saving": "$Sum_M_Saving",
                    "MD_Saving": "$Sum_MD_Saving",
                    
                    # Cost fields
                    "current_cost": 1,
                    "HC_cost": 1,
                    "M_cost": 1,
                    "M_D_cost": 1,
                    "HC_saving": "$Sum_H_Saving",
                    "M_saving": "$Sum_M_Saving",
                    "M_D_saving": "$Sum_MD_Saving",
                    
                    # Count fields
                    "recommendation_instances_count": 1,
                    "recommendation_regions_count": {
                        "$size": {
                            "$filter": {
                                "input": "$unique_zones",
                                "as": "z",
                                "cond": {"$and": [{"$ne": ["$$z", None]}, {"$ne": ["$$z", ""]}, {"$ne": ["$$z", "-"]}]}
                            }
                        }
                    }
                }}
            ]
    @staticmethod
    def get_eia_pipeline(portfolio_id):
        return [
                {"$match": {"portfolio_id": portfolio_id}},
                {"$group": {
                    "_id": None,
                    "Current_Vcpus": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$vCPU", "-"]},
                                {"$eq": ["$vCPU", None]},
                                {"$eq": ["$vCPU", ""]}
                            ]},
                            0,
                            {"$toDouble": "$vCPU"}
                        ]
                    }},

                    "HC_VCPUs": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$vCPU I", "-"]},
                                {"$eq": ["$vCPU I", None]},
                                {"$eq": ["$vCPU I", ""]}
                            ]},
                            0,
                            {"$toDouble": "$vCPU I"}
                        ]
                    }},

                    "M_VCPUs": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$vCPU II", "-"]},
                                {"$eq": ["$vCPU II", None]},
                                {"$eq": ["$vCPU II", ""]}
                            ]},
                            0,
                            {"$toDouble": "$vCPU II"}
                        ]
                    }},

                    "Sum_H_Saving": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Monthly Savings I", "-"]},
                                {"$eq": ["$Monthly Savings I", None]},
                                {"$eq": ["$Monthly Savings I", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Monthly Savings I"}
                        ]
                    }},

                    "Sum_M_Saving": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Monthly Savings II", "-"]},
                                {"$eq": ["$Monthly Savings II", None]},
                                {"$eq": ["$Monthly Savings II", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Monthly Savings II"}
                        ]
                    }},

                    "Sum_MD_Saving": {"$sum": 0},
                    "H_carboncount": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Instance Energy Consumption I (kwh)", "-"]},
                                {"$eq": ["$Instance Energy Consumption I (kwh)", None]},
                                {"$eq": ["$Instance Energy Consumption I (kwh)", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Instance Energy Consumption I (kwh)"}
                        ]
                    }},

                    "M_carboncount": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Instance Energy Consumption II (kwh)", "-"]},
                                {"$eq": ["$Instance Energy Consumption II (kwh)", None]},
                                {"$eq": ["$Instance Energy Consumption II (kwh)", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Instance Energy Consumption II (kwh)"}
                        ]
                    }},

                    "H_energycount": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Instance Emission I", "-"]},
                                {"$eq": ["$Instance Emission I", None]},
                                {"$eq": ["$Instance Emission I", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Instance Emission I"}
                        ]
                    }},

                    "M_energycount": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Instance Emission II", "-"]},
                                {"$eq": ["$Instance Emission II", None]},
                                {"$eq": ["$Instance Emission II", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Instance Emission II"}
                        ]
                    }},

                    "count_perfI": {"$sum": {"$cond": [{"$in": ["$Perf Enhancement I", ["-", None, ""]]}, 0, 1]}},
                    "count_perfII": {"$sum": {"$cond": [{"$in": ["$Perf Enhancement II", ["-", None, ""]]}, 0, 1]}},
                    "count_perfIII": {"$sum": 0},
                    "sum_perfI": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Perf Enhancement I", "-"]},
                                {"$eq": ["$Perf Enhancement I", None]},
                                {"$eq": ["$Perf Enhancement I", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Perf Enhancement I"}
                        ]
                    }},

                    "sum_perfII": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Perf Enhancement II", "-"]},
                                {"$eq": ["$Perf Enhancement II", None]},
                                {"$eq": ["$Perf Enhancement II", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Perf Enhancement II"}
                        ]
                    }},
                    "sum_perfIII": {"$sum": 0},
                    
                    # --- Monthly Costs ---
                    "current_cost": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Current Monthly Price", "-"]},
                                {"$eq": ["$Current Monthly Price", None]},
                                {"$eq": ["$Current Monthly Price", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Current Monthly Price"}
                        ]
                    }},
                    "O_cost": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Monthly Price I", "-"]},
                                {"$eq": ["$Monthly Price I", None]},
                                {"$eq": ["$Monthly Price I", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Monthly Price I"}
                        ]
                    }},
                    "G_cost": {"$sum": {
                        "$cond": [
                            {"$or": [
                                {"$eq": ["$Monthly Price II", "-"]},
                                {"$eq": ["$Monthly Price II", None]},
                                {"$eq": ["$Monthly Price II", ""]}
                            ]},
                            0,
                            {"$toDouble": "$Monthly Price II"}
                        ]
                    }},
                    
                    # --- Recommendation Counts ---
                    "recommendation_instances_count": {
                        "$sum": {
                            "$cond": [
                                {"$or": [
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price I", "-"]}, {"$eq": ["$Monthly Price I", None]}, {"$eq": ["$Monthly Price I", ""]}]},
                                            0, {"$toDouble": "$Monthly Price I"}
                                        ]
                                    }, 0]},
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price II", "-"]}, {"$eq": ["$Monthly Price II", None]}, {"$eq": ["$Monthly Price II", ""]}]},
                                            0, {"$toDouble": "$Monthly Price II"}
                                        ]
                                    }, 0]}
                                ]},
                                1, 0
                            ]
                        }
                    },
                    "unique_zones": {
                        "$addToSet": {
                            "$cond": [
                                {"$or": [
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price I", "-"]}, {"$eq": ["$Monthly Price I", None]}, {"$eq": ["$Monthly Price I", ""]}]},
                                            0, {"$toDouble": "$Monthly Price I"}
                                        ]
                                    }, 0]},
                                    {"$gt": [{
                                        "$cond": [
                                            {"$or": [{"$eq": ["$Monthly Price II", "-"]}, {"$eq": ["$Monthly Price II", None]}, {"$eq": ["$Monthly Price II", ""]}]},
                                            0, {"$toDouble": "$Monthly Price II"}
                                        ]
                                    }, 0]}
                                ]},
                                "$Zone",
                                None
                            ]
                        }
                    },
                    "recommendation_date": {"$min": "$created_at"},
                }},
                {"$project": {
                    "_id": 0,
                    "Current_Vcpus": 1,
                    "HC_VCPUs": 1,
                    "M_VCPUs": 1,
                    "H_carboncount": 1,
                    "M_carboncount": 1,
                    "H_energycount": 1,
                    "M_energycount": 1,
                    "recommendation_date": 1,
                    "H_Perf": {
                        "$cond": [
                            {"$eq": ["$count_perfI", 0]},
                            0,
                            {"$divide": ["$sum_perfI", "$count_perfI"]}
                        ]
                    },
                    "M_Perf": {
                        "$cond": [
                            {"$eq": ["$count_perfII", 0]},
                            0,
                            {"$divide": ["$sum_perfII", "$count_perfII"]}
                        ]
                    },
                    "MD_Perf": {"$literal": 0},
                    "Total_Perf": {
                        "$let": {
                            "vars": {
                                "sumPerf": {"$add": ["$sum_perfI", "$sum_perfII"]},
                                "countPerf": {"$add": ["$count_perfI", "$count_perfII"]}
                            },
                            "in": {
                                "$cond": [
                                    {"$eq": ["$countPerf", 0]},
                                    0,
                                    {"$divide": ["$sumPerf", "$countPerf"]}
                                ]
                            }
                        }
                    },
                    "H_Saving": "$Sum_H_Saving",
                    "M_Saving": "$Sum_M_Saving",
                    "MD_Saving": {"$literal": 0},
                    
                    # Cost fields
                    "current_cost": 1,
                    "O_cost": 1,
                    "G_cost": 1,
                    "O_saving": "$Sum_H_Saving",
                    "G_saving": "$Sum_M_Saving",
                    
                    # Count fields
                    "recommendation_instances_count": 1,
                    "recommendation_regions_count": {
                        "$size": {
                            "$filter": {
                                "input": "$unique_zones",
                                "as": "z",
                                "cond": {"$and": [{"$ne": ["$$z", None]}, {"$ne": ["$$z", ""]}, {"$ne": ["$$z", "-"]}]}
                            }
                        }
                    }
                }}
            ]
    
    def get_unsupported_instances(self,portfolio_id, app_name, user):
        condition = None
        if app_name.upper() == "EIA":
            condition = {
                "$and": [
                    {"$eq": ["$Recommendation I Instance", "-"]},
                    {"$eq": ["$Recommendation II Instance", "-"]},
                ]
            }
        else:
            condition = {
                "$and": [
                    {"$eq": ["$Recommendation I Instance", "-"]},
                    {"$eq": ["$Recommendation II Instance", "-"]},
                    {"$eq": ["$Recommendation III Instance", "-"]},
                ]
            }

        pipeline = [
            {"$match": {"portfolio_id": portfolio_id}},
            {"$project": {
                "CSP": 1,
                "Zone": 1,
                "Current Instance": 1,
                "created_at": 1,
                "Recommendation I Instance": 1,
                "Recommendation II Instance": 1,
                "Recommendation III Instance": 1,
            }},
            {"$addFields": {
                "matching_recs": condition
            }},
            {"$match": {"matching_recs": True}}
        ]

        etl_cursor = self.recommended_instance_collection.aggregate(pipeline)
        docs = list(etl_cursor)
        result = []
        for rec in docs:
            created_at = rec.get("created_at")
            result.append({
                "portfolio_id": portfolio_id,
                "cloud": rec.get("CSP"),
                "region": rec.get("Zone"),
                "current_instance": rec.get("Current Instance"),
                "created_at": created_at,
                "app_name": app_name,
                "user": user
            })
        return result

    def to_float_safe(self, v):
        """
        Convert Monthly Price I value to float if possible;
        return None for '-', '', None, invalid.
        """
        if v in ("", "-", None):
            return None
        try:
            return float(v)
        except (ValueError, TypeError):
            return None

    def is_valid_zone(self, z):
        """Check if a zone value is meaningful."""
        return z not in (None, "", "-")

    def build_portfolio_query(self, raw_pid):
        """
        Build a flexible query for portfolio_id.
        """
        clauses = []
        if raw_pid is None:
            return {"portfolio_id": None}

        if isinstance(raw_pid, str):
            clauses.append({"portfolio_id": raw_pid})
            trimmed = raw_pid.strip()
            if trimmed != raw_pid:
                clauses.append({"portfolio_id": trimmed})
            try:
                from bson.errors import InvalidId
                oid = ObjectId(trimmed)
                clauses.append({"portfolio_id": oid})
            except (InvalidId, Exception):
                pass
        else:
            clauses.append({"portfolio_id": raw_pid})

        if len(clauses) == 1:
            return clauses[0]
        return {"$or": clauses}



    def get_instance_and_region_counts(self, raw_pid, app_name):
        """
        Count recommendation instances and unique regions for a portfolio.

        Returns: (instances_count, regions_count)
        """
        query = self.build_portfolio_query(raw_pid)
        
        instances_count = 0
        regions = set()

        zone_field = "Zone I" if app_name == "CCA" else "Zone"

        inst_cursor = self.recommended_instance_collection.find(
            query,
            {
                "Monthly Price I": 1,
                zone_field: 1,
            }
        )

        for inst in inst_cursor:
            price_I = self.to_float_safe(inst.get("Monthly Price I"))
            zone = inst.get(zone_field)

            if price_I is not None and price_I > 0:
                instances_count += 1

                if self.is_valid_zone(zone):
                    regions.add(zone)

        return instances_count, len(regions)

    def etl_process_portfolio(self, portfolio_id):
        try:
            start_time = datetime.now()
            self.log_message("info", f"ETL Processing  for portfolio {portfolio_id} started.")
            doc = self.portfolio_collection.find_one(
                {"_id": ObjectId(portfolio_id)},
                {
                    "user_email": 1,
                    "cloud_provider": 1,
                    "created_at": 1,
                    "updated_at": 1,
                    "app_name": 1,
                    "name": 1
                }
            )
            if not doc:
                self.log_message("error", f"Portfolio {portfolio_id} not found, skipping.")
                return False

            app_name = doc.get("app_name", "").upper()
            user, org = self.extract_org_and_user_from_email(doc.get("user_email", ""))
            created_at_utc = self.convert_to_utc(doc.get("created_at"))
            portfolio_name = doc.get("name", "")

            # Single aggregation pipeline call that includes costs, savings, counts, vCPUs, and performance
            if app_name == "EIA":
                pipeline = self.get_eia_pipeline(portfolio_id)
            else:
                pipeline = self.get_cca_pipeline(portfolio_id)

            aggregate_result = next(self.recommended_instance_collection.aggregate(pipeline), None)

            if aggregate_result:
                if app_name == "EIA":
                    h_vcpus = aggregate_result.get("HC_VCPUs", 0)
                    m_vcpus = aggregate_result.get("M_VCPUs", 0)
                    total_vcpus_avg = (h_vcpus + m_vcpus) / 2 if (h_vcpus + m_vcpus) > 0 else 0
                    total_perf_avg = (aggregate_result.get("H_Perf", 0) + aggregate_result.get("M_Perf", 0)) / 2
                    final_result = {
                        "org": org,
                        "app_name": doc.get("app_name"),
                        "user_email": doc.get("user_email"),
                        "portfolio_id": portfolio_id,
                        "portfolio_name": portfolio_name,
                        "user": user,
                        "cloud": doc.get("cloud_provider"),
                        "portfolio_created_date": created_at_utc,
                        "recommendation_date": aggregate_result.get("recommendation_date", None),
                        "current_vCPUs": aggregate_result.get("Current_Vcpus", 0),
                        "O_vCPUs": h_vcpus,
                        "G_vCPUs": m_vcpus,
                        "total_vCPUs": total_vcpus_avg,
                        "O_perf": aggregate_result.get("H_Perf", 0),
                        "G_perf": aggregate_result.get("M_Perf", 0),
                        "total_perf": total_perf_avg,
                        "O_saving": aggregate_result.get("O_saving", 0),
                        "G_saving": aggregate_result.get("G_saving", 0),
                        "O_carboncount": aggregate_result.get("H_carboncount", 0),
                        "G_carboncount": aggregate_result.get("M_carboncount", 0),
                        "O_energycount": aggregate_result.get("H_energycount", 0),
                        "G_energycount": aggregate_result.get("M_energycount", 0),
                        "M_D_perf": 0,
                        "M_D_saving": 0,
                        
                        # New Fields (now directly from aggregate_result)
                        "current_cost": aggregate_result.get("current_cost", 0),
                        "O_cost": aggregate_result.get("O_cost", 0),
                        "G_cost": aggregate_result.get("G_cost", 0),
                        "recommendation_instances_count": aggregate_result.get("recommendation_instances_count", 0),
                        "recommendation_regions_count": aggregate_result.get("recommendation_regions_count", 0),
                    }
                else:
                    h_vcpus = aggregate_result.get("HC_VCPUs", 0)
                    m_vcpus = aggregate_result.get("M_VCPUs", 0)
                    md_vcpus = aggregate_result.get("MD_VCPUs", 0)
                    total_vcpus_avg = (h_vcpus + m_vcpus + md_vcpus) / 3 if (h_vcpus + m_vcpus + md_vcpus) > 0 else 0
                    total_perf_avg = (
                        aggregate_result.get("H_Perf", 0) +
                        aggregate_result.get("M_Perf", 0) +
                        aggregate_result.get("MD_Perf", 0)
                    ) / 3
                    final_result = {
                        "org": org,
                        "app_name": doc.get("app_name"),
                        "user_email": doc.get("user_email"),
                        "portfolio_id": portfolio_id,
                        "portfolio_name": portfolio_name,
                        "user": user,
                        "cloud": doc.get("cloud_provider"),
                        "portfolio_created_date": created_at_utc,
                        "recommendation_date": aggregate_result.get("recommendation_date", None),
                        "current_vCPUs": aggregate_result.get("Current_Vcpus", 0),
                        "HC_vCPUs": h_vcpus,
                        "M_vCPUs": m_vcpus,
                        "M_D_vCPUs": md_vcpus,
                        "total_vCPUs": total_vcpus_avg,
                        "HC_perf": aggregate_result.get("H_Perf", 0),
                        "M_perf": aggregate_result.get("M_Perf", 0),
                        "M_D_perf": aggregate_result.get("MD_Perf", 0),
                        "total_perf": total_perf_avg,
                        "HC_saving": aggregate_result.get("HC_saving", 0),
                        "M_saving": aggregate_result.get("M_saving", 0),
                        "M_D_saving": aggregate_result.get("M_D_saving", 0),
                        
                        # New Fields (now directly from aggregate_result)
                        "current_cost": aggregate_result.get("current_cost", 0),
                        "HC_cost": aggregate_result.get("HC_cost", 0),
                        "M_cost": aggregate_result.get("M_cost", 0),
                        "M_D_cost": aggregate_result.get("M_D_cost", 0),
                        "recommendation_instances_count": aggregate_result.get("recommendation_instances_count", 0),
                        "recommendation_regions_count": aggregate_result.get("recommendation_regions_count", 0),
                    }
                final_result["created_at"] = datetime.utcnow()
                # Insert new analytics entry (keeping historical entries)
                self.recommended_analytics.insert_one(final_result)
            else:
                self.log_message("warning", f"No recommendation data for portfolio {portfolio_id} for ETL")
            unsupported_instances = self.get_unsupported_instances(portfolio_id, app_name, user)
            if unsupported_instances:
                self.unsupported_recommendation_analytics.insert_many(unsupported_instances)
            end_time = datetime.now()
            self.log_message("info", f"ETL of Portfolio {portfolio_id} for {app_name} execution completed in {end_time - start_time}")
            return True
        except Exception as e:
            self.log_message(LevelType.ERROR, f"ETL processing error: {str(e)}")

    def handle_bulk_data(self):
        """handle bulk recommendation"""
        try:
            if check_cs_db_connection():
                self.log_message(LevelType.INFO, f"✅ DB connection OK CS_DB_URL : {CS_DB_URL}")
            else:
                self.log_message(LevelType.INFO, f"❌ DB connection failed CS_DB_URL : {CS_DB_URL}")
            portfolio_ids = self.fetch_and_sort_portfolioids()
            self.log_message(LevelType.INFO, f"QUEUED/INPROGRESS portfolio IDs: {portfolio_ids}")
            if not portfolio_ids:
                self.log_message(LevelType.INFO, "No portfolios with queued batches found.")
                return
            for portfolio_id in portfolio_ids:
                portfolio_doc = self.portfolio_collection.find_one({"_id": ObjectId(portfolio_id)})
                if not portfolio_doc:
                    self.log_message(LevelType.INFO, "No data found for portfolio")
                    continue
                user_email = portfolio_doc.get("user_email", "")
                portfolio_name = portfolio_doc.get("name", "")
                if not portfolio_doc:
                    self.log_message(LevelType.ERROR, f"Portfolio {portfolio_id} not found", ErrorCode=-1)
                    continue
                recommendation_status = portfolio_doc.get("recommendation_status","")
                if recommendation_status not in [RecommendationStatus.QUEUE, RecommendationStatus.IN_PROGRESS]:
                    self.log_message(LevelType.INFO, f"Skipping portfolio {portfolio_id} with status {recommendation_status}")
                    continue
                if recommendation_status == RecommendationStatus.QUEUE:
                    self.portfolio_collection.update_one({"_id": ObjectId(portfolio_id)}, {"$set": {"recommendation_status": RecommendationStatus.IN_PROGRESS, "recommendation_percentage": 0, "submittedForRecommendations" : True}})
                    self.log_message(LevelType.INFO, f"Portfolio {portfolio_id} status updated to INPROGRESS")
                cloud_csp = portfolio_doc.get("cloud_provider","")
                app_name = portfolio_doc.get("app_name","")
                failed = False
                if app_name.upper() not in ["CCA", "EIA"]:
                    self.log_message(LevelType.ERROR, f"Invalid application provided for portfolio {portfolio_id}")
                    failed = True
                if cloud_csp.upper() not in self.ALLOWED_PROVIDERS:
                    self.log_message(LevelType.ERROR, f"Invalid provider for portfolio {portfolio_id}")
                    failed = True
                if failed:
                    self.portfolio_collection.update_one({"_id": ObjectId(portfolio_id)}, {"$set": {"recommendation_status": RecommendationStatus.FAILED}})
                    if app_name.upper() == "CCA":
                        application_name = self.CCA_APP_NAME
                    else:
                        application_name = self.EIA_APP_NAME
                    user_subject = f"{application_name} Recommendation Failed for Portfolio {portfolio_id}"

                    # Get link for current app (fallback to EIA if not found)
                    support_url = self.support_links.get(app_name.upper())

                    # Build subject and body
                    user_subject = f"{application_name} Recommendation Failed for Portfolio: {portfolio_doc.get('name')}"

                    user_body = f"""
                    <html>
                    <body>
                        <p>Hello,</p>
                        <p>We regret to inform you that the recommendation generation for your portfolio <b>{portfolio_doc.get('name')}</b> has failed.</p>
                        <p>Please contact the 
                            <a href="{support_url}" target="_blank" style="color:#0078d4; text-decoration:none;">
                                EPYC Advisory Services
                            </a> 
                            team for assistance or more information.
                        </p>
                        <br/>
                        <p>Best regards,<br/>{SENDER}</p>
                    </body>
                    </html>
                    """

                    self.send_email(user_subject, portfolio_doc.get('user_email'), user_body)
                    self.log_message(LevelType.ERROR, f"Portfolio {portfolio_id} marked FAILED due to invalid application or provider selected")
                    continue
                self.log_message(LevelType.INFO, f"Starting batch processing for portfolio {portfolio_id}")
                email_notification = True
               
                # Get total number of batches for percentage calculation
                total_batches = self.get_total_batches_for_portfolio(portfolio_id)
                while True:
                    start_time=datetime.now()
                    next_batch = self.get_next_batch_for_portfolio(portfolio_id)
                    if not next_batch:
                        email_notification = self.update_portfolio_status_if_all_batches_completed(portfolio_id, rec_failed_to_emails, user_email, portfolio_name, app_name.upper())
                        self.log_message(LevelType.INFO, f"All batches processed for portfolio {portfolio_id}")
                        break
                    batch_id = next_batch["batch_id"]
                    rec_progress_id = next_batch["_id"]
                    batch_processed, remarks = self.process_batch(portfolio_id, batch_id, cloud_csp, portfolio_doc, rec_progress_id, app_name)
                    process_count = self.get_batch_count_on_status(portfolio_id)
                    percentage = int((process_count / total_batches) * 100)
                    if int(process_count)<int(total_batches):
                        self.portfolio_collection.update_one(
                            {"_id": ObjectId(portfolio_id)},
                            {"$set": {"recommendation_percentage": percentage}}
                        )
                        self.log_message(LevelType.INFO, f"Portfolio {portfolio_id} recommendation percentage updated to {percentage}%")
                   
                    if not batch_processed:
                        self.recommendation_tracking_collection.update_one({"_id": rec_progress_id}, {"$set": {"recommendation_status": "FAILED"}})
                        self.log_message(LevelType.ERROR, f"Batch {batch_id} processing failed for portfolio {portfolio_id}. Marked batch as FAILED.")
                    if remarks:
                        created_at = datetime.now(timezone.utc).isoformat()
                        remarks = [{**record, 'portfolio_id': portfolio_id, 'batch_id': batch_id, 'app_name': app_name.upper(), 'created_at': created_at} for record in remarks]
                        self.input_remarks_collection.delete_many({'portfolio_id': portfolio_id, 'batch_id': batch_id})
                        res = self.input_remarks_collection.insert_many(remarks)
                        self.log_message(LevelType.INFO, f"Inserted {len(res.inserted_ids)} error records for portfolio {portfolio_id} batch {batch_id}")
                    end_time = datetime.now()
                    self.log_message(LevelType.INFO, f"Execution time for portfolio: {portfolio_id} is {end_time - start_time}")
                if email_notification == True:
                    if app_name.upper() == "CCA":
                        notification_message = f"Recommendation of application {self.CCA_APP_NAME} for portfolio '{portfolio_doc.get('name')}' has been completed."
                        title = "EPYC Cloud Cost Advice Completed"
                        purpose="COST_ADVICE"
                    else:
                        notification_message = f"Recommendation of application {self.EIA_APP_NAME} for portfolio '{portfolio_doc.get('name')}' has been completed."
                        title = "EPYC Instance Advice Completed"
                        purpose="INSTANCE_ADVICE"
                    self.download_csv(portfolio_id, app_name.upper(), portfolio_doc.get('name'), portfolio_doc.get("user_email"))
                else:
                    if app_name.upper() == "CCA":
                        notification_message = f"Recommendation of application {self.CCA_APP_NAME} for portfolio '{portfolio_doc.get('name')}' is failed. Please contact EPYC Advisory team for further details."
                        title = "EPYC Cloud Cost Advice Failed"
                        purpose="COST_ADVICE"
                    else:
                        notification_message = f"Recommendation of application {self.EIA_APP_NAME} for portfolio '{portfolio_doc.get('name')}' is failed. Please contact EPYC Advisory team for further details."
                        title = "EPYC Instance Advice Failed"
                        purpose="INSTANCE_ADVICE"
                notif_id = self.add_notification(
                    user_email=portfolio_doc.get('user_email'),
                    app_name=app_name.upper(),
                    portfolio_id=portfolio_id,
                    portfolio_name=portfolio_doc.get('name'),
                    purpose=purpose,
                    message=notification_message,
                    title=title
                )
                self.etl_process_portfolio(portfolio_id)
                # Check counts of relevant statuses in one query per status list for performance
                failed_count = self.recommendation_tracking_collection.count_documents({
                    "portfolio_id": portfolio_id,
                    "recommendation_status": "FAILED"
                })
 
                in_progress_count = self.recommendation_tracking_collection.count_documents({
                    "portfolio_id": portfolio_id,
                    "recommendation_status": {"$in": ["QUEUE", "TO_PROCESS", "INPROGRESS"]}
                })
 
                if in_progress_count > 0:
                    self.log_message(LevelType.INFO, f"Portfolio {portfolio_id} recommendation_status yet to be processed")
                else:
                    update_fields = {
                        "recommendation_status": "FAILED" if failed_count > 0 else "COMPLETED",
                        "recommendation_percentage": 100,
                        "is_recommendation_available": failed_count == 0
                    }
 
                    self.portfolio_collection.update_one(
                        {"_id": ObjectId(portfolio_id)},
                        {"$set": update_fields}
                    )
 
                    self.log_message(LevelType.INFO, f"Portfolio {portfolio_id} recommendation_status set to {update_fields['recommendation_status']}")
                    self.log_message(LevelType.INFO, f"Notification added for portfolio: {portfolio_doc.get('name')} notification_id: {notif_id}")
                    self.log_message(LevelType.INFO, f"Cost or Instance advice completed for portfolio: {portfolio_doc.get('name')}, portfolio_id: {portfolio_id}")
            self.close_connection()
        except Exception as e:
            self.log_message(LevelType.ERROR, f"handle_bulk_data error: {str(e)}")

# ----------------- Example Usage -----------------
if __name__ == "__main__":
    reader = PortfolioReader()

    # Handle bulk portfolios
    reader.handle_bulk_data()
